# app_escuela/api/views.py
import math
import logging
import os
import base64
import tempfile
import random
import re

from copy import copy
from datetime import date, datetime, timedelta
from decimal import Decimal, ROUND_HALF_UP, InvalidOperation
from io import BytesIO
from django.conf import settings
from django.contrib.auth import authenticate
from django.contrib.auth.models import update_last_login
from django.db import IntegrityError, models, transaction
from django.db.models import (
    F,
    Q,
    Sum,
    Count,
    FloatField,
    Prefetch,
    Exists,
    OuterRef,
    Subquery,
)
from django.db.models.functions import Cast, TruncMonth
from django.http import HttpResponse
from django.shortcuts import get_object_or_404
from django.utils import timezone
from django.utils.dateparse import parse_date
from django_filters.rest_framework import DjangoFilterBackend
from rest_framework import status, viewsets, serializers
from rest_framework.authtoken.models import Token
from rest_framework.decorators import (
    action,
    api_view,
    authentication_classes,
    permission_classes,
    throttle_classes,
)
from rest_framework.parsers import MultiPartParser, FormParser, JSONParser
from rest_framework.permissions import IsAuthenticated, AllowAny
from rest_framework.response import Response
from rest_framework.throttling import SimpleRateThrottle
from rest_framework.views import APIView
from openpyxl import Workbook, load_workbook
from openpyxl.drawing.image import Image as ExcelImage
from openpyxl.styles import Font, PatternFill, Border, Side, Alignment
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from .pagination import PaginacionOpcional
from ..models import (
    Rol,
    Usuario,
    Estudiante,
    Instructor,
    CategoriaVehiculo,
    PlanEstudio,
    Matricula,
    Recibo,
    Calendario,
    Asistencia,
    Notas,
    ValorCurso,
    PreguntaExamenTeorico,
    OpcionPreguntaExamenTeorico,
    ExamenTeorico,
    IntentoExamenTeorico,
    PreguntaIntentoExamenTeorico,
    RespuestaExamenTeorico,
    PagoInstructor,
    CargoInstitucional,
    ProgresoTema,
    ProgresoClaseTema,
    Notificacion,
    HistorialPlanEstudio,
)
from .serializers import (
    RolSerializer,
    UserSerializer,
    CargoInstitucionalSerializer,
    EstudianteSerializer,
    InstructorSerializer,
    InstructorListSerializer,
    InstructorCalendarioSerializer,
    CategoriaVehiculoSerializer,
    MatriculaSerializer,
    ReciboSerializer,
    CalendarioSerializer,
    CrearBloqueCitasSerializer,
    CrearCalendarioManualSerializer,
    AsistenciaSerializer,
    NotasSerializer,
    ValorCursoSerializer,
    PreguntaExamenTeoricoSerializer,
    ExamenTeoricoSerializer,
    PreguntaExamenEstudianteSerializer,
    RespuestaEnviarExamenSerializer,
    RespuestaExamenTeoricoSerializer,
    PagoInstructorSerializer,
    PlanEstudioSerializer,
    ProgresoTemaSerializer,
    NotificacionSerializer,
    actualizar_estado_matricula_por_notas,
)
from .permissions import (
    es_administrativo as es_admin,
    es_estudiante,
    es_instructor,
    obtener_rol_usuario as obtener_rol,
)

logger = logging.getLogger(__name__)
logging.getLogger("PIL").setLevel(logging.WARNING)

MAX_KILOMETRAJE = Decimal('99999999.99')

def validar_kilometraje(valor, nombre_campo):
    try:
        kilometraje = Decimal(
            str(valor).strip()
        )
    except (
        InvalidOperation,
        ValueError,
        TypeError,
    ):
        raise ValueError(
            f'El {nombre_campo} debe ser numérico.'
        )

    if not kilometraje.is_finite():
        raise ValueError(
            f'El {nombre_campo} debe ser un número válido.'
        )

    if kilometraje < 0:
        raise ValueError(
            f'El {nombre_campo} no puede ser negativo.'
        )

    if kilometraje > MAX_KILOMETRAJE:
        raise ValueError(
            f'El {nombre_campo} supera el valor permitido.'
        )

    return kilometraje

class LoginRateThrottle(SimpleRateThrottle):
    """
    Limita los intentos de inicio de sesión por dirección IP.
    La cantidad permitida se configura en settings.py
    mediante la clave 'login'.
    """

    scope = 'login'

    def get_cache_key(self, request, view):
        return self.cache_format % {
            'scope': self.scope,
            'ident': self.get_ident(request),
        }

def matricula_usa_checks(matricula):
    """
    Solamente los estudiantes del curso Principiante
    utilizan checks y progreso del plan de estudio.
    """

    tipo_curso = str(
        getattr(matricula, 'tipo_curso', '') or ''
    ).strip().lower()

    return tipo_curso == 'principiante'

def obtener_progresos_plan_actual(matricula):
    """
    Devuelve solamente los progresos pertenecientes al plan
    que está asignado actualmente a la matrícula.
    """

    if not matricula.plan_de_estudio_id:
        return ProgresoTema.objects.none()

    return ProgresoTema.objects.filter(
        matricula=matricula,
        tema__plan_estudio_id=matricula.plan_de_estudio_id,
        tema__activo=True,
    )

def validar_plan_completado_para_examen(matricula):
    """
    Devuelve información de seguimiento del plan de estudio.
    Los checks no habilitan ni bloquean el examen teórico.
    Solamente sirven como control visual/académico.
    """

    tipo = str(
        getattr(matricula, 'tipo_curso', '') or ''
    ).strip().lower()

    if tipo in ['intermedio', 'avanzado']:
        return {
            'completo': True,
            'error': None,
            'progreso': None,
            'usa_checks': False,
        }

    if tipo != 'principiante':
        return {
            'completo': True,
            'error': None,
            'progreso': '0/0',
            'usa_checks': False,
        }

    progresos = obtener_progresos_plan_actual(
        matricula
    )

    total_temas = progresos.count()

    if total_temas == 0:
        return {
            'completo': True,
            'error': None,
            'progreso': '0/0',
            'usa_checks': True,
        }

    completados = progresos.filter(
        Q(completado=True)
        |
        Q(
            estudiante_completado=True,
            instructor_completado=True,
        )
    ).distinct().count()

    return {
        'completo': True,
        'error': None,
        'progreso': f'{completados}/{total_temas}',
        'usa_checks': True,
    }

def obtener_rango_horario(matricula):
    mapeo = {
        '06AM': ('06:00', '08:00'),
        '08AM': ('08:00', '10:00'),
        '10AM': ('10:00', '12:00'),
        '12PM': ('12:00', '14:00'),
        '02PM': ('14:00', '16:00'),
        '04PM': ('16:00', '18:00'),
    }

    return mapeo.get(matricula.horario)

def crear_clase_recuperacion(clase_faltada):
    matricula = clase_faltada.matricula

    modalidad = str(
        matricula.modalidad or ''
    ).strip().lower()

    def siguiente_fecha_valida(fecha_base):
        nueva_fecha = fecha_base + timedelta(days=1)

        while True:
            es_fin_semana = (
                nueva_fecha.weekday() >= 5
            )

            if modalidad == 'extraordinario':
                if es_fin_semana:
                    return nueva_fecha

            elif modalidad == 'mixto':
                return nueva_fecha

            elif not es_fin_semana:
                return nueva_fecha

            nueva_fecha += timedelta(days=1)

    ultima_clase_numero = (
        Calendario.objects
        .filter(
            matricula=matricula,
            es_examen=False,
        )
        .order_by(
            '-numero_clase',
            '-id',
        )
        .first()
    )

    ultimo_numero = (
        ultima_clase_numero.numero_clase
        if ultima_clase_numero
        else 0
    )

    ultima_clase_fecha = (
        Calendario.objects
        .filter(
            matricula=matricula,
            es_examen=False,
        )
        .order_by(
            '-fecha',
            '-hora_inicio',
            '-id',
        )
        .first()
    )

    fecha_base = (
        ultima_clase_fecha.fecha
        if ultima_clase_fecha
        else clase_faltada.fecha
    )

    fecha_recuperacion = siguiente_fecha_valida(
        fecha_base
    )

    while (
        Calendario.objects
        .filter(
            instructor=clase_faltada.instructor,
            fecha=fecha_recuperacion,
            hora_inicio__lt=clase_faltada.hora_fin,
            hora_fin__gt=clase_faltada.hora_inicio,
        )
        .exclude(
            estado='cancelada'
        )
        .exists()
    ):
        fecha_recuperacion = siguiente_fecha_valida(
            fecha_recuperacion
        )

    return Calendario.objects.create(
        matricula=matricula,
        instructor=clase_faltada.instructor,
        fecha=fecha_recuperacion,
        hora_inicio=clase_faltada.hora_inicio,
        hora_fin=clase_faltada.hora_fin,
        numero_clase=ultimo_numero + 1,
        estado='pendiente',
        es_examen=False,
    )

def desactivar_usuario(usuario):
    if not usuario:
        return

    usuario.is_active = False
    usuario.save(update_fields=['is_active'])

    Token.objects.filter(user=usuario).delete()

def desactivar_usuarios_estudiante(estudiante):
    usuarios = estudiante.usuarios.all()

    for usuario in usuarios:
        desactivar_usuario(usuario)

    estudiante.activo = False
    estudiante.save(update_fields=['activo'])

def desactivar_usuarios_instructor(instructor):
    usuarios = instructor.usuarios.all()

    for usuario in usuarios:
        desactivar_usuario(usuario)

def obtener_ids_matriculas_egresadas():
    return (
        Matricula.objects
        .filter(
            estado='finalizado'
        )
        .values_list(
            'id',
            flat=True,
        )
    )

def matricula_tiene_notas_teorica_y_practica(matricula):
    tipos_notas = set(
        Notas.objects
        .filter(
            matricula=matricula,
            tipo_nota__in=[
                'teorico',
                'practico',
            ],
        )
        .values_list(
            'tipo_nota',
            flat=True,
        )
    )

    return {
        'teorico',
        'practico',
    }.issubset(tipos_notas)

TIPOS_RECIBO_INGRESO = [
    'completo',
    'anticipo',
    'beneficio',
]

def obtener_rango_mes(
    anio,
    mes,
):
    inicio_mes = date(
        anio,
        mes,
        1,
    )

    if mes == 12:
        inicio_mes_siguiente = date(
            anio + 1,
            1,
            1,
        )
    else:
        inicio_mes_siguiente = date(
            anio,
            mes + 1,
            1,
        )

    return (
        inicio_mes,
        inicio_mes_siguiente,
    )

def obtener_foto_instructor(instructor):
    if not instructor:
        return None

    return getattr(instructor, "foto_base64", None)

def crear_archivo_temporal_desde_base64(foto_base64):
    if not foto_base64:
        return None

    try:
        formato, imgstr = foto_base64.split(";base64,")
        extension = formato.split("/")[-1]

        archivo = tempfile.NamedTemporaryFile(
            delete=False,
            suffix=f".{extension}"
        )

        archivo.write(base64.b64decode(imgstr))
        archivo.close()

        return archivo.name

    except Exception:
        return None

def obtener_ruta_foto_instructor_para_excel(instructor):
    if not instructor:
        return None

    return crear_archivo_temporal_desde_base64(
        getattr(instructor, "foto_base64", None)
    )

class RolViewSet(viewsets.ModelViewSet):
    queryset = Rol.objects.all()
    serializer_class = RolSerializer
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)

    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

class CategoriaVehiculoViewSet(viewsets.ModelViewSet):
    queryset = CategoriaVehiculo.objects.all()
    serializer_class = CategoriaVehiculoSerializer
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)


    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

class EstudianteViewSet(viewsets.ModelViewSet):
    queryset = Estudiante.objects.all()
    serializer_class = EstudianteSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        queryset = (
            Estudiante.objects
            .prefetch_related(
                Prefetch(
                    'usuarios',
                    queryset=(
                        Usuario.objects
                        .filter(
                            rol__nombre__iexact='estudiante'
                        )
                        .select_related('rol')
                        .order_by('id')
                    ),
                    to_attr='usuarios_estudiante_precargados',
                ),
                Prefetch(
                    'matriculas',
                    queryset=(
                        Matricula.objects
                        .only(
                            'id',
                            'estudiante_id',
                            'estado',
                        )
                        .order_by('-id')
                    ),
                    to_attr='matriculas_precargadas',
                ),
            )
            .order_by('-id')
        )

        buscar = self.request.query_params.get('buscar')

        if es_admin(user):
            pass

        elif es_instructor(user):
            estudiantes_ids = Calendario.objects.filter(
                instructor_id=user.instructor_id,
                es_examen=False,
                matricula__estado__in=[
                    'pendiente',
                    'matriculado',
                ],
            ).exclude(
                estado='cancelada'
            ).values_list(
                'matricula__estudiante_id',
                flat=True
            ).distinct()

            queryset = queryset.filter(
                id__in=estudiantes_ids
            )

        elif es_estudiante(user):
            queryset = queryset.filter(id=user.estudiante_id)

        else:
            return Estudiante.objects.none()

        if buscar:
            buscar = buscar.strip()

            queryset = queryset.filter(
                Q(nombre__icontains=buscar)
                | Q(apellido__icontains=buscar)
                | Q(cedula__icontains=buscar)
                | Q(correo_electronico__icontains=buscar)
                | Q(telefono_movil__icontains=buscar)
            )

        return queryset

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'registrar estudiantes.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(
            request,
            *args,
            **kwargs
        )

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'editar estudiantes.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(
            request,
            *args,
            **kwargs
        )

    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'editar estudiantes.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(
            request,
            *args,
            **kwargs
        )

    @action(detail=False, methods=['get'], url_path='resumen')
    def resumen(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo Administración puede consultar '
                        'el resumen de estudiantes.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        total = Estudiante.objects.count()
        activos = Estudiante.objects.filter(
            activo=True
        ).count()
        inactivos = Estudiante.objects.filter(
            activo=False
        ).count()

        return Response({
            'total': total,
            'activos': activos,
            'inactivos': inactivos,
        })

class PlanEstudioViewSet(viewsets.ModelViewSet):
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)


    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

    queryset = PlanEstudio.objects.prefetch_related(
        'temas',
        'temas__subtemas'
    ).all()

    serializer_class = PlanEstudioSerializer

    @action(detail=False, methods=['get'], url_path='tipos-curso')
    def tipos_curso(self, request):
        return Response([
            {'value': 'Principiante', 'label': 'Principiante'},
            {'value': 'Intermedio', 'label': 'Intermedio'},
            {'value': 'Avanzado', 'label': 'Avanzado'},
        ])

class ValorCursoViewSet(viewsets.ModelViewSet):
    queryset = ValorCurso.objects.all().order_by('-fecha_modificacion')
    serializer_class = ValorCursoSerializer
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)

    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

    def get_queryset(self):
        queryset = ValorCurso.objects.all().order_by('-fecha_modificacion')
        activo = self.request.query_params.get('activo')
        tipo_curso = self.request.query_params.get('tipo_curso')

        if activo is not None:
            if activo.lower() == 'true':
                queryset = queryset.filter(activo=True)
            elif activo.lower() == 'false':
                queryset = queryset.filter(activo=False)

        if tipo_curso:
            queryset = queryset.filter(tipo_curso=tipo_curso)

        return queryset


class InstructorViewSet(viewsets.ModelViewSet):
    queryset = Instructor.objects.all()
    serializer_class = InstructorSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]
    parser_classes = [JSONParser, MultiPartParser, FormParser]

    def get_serializer_class(self):
        if es_estudiante(self.request.user):
            return InstructorCalendarioSerializer

        if self.action == 'list':
            return InstructorListSerializer

        return InstructorSerializer

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)


    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

    def get_queryset(self):
        user = self.request.user
        queryset = Instructor.objects.order_by('-id')

        if self.action == 'list':
            queryset = queryset.defer('foto_base64')

        if es_admin(user):
            pass

        elif es_instructor(user):
            queryset = queryset.filter(
                id=user.instructor_id
            )

        elif es_estudiante(user):
            instructores_ids = Calendario.objects.filter(
                matricula__estudiante_id=user.estudiante_id
            ).exclude(
                estado='cancelada'
            ).values_list(
                'instructor_id',
                flat=True
            ).distinct()

            queryset = queryset.filter(
                id__in=instructores_ids
            )

        else:
            return Instructor.objects.none()

        buscar = self.request.query_params.get('buscar')

        if buscar:
            buscar = buscar.strip()

            queryset = queryset.filter(
                Q(nombre__icontains=buscar)
                | Q(apellido__icontains=buscar)
                | Q(cedula__icontains=buscar)
                | Q(numero_telefono__icontains=buscar)
                | Q(categoria_instructor__icontains=buscar)
            )

        activo = self.request.query_params.get('activo')

        if activo is not None:
            activo = activo.strip().lower()

            if activo == 'true':
                queryset = queryset.filter(activo=True)
            elif activo == 'false':
                queryset = queryset.filter(activo=False)

        return queryset

    @action(detail=True, methods=['post'], url_path='desactivar')
    def despedir(self, request, pk=None):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'desactivar instructores.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        instructor = self.get_object()

        fecha_salida = request.data.get('fecha_salida') or timezone.now().date()
        motivo_salida = request.data.get('motivo_salida') or 'Instructor desactivado por administración.'

        instructor.fecha_salida = fecha_salida
        instructor.motivo_salida = motivo_salida
        instructor.activo = False
        instructor.save(update_fields=['fecha_salida', 'motivo_salida', 'activo'])

        desactivar_usuarios_instructor(instructor)

        return Response({
            'message': 'Instructor desactivado correctamente. Su usuario ya no puede acceder.',
            'instructor': self.get_serializer(instructor).data
        })


class MatriculaViewSet(viewsets.ModelViewSet):
    queryset = Matricula.objects.select_related('estudiante').all()
    serializer_class = MatriculaSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        queryset = (
            Matricula.objects
            .select_related(
                'estudiante',
                'categoria',
                'plan_de_estudio',
            )
            .prefetch_related(
                Prefetch(
                    'estudiante__usuarios',
                    queryset=(
                        Usuario.objects
                        .only(
                            'id',
                            'estudiante_id',
                        )
                    ),
                    to_attr='usuarios_precargados',
                )
            )
            .order_by('-id')
        )

        if es_admin(user):
            pass

        elif es_instructor(user):
            matriculas_ids = Calendario.objects.filter(
                instructor_id=user.instructor_id,
                es_examen=False,
            ).exclude(
                estado='cancelada'
            ).values_list(
                'matricula_id',
                flat=True
            ).distinct()

            queryset = queryset.filter(
                id__in=matriculas_ids
            )

        elif es_estudiante(user):
            queryset = queryset.filter(
                estudiante_id=user.estudiante_id
            )

        else:
            return Matricula.objects.none()

        buscar = self.request.query_params.get('buscar')
        solo_activos = self.request.query_params.get(
            'solo_activos'
        )
        estado = self.request.query_params.get('estado')

        if buscar:
            buscar = buscar.strip()

            queryset = queryset.filter(
                Q(estudiante__cedula__icontains=buscar)
                | Q(estudiante__nombre__icontains=buscar)
                | Q(estudiante__apellido__icontains=buscar)
                | Q(estudiante__telefono_movil__icontains=buscar)
                | Q(tipo_curso__icontains=buscar)
                | Q(categoria__nombre__icontains=buscar)
                | Q(estado__icontains=buscar)
            )

        if (
            solo_activos
            and solo_activos.strip().lower() == 'true'
        ):
            queryset = queryset.exclude(
                estado='finalizado'
            )

        if estado:
            queryset = queryset.filter(
                estado=estado
            )

        return queryset

    @action(detail=False, methods=['get'], url_path='resumen')
    def resumen(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo Administración puede consultar '
                        'el resumen de matrículas.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        queryset = Matricula.objects.all()

        return Response({
            'total': queryset.count(),
            'matriculados': queryset.filter(
                estado='matriculado'
            ).count(),
            'pendientes': queryset.filter(
                estado='pendiente'
            ).count(),
            'finalizados': queryset.filter(
                estado='finalizado'
            ).count(),
        })

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede crear matrículas.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar matrículas.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)


    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar matrículas.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

    @action(detail=False, methods=['get'], url_path='buscar-estudiante')
    def buscar_estudiante(self, request):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para buscar estudiantes.'},
                status=status.HTTP_403_FORBIDDEN
            )
        q = request.query_params.get('q')

        if not q:
            return Response(
                {'error': 'Debe enviar un nombre o una cédula para buscar.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        estudiantes = Estudiante.objects.filter(
            Q(cedula__icontains=q) |
            Q(nombre__icontains=q) |
            Q(apellido__icontains=q)
        ).order_by('-id')[:10]

        resultados = []

        for estudiante in estudiantes:
            resultados.append({
                'id': estudiante.id,
                'nombre': estudiante.nombre,
                'apellido': estudiante.apellido,
                'cedula': estudiante.cedula,
                'edad': estudiante.edad,
                'sexo': estudiante.sexo,
                'nacionalidad': estudiante.nacionalidad,
                'fecha_nacimiento': estudiante.fecha_nacimiento,
                'direccion': estudiante.direccion,
                'correo_electronico': estudiante.correo_electronico,
                'telefono_movil': estudiante.telefono_movil,
                'nivel_educativo': estudiante.nivel_educativo,
                'nombre_emergencia': estudiante.nombre_emergencia,
                'telefono_emergencia': estudiante.telefono_emergencia,
                'tiene_usuario': estudiante.usuarios.exists(),
            })

        return Response(resultados)

    @action(detail=False, methods=['get'], url_path='para-examen')
    def para_examen(self, request):
        user = request.user

        if not es_instructor(user):
            return Response(
                [],
                status=status.HTTP_200_OK
            )

        matriculas_ids = (
            Calendario.objects
            .filter(
                instructor_id=user.instructor_id,
                es_examen=False,
                matricula__estado__in=[
                    'matriculado',
                    'finalizado',
                ],
            )
            .exclude(estado='cancelada')
            .values_list(
                'matricula_id',
                flat=True
            )
            .distinct()
        )

        clases_regulares = (
            Calendario.objects
            .filter(
                matricula_id=OuterRef('pk'),
                es_examen=False,
            )
            .exclude(estado='cancelada')
        )

        clases_pendientes = clases_regulares.exclude(
            estado='completada'
        )

        examenes = Calendario.objects.filter(
            matricula_id=OuterRef('pk'),
            es_examen=True,
        )

        examenes_vigentes = examenes.exclude(
            estado='cancelada'
        )

        notas_teoricas = Notas.objects.filter(
            matricula_id=OuterRef('pk'),
            tipo_nota='teorico',
        )

        notas_practicas = Notas.objects.filter(
            matricula_id=OuterRef('pk'),
            tipo_nota='practico',
        )

        matriculas = (
            Matricula.objects
            .select_related('estudiante')
            .filter(
                id__in=matriculas_ids,
                estado__in=[
                    'matriculado',
                    'finalizado',
                ],
            )
            .filter(
                Q(tipo_curso__iexact='Principiante')
                | (
                    Q(tipo_curso__iexact='Intermedio')
                    & Q(incluye_examen_policial=True)
                )
                | (
                    Q(tipo_curso__iexact='Avanzado')
                    & Q(incluye_examen_policial=True)
                )
            )
            .annotate(
                tiene_clases_regulares=Exists(
                    clases_regulares
                ),
                tiene_clases_pendientes=Exists(
                    clases_pendientes
                ),
                tiene_nota_teorica=Exists(
                    notas_teoricas
                ),
                tiene_nota_practica=Exists(
                    notas_practicas
                ),
                total_examenes_asignados=Count(
                    'clases',
                    filter=Q(
                        clases__es_examen=True
                    ),
                    distinct=True,
                ),
                tiene_examen_vigente=Exists(
                    examenes_vigentes
                ),
            )
            .filter(
                tiene_clases_regulares=True,
                tiene_clases_pendientes=False,
                tiene_nota_teorica=True,
                tiene_nota_practica=True,
                total_examenes_asignados__lt=3,
                tiene_examen_vigente=False,
            )
            .order_by(
                'estudiante__nombre',
                'estudiante__apellido',
            )
        )

        resultados = []

        for matricula in matriculas:
            resultados.append({
                'id': matricula.id,
                'estudiante_nombre': (
                    f'{matricula.estudiante.nombre} '
                    f'{matricula.estudiante.apellido}'
                ).strip(),
                'estudiante_cedula': (
                    matricula.estudiante.cedula
                ),
                'tipo_curso': matricula.tipo_curso,
                'incluye_examen_policial': (
                    matricula.incluye_examen_policial
                ),
                'examenes_asignados': (
                    matricula.total_examenes_asignados
                ),
                'examenes_disponibles': (
                    3 - matricula.total_examenes_asignados
                ),
            })

        return Response(resultados)

    @action(detail=False, methods=['get'], url_path='asignadas-instructor')
    def asignadas_instructor(self, request):
        user = request.user

        if not es_admin(user) and not es_instructor(user):
            return Response(
                {
                    'error': (
                        'Solo el administrador o un instructor '
                        'pueden consultar las matrículas asignadas.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        examenes = Calendario.objects.filter(
            matricula_id=OuterRef('pk'),
            es_examen=True,
        )

        matriculas = (
            Matricula.objects
            .select_related('estudiante')
            .filter(estado='matriculado')
            .annotate(
                tiene_examen=Exists(examenes)
            )
        )

        if es_instructor(user):
            if not user.instructor_id:
                return Response([], status=200)

            matriculas_ids = Calendario.objects.filter(
                instructor_id=user.instructor_id,
                es_examen=False
            ).values_list(
                'matricula_id',
                flat=True
            ).distinct()

            matriculas = matriculas.filter(id__in=matriculas_ids)

        resultados = []

        for matricula in matriculas:

            resultados.append({
                'id': matricula.id,
                'estudiante_nombre': (
                    f"{matricula.estudiante.nombre} "
                    f"{matricula.estudiante.apellido}"
                ).strip(),
                'estudiante_cedula': matricula.estudiante.cedula,
                'tiene_examen': matricula.tiene_examen,
            })

        return Response(resultados)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def saldo(request):
    matricula_id = request.query_params.get('matricula')

    if not matricula_id:
        return Response(
            {'error': 'Debe enviar el ID de la matrícula.'},
            status=400
        )

    try:
        matricula = Matricula.objects.select_related('estudiante').get(id=matricula_id)
    except Matricula.DoesNotExist:
        return Response(
            {'error': 'La matrícula no existe.'},
            status=404
        )

    if not es_admin(request.user):
        es_propietario = (
            es_estudiante(request.user)
            and matricula.estudiante_id == request.user.estudiante_id
        )

        if not es_propietario:
            return Response(
                {
                    'error': (
                        'No tienes permiso para consultar '
                        'el saldo de esta matrícula.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

    def calcular_monto_total(matricula):
        primer_recibo = Recibo.objects.filter(
            matricula=matricula,
            valor_curso__isnull=False
        ).select_related(
            'valor_curso'
        ).order_by(
            'id'
        ).first()

        if primer_recibo:
            valor_curso = primer_recibo.valor_curso
        else:
            valor_curso = ValorCurso.objects.filter(
                tipo_curso=matricula.tipo_curso,
                activo=True
            ).order_by(
                '-fecha_modificacion',
                '-id'
            ).first()

        if not valor_curso:
            raise ValueError(
                f'No existe un valor activo para el curso '
                f'{matricula.tipo_curso}.'
            )

        if matricula.tipo_curso == 'Principiante':
            monto_total = Decimal(
                str(valor_curso.precio_total)
            )

        elif matricula.tipo_curso in [
            'Intermedio',
            'Avanzado',
        ]:
            horas = matricula.horas_reforzamiento

            if not horas:
                raise ValueError(
                    f'La matrícula del curso '
                    f'{matricula.tipo_curso} '
                    f'no tiene horas asignadas.'
                )

            monto_total = (
                Decimal(str(horas))
                * Decimal(str(valor_curso.precio_hora))
            )

        else:
            monto_total = Decimal('0')

        return monto_total.quantize(
            Decimal('1'),
            rounding=ROUND_HALF_UP
        )

    try:
        monto_total = calcular_monto_total(matricula)
    except ValueError as error:
        return Response(
            {'error': str(error)},
            status=400
        )

    total_pagado = Recibo.objects.filter(
        matricula=matricula,
    ).aggregate(
        total=models.Sum('monto_pagado')
    )['total'] or Decimal('0')

    cantidad_pagos = Recibo.objects.filter(matricula=matricula).count()

    saldo_pendiente = monto_total - total_pagado

    return Response({
        'matricula_id': matricula.id,
        'nombre': matricula.estudiante.nombre,
        'apellido': matricula.estudiante.apellido,
        'cedula': matricula.estudiante.cedula,
        'tipo_curso': matricula.tipo_curso,
        'horas_reforzamiento': matricula.horas_reforzamiento,
        'monto_total': float(monto_total),
        'total_pagado': float(total_pagado),
        'saldo_pendiente': float(saldo_pendiente),
        'cantidad_pagos': cantidad_pagos,
        'pagos_permitidos': 2,
    })

class ReciboViewSet(viewsets.ModelViewSet):
    queryset = Recibo.objects.all()
    serializer_class = ReciboSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        queryset = Recibo.objects.select_related(
            'matricula',
            'matricula__estudiante',
            'valor_curso',
        ).all().order_by('-id')

        buscar = self.request.query_params.get(
            'buscar'
        )

        if buscar:
            buscar = buscar.strip()

            queryset = queryset.filter(
                Q(numero_recibo__icontains=buscar)
                | Q(
                    matricula__estudiante__nombre__icontains=buscar
                )
                | Q(
                    matricula__estudiante__apellido__icontains=buscar
                )
                | Q(
                    matricula__estudiante__cedula__icontains=buscar
                )
            )

        if es_admin(user):
            return queryset

        if es_estudiante(user):
            return queryset.filter(
                matricula__estudiante_id=user.estudiante_id
            )

        return Recibo.objects.none()

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede registrar recibos.'},
                status=status.HTTP_403_FORBIDDEN
            )
        try:
            return super().create(
                request,
                *args,
                **kwargs
            )
        except IntegrityError:
            return Response(
                {
                    'numero_recibo': [
                        (
                            'Ya existe un recibo con este '
                            'número. Verifique el número '
                            'e inténtelo nuevamente.'
                        )
                    ]
                },
                status=status.HTTP_400_BAD_REQUEST
            )

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar recibos.'},
                status=status.HTTP_403_FORBIDDEN
            )
        try:
            return super().update(
                request,
                *args,
                **kwargs
            )
        except IntegrityError:
            return Response(
                {
                    'numero_recibo': [
                        (
                            'Ya existe otro recibo con '
                            'este número.'
                        )
                    ]
                },
                status=status.HTTP_400_BAD_REQUEST
            )


    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar recibos.'},
                status=status.HTTP_403_FORBIDDEN
            )
        try:
            return super().partial_update(
                request,
                *args,
                **kwargs
            )
        except IntegrityError:
            return Response(
                {
                    'numero_recibo': [
                        (
                            'Ya existe otro recibo con '
                            'este número.'
                        )
                    ]
                },
                status=status.HTTP_400_BAD_REQUEST
            )

    @action(detail=False, methods=['get'], url_path='resumen')
    def resumen(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede ver '
                        'el resumen de recibos.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        hoy = timezone.localdate()

        inicio_mes, inicio_mes_siguiente = obtener_rango_mes(
            hoy.year,
            hoy.month,
        )

        queryset = self.get_queryset().filter(
            tipo_pago__in=TIPOS_RECIBO_INGRESO,
        )

        ingresos_mes = (
            queryset
            .filter(
                fecha_pago__gte=inicio_mes,
                fecha_pago__lt=inicio_mes_siguiente,
            )
            .aggregate(
                total=Sum('monto_pagado')
            )['total']
            or Decimal('0')
        )

        ingresos_totales = (
            queryset
            .aggregate(
                total=Sum('monto_pagado')
            )['total']
            or Decimal('0')
        )

        recibos_mes = (
            queryset
            .filter(
                fecha_pago__gte=inicio_mes,
                fecha_pago__lt=inicio_mes_siguiente,
            )
            .count()
        )

        return Response({
            'ingresos_mes': float(ingresos_mes),
            'ingresos_totales': float(ingresos_totales),
            'recibos_mes': recibos_mes,
        })

class UserViewSet(viewsets.ModelViewSet):
    serializer_class = UserSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        queryset = (
            Usuario.objects
            .select_related(
                'rol',
                'estudiante',
                'instructor',
            )
            .prefetch_related(
                Prefetch(
                    'estudiante__matriculas',
                    queryset=(
                        Matricula.objects
                        .only(
                            'id',
                            'estudiante_id',
                            'estado',
                        )
                        .order_by('-id')
                    ),
                    to_attr='matriculas_usuario_precargadas',
                )
            )
            .order_by('id')
        )

        if not es_admin(user):
            return queryset.filter(id=user.id)

        rol_param = str(
            self.request.query_params.get('rol') or ''
        ).strip()

        buscar = str(
            self.request.query_params.get('buscar') or ''
        ).strip()[:100]

        if rol_param:
            queryset = queryset.filter(
                rol__nombre__iexact=rol_param
            )

        if buscar:
            queryset = queryset.filter(
                Q(username__icontains=buscar)
                | Q(email__icontains=buscar)
                | Q(first_name__icontains=buscar)
                | Q(last_name__icontains=buscar)
                | Q(estudiante__nombre__icontains=buscar)
                | Q(estudiante__apellido__icontains=buscar)
            )

        return queryset

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede crear usuarios.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar usuarios.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().update(request, *args, **kwargs)

    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede editar usuarios.'},
                status=status.HTTP_403_FORBIDDEN
            )

        return super().partial_update(request, *args, **kwargs)

    @action(detail=False, methods=['post'], url_path='crear-estudiante')
    def crear_usuario_estudiante(self, request):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede crear usuarios de estudiantes.'},
                status=status.HTTP_403_FORBIDDEN
            )

        matricula_id = request.data.get('matricula_id')

        if not matricula_id:
            return Response(
                {'error': 'Debe enviar la matrícula.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            matricula = Matricula.objects.select_related('estudiante').get(id=matricula_id)
        except Matricula.DoesNotExist:
            return Response(
                {'error': 'Matrícula no encontrada.'},
                status=status.HTTP_404_NOT_FOUND
            )

        if matricula.estado == 'finalizado':
            return Response(
                {
                    'error': (
                        'No se puede crear usuario para una '
                        'matrícula finalizada.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        usuario_estudiante = (
            matricula.estudiante.usuarios
            .filter(
                rol__nombre__iexact='estudiante'
            )
            .order_by(
                '-id'
            )
            .first()
        )

        if usuario_estudiante:
            if usuario_estudiante.is_active:
                return Response(
                    {
                        'error': (
                            'Este estudiante ya tiene un usuario activo.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            password = request.data.get('password')

            if password:
                usuario_estudiante.set_password(password)

            usuario_estudiante.is_active = True

            if not usuario_estudiante.estudiante_id:
                usuario_estudiante.estudiante = matricula.estudiante

            usuario_estudiante.save()

            if not matricula.estudiante.activo:
                matricula.estudiante.activo = True
                matricula.estudiante.save(
                    update_fields=[
                        'activo',
                    ]
                )

            Token.objects.filter(
                user=usuario_estudiante
            ).delete()

            return Response(
                self.get_serializer(usuario_estudiante).data,
                status=status.HTTP_200_OK
            )

        data = request.data.copy()
        data['matricula_id'] = matricula.id
        data.setdefault('first_name', matricula.estudiante.nombre)
        data.setdefault('last_name', matricula.estudiante.apellido)
        correo_estudiante = str(
            matricula.estudiante.correo_electronico
            or ''
        ).strip().lower()

        if correo_estudiante:
            data['email'] = correo_estudiante
        else:
            data.pop('email', None)

        serializer = self.get_serializer(data=data)
        serializer.is_valid(raise_exception=True)
        usuario = serializer.save()

        return Response(
            self.get_serializer(usuario).data,
            status=status.HTTP_201_CREATED
        )


class CalendarioViewSet(viewsets.ModelViewSet):
    queryset = (
        Calendario.objects
        .select_related(
            'matricula',
            'matricula__estudiante',
            'matricula__categoria',
            'instructor',
        )
        .defer(
            'instructor__foto_base64'
        )
    )
    serializer_class = CalendarioSerializer
    filter_backends = [DjangoFilterBackend]
    filterset_fields = ['instructor', 'fecha']
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user
        qs = self.queryset

        mes = self.request.query_params.get('mes')
        instructor_param = self.request.query_params.get('instructor')

        if mes:
            mes = mes.strip()

            if not re.fullmatch(
                r'\d{4}-(0[1-9]|1[0-2])',
                mes,
            ):
                raise serializers.ValidationError({
                    'mes': (
                        'El mes debe tener el formato '
                        'AAAA-MM, por ejemplo 2026-07.'
                    )
                })

            anio, numero_mes = map(int, mes.split('-'))
            inicio = date(anio, numero_mes, 1)

            if numero_mes == 12:
                fin = date(anio + 1, 1, 1)
            else:
                fin = date(anio, numero_mes + 1, 1)

            qs = qs.filter(
                fecha__gte=inicio,
                fecha__lt=fin,
            )

        if instructor_param and instructor_param != 'all':
            try:
                instructor_id = int(instructor_param)
            except (TypeError, ValueError):
                raise serializers.ValidationError({
                    'instructor': (
                        'El instructor debe ser un identificador válido.'
                    )
                })

            if instructor_id <= 0:
                raise serializers.ValidationError({
                    'instructor': (
                        'El instructor debe ser un identificador válido.'
                    )
                })

            qs = qs.filter(instructor_id=instructor_id)

        if es_admin(user):
            return qs.order_by(
                'fecha',
                'hora_inicio'
            )

        if getattr(user, 'instructor_id', None):
            return qs.filter(instructor_id=user.instructor_id).order_by('fecha', 'hora_inicio')

        if getattr(user, 'estudiante_id', None):
            return qs.filter(matricula__estudiante_id=user.estudiante_id).order_by('fecha', 'hora_inicio')

        return qs.none()

    @transaction.atomic
    def crear_clases_regulares_seguras(
        self,
        matricula_id,
        instructor_id,
        fechas,
        duraciones_clases,
        hora_inicio,
    ):
        matricula = (
            Matricula.objects
            .select_for_update()
            .select_related(
                'estudiante',
                'categoria',
            )
            .get(id=matricula_id)
        )

        instructor = (
            Instructor.objects
            .select_for_update()
            .get(id=instructor_id)
        )

        if not instructor.activo:
            raise serializers.ValidationError({
                'error': (
                    'El instructor está inactivo y no '
                    'puede recibir nuevas clases.'
                )
            })

        calendario_existente = (
            Calendario.objects
            .filter(
                matricula=matricula,
                es_examen=False,
            )
            .exclude(
                estado='cancelada'
            )
            .exists()
        )

        if calendario_existente:
            raise serializers.ValidationError({
                'error': (
                    'La matrícula ya tiene un calendario '
                    'regular asignado.'
                )
            })

        planes = []

        for fecha_clase, duracion_clase in zip(
            fechas,
            duraciones_clases,
        ):
            hora_fin_clase = (
                datetime.combine(
                    date.today(),
                    hora_inicio,
                )
                + timedelta(
                    hours=duracion_clase
                )
            ).time()

            choque = (
                Calendario.objects
                .filter(
                    instructor=instructor,
                    fecha=fecha_clase,
                    hora_inicio__lt=hora_fin_clase,
                    hora_fin__gt=hora_inicio,
                )
                .exclude(
                    estado='cancelada'
                )
                .exists()
            )

            if choque:
                raise serializers.ValidationError({
                    'error': (
                        'El instructor ya tiene ocupado '
                        f'el horario '
                        f'{hora_inicio.strftime("%H:%M")} - '
                        f'{hora_fin_clase.strftime("%H:%M")} '
                        f'el día '
                        f'{fecha_clase.strftime("%d/%m/%Y")}.'
                    )
                })

            planes.append({
                'fecha': fecha_clase,
                'hora_fin': hora_fin_clase,
            })

        creadas = []

        for numero_clase, plan in enumerate(
            planes,
            start=1,
        ):
            clase = Calendario.objects.create(
                matricula=matricula,
                instructor=instructor,
                fecha=plan['fecha'],
                hora_inicio=hora_inicio,
                hora_fin=plan['hora_fin'],
                numero_clase=numero_clase,
                estado='pendiente',
                es_examen=False,
            )

            creadas.append(clase)

        if matricula_usa_checks(matricula):
            ProgresoTema.objects.filter(
                matricula=matricula,
                completado=False,
            ).update(
                desbloqueado=False
            )

        return creadas

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'crear calendarios.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        return Response(
            {
                'error': (
                    'No se permite crear encuentros directamente '
                    'desde esta ruta. Utilice la creación por bloque '
                    'o la creación manual del calendario.'
                )
            },
            status=status.HTTP_405_METHOD_NOT_ALLOWED
        )

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede '
                        'modificar el calendario.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        return Response(
            {
                'error': (
                    'La edición completa mediante PUT no está '
                    'permitida. Utilice PATCH para modificar '
                    'fecha, horario o instructor.'
                )
            },
            status=status.HTTP_405_METHOD_NOT_ALLOWED
        )

    @action(detail=False, methods=['get'], url_path='hoy')
    def citas_hoy(self, request):
        hoy = date.today()
        citas = self.get_queryset().filter(fecha=hoy)

        serializer = self.get_serializer(citas, many=True)

        return Response({
            'results': serializer.data,
            'count': citas.count(),
            'fecha': hoy.isoformat(),
        })

    @action(detail=True, methods=['post'], url_path='resultado-examen')
    def resultado_examen(self, request, pk=None):
        calendario = self.get_object()
        user = request.user

        if not calendario.es_examen:
            return Response(
                {
                    'error': (
                        'Esta cita no corresponde a un '
                        'examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        tiene_permiso = (
            es_admin(user)
            or (
                es_instructor(user)
                and calendario.instructor_id == user.instructor_id
            )
        )

        if not tiene_permiso:
            return Response(
                {
                    'error': (
                        'No tienes permiso para registrar '
                        'el resultado de este examen.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        resultado = str(
            request.data.get('resultado') or ''
        ).strip().lower()

        if resultado not in ['asistieron', 'cancelado']:
            return Response(
                {
                    'error': (
                        'Resultado no válido. Debe enviar '
                        '"asistieron" o "cancelado".'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if calendario.estado != 'pendiente':
            return Response(
                {
                    'error': (
                        'Este examen ya fue procesado '
                        'y no puede modificarse nuevamente.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        with transaction.atomic():
            if resultado == 'cancelado':
                calendario.estado = 'cancelada'
                calendario.save(
                    update_fields=[
                        'estado',
                    ]
                )

                mensaje = (
                    'El examen policial fue cancelado. '
                    'Este resultado ya no puede modificarse. '
                    'El estudiante puede volver a ser programado '
                    'si todavía no alcanzó las 3 asignaciones.'
                )

            else:
                calendario.estado = 'completada'
                calendario.save(
                    update_fields=[
                        'estado',
                    ]
                )

                matricula = calendario.matricula

                desactivar_usuarios_estudiante(
                    matricula.estudiante
                )

                mensaje = (
                    'El examen policial fue marcado como asistido. '
                    'El acceso del estudiante fue desactivado '
                    'correctamente.'
                )

        calendario.refresh_from_db()

        return Response({
            'message': mensaje,
            'calendario': CalendarioSerializer(
                calendario
            ).data,
        })

    @action(detail=False, methods=['post'], url_path='crear-bloque')
    def crear_bloque_citas(self, request):
        serializer = CrearBloqueCitasSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        data = serializer.validated_data

        matricula = Matricula.objects.select_related(
            'estudiante',
            'categoria',
        ).get(id=data['matricula_id'])

        try:
            instructor = Instructor.objects.get(
                id=data['instructor_id'],
                activo=True,
            )
        except Instructor.DoesNotExist:
            return Response(
                {
                    'error': (
                        'El instructor seleccionado no existe '
                        'o se encuentra desactivado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        rango = obtener_rango_horario(matricula)

        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede crear bloques de citas.'},
                status=status.HTTP_403_FORBIDDEN
            )


        if not rango:
            return Response(
                {'error': 'La matrícula no tiene un horario válido.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        horas_por_dia = int(data.get('horas_por_dia', 2))

        if horas_por_dia <= 0:
            return Response(
                {'error': 'Las horas por día deben ser mayores a cero.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        hora_inicio = datetime.strptime(
            rango[0],
            '%H:%M'
        ).time()

        if matricula.tipo_curso in ['Intermedio', 'Avanzado']:
            horas_totales = int(
                matricula.horas_reforzamiento or 0
            )

            if horas_totales <= 0:
                return Response(
                    {
                        'error': (
                            'La matrícula no tiene horas asignadas '
                            'para este curso.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if matricula.incluye_examen_policial:
                horas_totales -= 2

        else:
            # Principiante continúa exactamente como está.
            horas_totales = 16

        num_clases = int(horas_totales) // horas_por_dia

        if int(horas_totales) % horas_por_dia != 0:
            num_clases += 1

        # Duración individual de cada encuentro. Permite que el último encuentro dure menos horas cuando las horas totales no son divisibles exactamente.
        duraciones_clases = []
        horas_restantes = int(horas_totales)

        for _ in range(num_clases):
            duracion_clase = min(
                horas_por_dia,
                horas_restantes
            )

            duraciones_clases.append(duracion_clase)
            horas_restantes -= duracion_clase

        fechas = []
        actual = data['fecha_inicio']
        es_extraordinario = str(matricula.modalidad).lower() == 'extraordinario'

        while len(fechas) < num_clases:
            es_fin_semana = actual.weekday() >= 5

            if es_extraordinario and es_fin_semana:
                fechas.append(actual)

            if not es_extraordinario and not es_fin_semana:
                fechas.append(actual)

            actual += timedelta(days=1)

        creadas = (
            self.crear_clases_regulares_seguras(
                matricula_id=matricula.id,
                instructor_id=instructor.id,
                fechas=fechas,
                duraciones_clases=duraciones_clases,
                hora_inicio=hora_inicio,
            )
        )

        return Response(
            {
                'message': (
                    f'Bloque de {len(creadas)} clases '
                    f'creado correctamente.'
                ),
                'fecha_inicio': creadas[0].fecha if creadas else None,
                'fecha_fin': creadas[-1].fecha if creadas else None,
                'clases_creadas': len(creadas),
                'horas_totales': horas_totales,
                'hora_inicio': hora_inicio.strftime('%H:%M'),
                'hora_fin_ultima_clase': (
                    creadas[-1].hora_fin.strftime('%H:%M')
                    if creadas
                    else None
                ),
                'citas': CalendarioSerializer(
                    creadas,
                    many=True
                ).data,
            },
            status=status.HTTP_201_CREATED
        )

    @action(detail=False, methods=['post'], url_path='crear-manual')
    def crear_calendario_manual(self, request):
        if not es_admin(request.user):
            return Response(
                {'error': 'Solo el administrador puede crear calendario manual.'},
                status=status.HTTP_403_FORBIDDEN
            )

        serializer = CrearCalendarioManualSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        data = serializer.validated_data

        matricula = Matricula.objects.select_related(
            'estudiante',
            'categoria',
        ).get(id=data['matricula_id'])

        try:
            instructor = Instructor.objects.get(
                id=data['instructor_id'],
                activo=True,
            )
        except Instructor.DoesNotExist:
            return Response(
                {
                    'error': (
                        'El instructor seleccionado no existe '
                        'o se encuentra desactivado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        rango = obtener_rango_horario(matricula)

        if not rango:
            return Response(
                {'error': 'La matrícula no tiene un horario válido.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        horas_por_dia = int(data.get('horas_por_dia', 2))

        if horas_por_dia <= 0:
            return Response(
                {'error': 'Las horas por día deben ser mayores a cero.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        hora_inicio = datetime.strptime(
            rango[0],
            '%H:%M'
        ).time()

        fechas = sorted(
            data['fechas']
        )

        if matricula.tipo_curso in [
            'Intermedio',
            'Avanzado',
        ]:
            horas_totales = int(
                matricula.horas_reforzamiento or 0
            )

            if horas_totales <= 0:
                return Response(
                    {
                        'error': (
                            'La matrícula no tiene horas asignadas '
                            'para este curso.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if matricula.incluye_examen_policial:
                horas_totales -= 2

        else:
            # Principiante mantiene 16 horas operativas en calendario.
            horas_totales = 16

        if horas_totales <= 0:
            return Response(
                {
                    'error': (
                        'Las horas de práctica deben ser mayores a cero.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        duraciones_clases = []
        horas_restantes = int(
            horas_totales
        )

        for _ in fechas:
            if horas_restantes <= 0:
                break

            duracion_clase = min(
                horas_por_dia,
                horas_restantes
            )

            duraciones_clases.append(
                duracion_clase
            )

            horas_restantes -= duracion_clase

        if horas_restantes > 0:
            return Response(
                {
                    'error': (
                        f'Las fechas seleccionadas no cubren las {horas_totales} '
                        f'horas requeridas. Selecciona más fechas.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if len(fechas) > len(duraciones_clases):
            return Response(
                {
                    'error': (
                        f'Seleccionaste {len(fechas)} fechas, pero solo se necesitan '
                        f'{len(duraciones_clases)} clases para cubrir las '
                        f'{horas_totales} horas requeridas.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        creadas = (
            self.crear_clases_regulares_seguras(
                matricula_id=matricula.id,
                instructor_id=instructor.id,
                fechas=fechas,
                duraciones_clases=duraciones_clases,
                hora_inicio=hora_inicio,
            )
        )

        return Response(
            {
                'message': (
                    f'Calendario manual de {len(creadas)} clases '
                    f'creado correctamente.'
                ),
                'fecha_inicio': creadas[0].fecha if creadas else None,
                'fecha_fin': creadas[-1].fecha if creadas else None,
                'clases_creadas': len(creadas),
                'horas_totales': horas_totales,
                'hora_inicio': hora_inicio.strftime('%H:%M'),
                'hora_fin_ultima_clase': (
                    creadas[-1].hora_fin.strftime('%H:%M')
                    if creadas
                    else None
                ),
                'citas': CalendarioSerializer(
                    creadas,
                    many=True
                ).data,
            },
            status=status.HTTP_201_CREATED
        )

    def partial_update(self, request, *args, **kwargs):
        def error(
            mensaje,
            codigo=status.HTTP_400_BAD_REQUEST
        ):
            return Response(
                {'error': mensaje},
                status=codigo
            )

        if not es_admin(request.user):
            return error(
                'Solo el administrador puede modificar '
                'el calendario.',
                status.HTTP_403_FORBIDDEN,
            )

        instance = self.get_object()

        if instance.es_examen:
            return error(
                'Los exámenes policiales no pueden '
                'modificarse desde la edición normal.'
            )

        aplicar_a = str(
            request.data.get(
                'aplicar_a',
                'solo'
            )
            or 'solo'
        ).strip().lower()

        if aplicar_a not in [
            'solo',
            'pendientes',
        ]:
            return error(
                'La opción aplicar_a debe ser '
                '"solo" o "pendientes".'
            )

        instructor_id = request.data.get(
            'instructor'
        )
        nueva_fecha = request.data.get(
            'fecha'
        )
        nueva_hora_inicio = request.data.get(
            'hora_inicio'
        )
        nueva_hora_fin = request.data.get(
            'hora_fin'
        )

        if instructor_id in ['', None]:
            instructor_id = None
        else:
            try:
                instructor_id = int(
                    instructor_id
                )
            except (TypeError, ValueError):
                return error(
                    'El instructor enviado no es válido.'
                )

            if not Instructor.objects.filter(
                id=instructor_id,
                activo=True,
            ).exists():
                return error(
                    (
                        'El instructor seleccionado no existe '
                        'o se encuentra desactivado.'
                    ),
                    status.HTTP_400_BAD_REQUEST,
                )

        fecha_obj = None

        if nueva_fecha not in [None, '']:
            fecha_obj = parse_date(
                str(nueva_fecha)
            )

            if not fecha_obj:
                return error(
                    'La fecha debe tener formato '
                    'YYYY-MM-DD.'
                )

        def convertir_hora(valor):
            if valor in [None, '']:
                return None

            valor = str(valor).strip()

            formato = (
                '%H:%M'
                if len(valor) == 5
                else '%H:%M:%S'
            )

            try:
                return datetime.strptime(
                    valor,
                    formato
                ).time()
            except ValueError:
                return None

        hora_inicio_obj = convertir_hora(
            nueva_hora_inicio
        )
        hora_fin_obj = convertir_hora(
            nueva_hora_fin
        )

        envio_alguna_hora = (
            nueva_hora_inicio not in [None, '']
            or nueva_hora_fin not in [None, '']
        )

        if (
            envio_alguna_hora
            and (
                hora_inicio_obj is None
                or hora_fin_obj is None
            )
        ):
            return error(
                'Debe enviar hora de inicio y hora fin '
                'con formato HH:MM.'
            )

        if (
            envio_alguna_hora
            and hora_fin_obj <= hora_inicio_obj
        ):
            return error(
                'La hora fin debe ser mayor '
                'que la hora inicio.'
            )

        with transaction.atomic():
            Matricula.objects.select_for_update().get(
                id=instance.matricula_id
            )

            ids_instructores = set(
                Calendario.objects.filter(
                    matricula_id=(
                        instance.matricula_id
                    ),
                    es_examen=False,
                    estado='pendiente',
                    numero_clase__gte=(
                        instance.numero_clase
                    ),
                ).values_list(
                    'instructor_id',
                    flat=True
                )
            )

            if instructor_id is not None:
                ids_instructores.add(
                    instructor_id
                )

            list(
                Instructor.objects
                .select_for_update()
                .filter(
                    id__in=sorted(
                        ids_instructores
                    )
                )
                .order_by('id')
            )

            instance = (
                Calendario.objects
                .select_for_update()
                .select_related('matricula')
                .get(id=instance.id)
            )

            if instance.es_examen:
                return error(
                    'Los exámenes policiales no pueden '
                    'modificarse desde la edición normal.'
                )

            if instance.estado != 'pendiente':
                return error(
                    'Solo se pueden modificar encuentros '
                    'que todavía estén pendientes.'
                )

            if Asistencia.objects.filter(
                As_calendario=instance
            ).exists():
                return error(
                    'Este encuentro ya tiene una '
                    'asistencia registrada.'
                )

            cambia_fecha = (
                fecha_obj is not None
                and fecha_obj != instance.fecha
            )

            incluir_pendientes = (
                aplicar_a == 'pendientes'
                or cambia_fecha
            )

            if incluir_pendientes:
                clases = list(
                    Calendario.objects
                    .select_for_update()
                    .filter(
                        matricula=instance.matricula,
                        es_examen=False,
                        estado='pendiente',
                        numero_clase__gte=(
                            instance.numero_clase
                        ),
                    )
                    .order_by(
                        'numero_clase',
                        'fecha',
                        'hora_inicio',
                        'id',
                    )
                )
            else:
                clases = [instance]

            if not clases:
                return error(
                    'No hay encuentros pendientes disponibles para actualizar.'
                )

            ids_clases = [
                clase.id
                for clase in clases
            ]

            if Asistencia.objects.filter(
                As_calendario_id__in=ids_clases
            ).exists():
                return error(
                    'Uno de los encuentros pendientes '
                    'ya tiene asistencia registrada.'
                )

            fechas_destino = {}

            if cambia_fecha:
                modalidad = str(
                    instance.matricula.modalidad
                    or ''
                ).strip().lower()

                def fecha_permitida(fecha):
                    if modalidad == 'regular':
                        return fecha.weekday() < 5

                    if modalidad == 'extraordinario':
                        return fecha.weekday() >= 5

                    # La modalidad Mixto permite
                    # fechas entre semana y fines de semana.
                    return True

                if not fecha_permitida(fecha_obj):
                    if modalidad == 'regular':
                        return error(
                            'La modalidad Regular solamente '
                            'permite clases de lunes a viernes.'
                        )

                    if modalidad == 'extraordinario':
                        return error(
                            'La modalidad Extraordinario solamente '
                            'permite clases los sábados y domingos.'
                        )

                if modalidad in [
                    'regular',
                    'extraordinario',
                ]:
                    fecha_destino = fecha_obj

                    for indice, clase in enumerate(
                        clases
                    ):
                        if indice > 0:
                            fecha_destino += timedelta(
                                days=1
                            )

                            while not fecha_permitida(
                                fecha_destino
                            ):
                                fecha_destino += timedelta(
                                    days=1
                                )

                        fechas_destino[
                            clase.id
                        ] = fecha_destino

                else:
                    # En modalidad Mixto se conserva
                    # la separación original entre las
                    # fechas seleccionadas manualmente.
                    desplazamiento = (
                        fecha_obj
                        - instance.fecha
                    )

                    fechas_destino = {
                        clase.id: (
                            clase.fecha
                            + desplazamiento
                        )
                        for clase in clases
                    }

            planes = []

            for clase in clases:
                aplicar_datos = (
                    aplicar_a == 'pendientes'
                    or clase.id == instance.id
                )

                instructor_destino = (
                    instructor_id
                    if (
                        aplicar_datos
                        and instructor_id is not None
                    )
                    else clase.instructor_id
                )

                if not instructor_destino:
                    return error(
                        'Uno de los encuentros no tiene '
                        'instructor asignado.'
                    )

                planes.append({
                    'clase': clase,
                    'instructor_id': (
                        instructor_destino
                    ),
                    'fecha': fechas_destino.get(
                        clase.id,
                        clase.fecha
                    ),
                    'hora_inicio': (
                        hora_inicio_obj
                        if (
                            aplicar_datos
                            and envio_alguna_hora
                        )
                        else clase.hora_inicio
                    ),
                    'hora_fin': (
                        hora_fin_obj
                        if (
                            aplicar_datos
                            and envio_alguna_hora
                        )
                        else clase.hora_fin
                    ),
                })

            # Evita choques entre los propios
            # encuentros que serán desplazados.
            for indice, plan in enumerate(planes):
                for otro in planes[indice + 1:]:
                    chocan_entre_si = (
                        plan['instructor_id']
                        == otro['instructor_id']
                        and plan['fecha']
                        == otro['fecha']
                        and plan['hora_inicio']
                        < otro['hora_fin']
                        and plan['hora_fin']
                        > otro['hora_inicio']
                    )

                    if chocan_entre_si:
                        return error(
                            'Los cambios producirían dos '
                            'encuentros en el mismo horario '
                            f'el día {plan["fecha"]}.'
                        )

            # Comprueba choques con otras matrículas
            # y con exámenes policiales.
            for plan in planes:
                choque = (
                    Calendario.objects
                    .select_for_update()
                    .filter(
                        instructor_id=(
                            plan['instructor_id']
                        ),
                        fecha=plan['fecha'],
                        hora_inicio__lt=(
                            plan['hora_fin']
                        ),
                        hora_fin__gt=(
                            plan['hora_inicio']
                        ),
                    )
                    .exclude(
                        id__in=ids_clases
                    )
                    .exclude(
                        estado='cancelada'
                    )
                    .exists()
                )

                if choque:
                    inicio = (
                        plan['hora_inicio']
                        .strftime('%H:%M')
                    )
                    fin = (
                        plan['hora_fin']
                        .strftime('%H:%M')
                    )

                    return error(
                        'El instructor ya tiene ocupado '
                        f'el horario {inicio} - {fin} '
                        f'el día {plan["fecha"]}.'
                    )

            clases_actualizadas = 0
            fechas_desplazadas = 0

            for plan in planes:
                clase = plan['clase']
                campos = []

                if (
                    clase.instructor_id
                    != plan['instructor_id']
                ):
                    clase.instructor_id = (
                        plan['instructor_id']
                    )
                    campos.append('instructor')

                if clase.fecha != plan['fecha']:
                    clase.fecha = plan['fecha']
                    campos.append('fecha')
                    fechas_desplazadas += 1

                if (
                    clase.hora_inicio
                    != plan['hora_inicio']
                ):
                    clase.hora_inicio = (
                        plan['hora_inicio']
                    )
                    campos.append('hora_inicio')

                if (
                    clase.hora_fin
                    != plan['hora_fin']
                ):
                    clase.hora_fin = (
                        plan['hora_fin']
                    )
                    campos.append('hora_fin')

                if campos:
                    clase.save(
                        update_fields=campos
                    )
                    clases_actualizadas += 1

            # Conserva el instructor histórico utilizado
            # por el reporte de inducción.
            if (
                instructor_id is not None
                and aplicar_a == 'pendientes'
            ):
                Notas.objects.filter(
                    matricula=instance.matricula,
                    tipo_nota='practico',
                ).exclude(
                    instructor_id=instructor_id
                ).update(
                    instructor_id=instructor_id
                )

        instance.refresh_from_db()

        return Response({
            'message': (
                'Calendario actualizado correctamente.'
            ),
            'clases_actualizadas': (
                clases_actualizadas
            ),
            'fechas_desplazadas': (
                fechas_desplazadas
            ),
            'calendario': self.get_serializer(
                instance
            ).data,
        })

    @action(detail=False, methods=['post'], url_path='crear-examen')
    def crear_examen(self, request):
        user = request.user

        if not es_instructor(user):
            return Response(
                {
                    'error': 'Solo un instructor puede programar un examen policial.'
                },
                status=status.HTTP_403_FORBIDDEN
            )

        matricula_id = request.data.get('matricula_id')
        fecha = request.data.get('fecha')

        if not matricula_id:
            return Response(
                {
                    'error': 'Debe seleccionar un estudiante.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if not fecha:
            return Response(
                {
                    'error': (
                        'Debe seleccionar la fecha del examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        fecha_examen = parse_date(str(fecha))

        if not fecha_examen:
            return Response(
                {
                    'error': (
                        'La fecha del examen no tiene un formato válido.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        dia_semana = fecha_examen.weekday()

        # Lunes a viernes: 2:00 p. m. a 4:00 p. m.
        if dia_semana in [0, 1, 2, 3, 4]:
            horario_examen = '14_16'
            hora_inicio = datetime.strptime(
                '14:00',
                '%H:%M'
            ).time()
            hora_fin = datetime.strptime(
                '16:00',
                '%H:%M'
            ).time()

        # Sábado: 8:00 a. m. a 10:00 a. m.
        elif dia_semana == 5:
            horario_examen = '08_10'
            hora_inicio = datetime.strptime(
                '08:00',
                '%H:%M'
            ).time()
            hora_fin = datetime.strptime(
                '10:00',
                '%H:%M'
            ).time()

        else:
            return Response(
                {
                    'error': (
                        'Los exámenes policiales no pueden programarse los domingos.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            matricula = Matricula.objects.select_related(
                'estudiante'
            ).get(
                id=matricula_id
            )

        except Matricula.DoesNotExist:
            return Response(
                {
                    'error': 'Matrícula no encontrada.'
                },
                status=status.HTTP_404_NOT_FOUND
            )

        try:
            instructor = Instructor.objects.get(
                id=user.instructor_id
            )

        except Instructor.DoesNotExist:
            return Response(
                {
                    'error': 'Instructor no encontrado.'
                },
                status=status.HTTP_404_NOT_FOUND
            )

        if matricula.estado not in [
            'matriculado',
            'finalizado',
        ]:
            return Response(
                {
                    'error': (
                        'No se puede programar el examen porque la matrícula no está aprobada.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if not matricula.estudiante.activo:
            return Response(
                {
                    'error': (
                        'El estudiante se encuentra inactivo.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        tipo_curso = str(
            matricula.tipo_curso or ''
        ).strip().lower()

        if tipo_curso not in [
            'principiante',
            'intermedio',
            'avanzado',
        ]:
            return Response(
                {
                    'error': (
                        'El tipo de curso de la matrícula no es válido.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if (
            tipo_curso in ['intermedio', 'avanzado']
            and not matricula.incluye_examen_policial
        ):
            return Response(
                {
                    'error': (
                        'Esta matrícula no incluye acompañamiento al examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        usuario_activo = matricula.estudiante.usuarios.filter(
            rol__nombre__iexact='estudiante',
            is_active=True,
        ).exists()

        if not usuario_activo:
            return Response(
                {
                    'error': (
                        'No se puede programar el examen porque el estudiante no tiene un usuario activo.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        clases_regulares = Calendario.objects.filter(
            matricula=matricula,
            es_examen=False,
        ).exclude(
            estado='cancelada'
        )

        if not clases_regulares.exists():
            return Response(
                {
                    'error': (
                        'El estudiante no tiene clases regulares asignadas.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        esta_asignado_al_instructor = (
            clases_regulares.filter(
                instructor=instructor
            ).exists()
        )

        if not esta_asignado_al_instructor:
            return Response(
                {
                    'error': (
                        'Este estudiante no está asignado al instructor actual.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        tiene_clases_pendientes = clases_regulares.exclude(
            estado='completada'
        ).exists()

        if tiene_clases_pendientes:
            return Response(
                {
                    'error': (
                        'El estudiante todavía no ha completado todas sus clases.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if not matricula_tiene_notas_teorica_y_practica(
            matricula
        ):
            return Response(
                {
                    'error': (
                        'El estudiante necesita tener registrada '
                        'la nota teórica y la nota práctica antes de programar el examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        examenes_del_estudiante = Calendario.objects.filter(
            matricula=matricula,
            es_examen=True,
        )

        total_examenes_asignados = (
            examenes_del_estudiante.count()
        )

        # Los exámenes cancelados también cuentan dentro de las tres asignaciones.
        if total_examenes_asignados >= 3:
            return Response(
                {
                    'error': (
                        'El estudiante ya alcanzó el máximo de 3 asignaciones para examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        tiene_examen_vigente = (
            examenes_del_estudiante.exclude(
                estado='cancelada'
            ).exists()
        )

        if tiene_examen_vigente:
            return Response(
                {
                    'error': (
                        'El estudiante ya tiene un examen policial pendiente o completado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        clases_reprogramadas = []
        nuevas_clases_recuperacion = []

        with transaction.atomic():
            matricula = (
                Matricula.objects
                .select_for_update()
                .select_related('estudiante')
                .get(id=matricula.id)
            )

            instructor = (
                Instructor.objects
                .select_for_update()
                .get(id=instructor.id)
            )

            examenes_del_estudiante = (
                Calendario.objects
                .filter(
                    matricula=matricula,
                    es_examen=True,
                )
            )

            total_examenes_asignados = (
                examenes_del_estudiante.count()
            )

            if total_examenes_asignados >= 3:
                return Response(
                    {
                        'error': (
                            'El estudiante ya alcanzó el máximo de 3 asignaciones para examen policial.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if (
                examenes_del_estudiante
                .exclude(estado='cancelada')
                .exists()
            ):
                return Response(
                    {
                        'error': (
                            'El estudiante ya tiene un examen policial pendiente o completado.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            examenes_instructor_dia = (
                Calendario.objects
                .filter(
                    instructor=instructor,
                    fecha=fecha_examen,
                    es_examen=True,
                )
                .exclude(
                    estado='cancelada'
                )
                .count()
            )

            if examenes_instructor_dia >= 10:
                return Response(
                    {
                        'error': (
                            'El instructor ya tiene asignados 10 estudiantes para examen policial en esa fecha.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            clases_en_conflicto = list(
                Calendario.objects
                .select_for_update()
                .filter(
                    instructor=instructor,
                    fecha=fecha_examen,
                    es_examen=False,
                    hora_inicio__lt=hora_fin,
                    hora_fin__gt=hora_inicio,
                )
                .exclude(
                    estado='cancelada'
                )
                .order_by(
                    'hora_inicio',
                    'id',
                )
            )

            clases_no_reprogramables = [
                clase
                for clase in clases_en_conflicto
                if clase.estado != 'pendiente'
            ]

            if clases_no_reprogramables:
                return Response(
                    {
                        'error': (
                            'El instructor tiene una clase que ya fue procesada en ese horario '
                            'y no puede reprogramarse automáticamente.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            for clase in clases_en_conflicto:
                Asistencia.objects.update_or_create(
                    As_calendario=clase,
                    defaults={
                        'As_estudiante': (
                            clase.matricula.estudiante
                        ),
                        'estado': 'justificado',
                        'observacion': (
                            'Encuentro justificado '
                            'automáticamente porque el '
                            'instructor fue asignado a un '
                            'examen policial en el mismo horario.'
                        ),
                        'justificado_por_admin': True,
                        'km_inicial': None,
                        'km_final': None,
                    },
                )

                clase.estado = 'reprogramada'
                clase.save(
                    update_fields=['estado']
                )

                nueva_clase = crear_clase_recuperacion(
                    clase
                )

                clases_reprogramadas.append(
                    clase.id
                )

                nuevas_clases_recuperacion.append(
                    nueva_clase.id
                )

            examen = Calendario.objects.create(
                matricula=matricula,
                instructor=instructor,
                fecha=fecha_examen,
                hora_inicio=hora_inicio,
                hora_fin=hora_fin,
                numero_clase=99,
                estado='pendiente',
                es_examen=True,
            )

            return Response(
                {
                    'message': (
                        'Examen policial programado correctamente.'
                        if not clases_reprogramadas
                        else (
                            'Examen policial programado correctamente. '
                            f'Se justificaron y reprogramaron '
                            f'{len(clases_reprogramadas)} '
                            f'encuentro(s) regular(es).'
                        )
                    ),
                    'fecha': fecha_examen,
                    'horario_examen': horario_examen,
                    'hora_inicio': hora_inicio.strftime(
                        '%H:%M'
                    ),
                    'hora_fin': hora_fin.strftime(
                        '%H:%M'
                    ),
                    'cupo_instructor': {
                        'asignados': (
                            examenes_instructor_dia + 1
                        ),
                        'maximo': 10,
                    },
                    'asignaciones_estudiante': {
                        'utilizadas': (
                            total_examenes_asignados + 1
                        ),
                        'maximo': 3,
                    },
                    'reprogramacion_automatica': {
                        'cantidad': len(
                            clases_reprogramadas
                        ),
                        'clases_justificadas': (
                            clases_reprogramadas
                        ),
                        'nuevas_clases': (
                            nuevas_clases_recuperacion
                        ),
                    },
                    'examen': CalendarioSerializer(
                        examen
                    ).data,
                },
                status=status.HTTP_201_CREATED
            )

class AsistenciaViewSet(viewsets.GenericViewSet):
    queryset = Asistencia.objects.select_related(
        'As_estudiante',
        'As_calendario',
        'As_calendario__matricula',
        'As_calendario__matricula__estudiante',
        'As_calendario__instructor',
    ).all()

    serializer_class = AsistenciaSerializer
    permission_classes = [IsAuthenticated]

    http_method_names = [
        'get',
        'post',
        'head',
        'options',
    ]

    def get_queryset(self):
        """
        Limita las asistencias según el usuario autenticado.
        Administración y Secretaría pueden consultar todas.
        El instructor solamente puede consultar sus clases.
        El estudiante solamente puede consultar sus asistencias.
        """
        queryset = super().get_queryset()
        user = self.request.user

        if es_admin(user):
            return queryset

        if es_instructor(user):
            return queryset.filter(
                As_calendario__instructor_id=user.instructor_id
            )

        if es_estudiante(user):
            return queryset.filter(
                As_estudiante_id=user.estudiante_id
            )

        return queryset.none()

    def list(self, request, *args, **kwargs):
        user = request.user
        hoy = timezone.localdate()

        es_admin_asistencia = es_admin(user)
        es_usuario_instructor = es_instructor(user)
        es_usuario_estudiante = es_estudiante(user)
        fecha_param = request.query_params.get('fecha')
        fecha_inicio_param = request.query_params.get('fecha_inicio')
        fecha_fin_param = request.query_params.get('fecha_fin')

        if fecha_param:
            fecha_inicio_param = fecha_param
            fecha_fin_param = fecha_param

        if not fecha_inicio_param:
            fecha_inicio_param = hoy.isoformat()

        if not fecha_fin_param:
            fecha_fin_param = fecha_inicio_param

        try:
            fecha_inicio = datetime.strptime(
                fecha_inicio_param,
                '%Y-%m-%d'
            ).date()

            fecha_fin = datetime.strptime(
                fecha_fin_param,
                '%Y-%m-%d'
            ).date()

        except ValueError:
            return Response(
                {
                    'error': 'Las fechas deben tener el formato YYYY-MM-DD.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if fecha_fin < fecha_inicio:
            return Response(
                {
                    'error': 'La fecha final no puede ser menor que la fecha inicial.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        clases_base = Calendario.objects.select_related(
            'matricula',
            'matricula__estudiante',
            'instructor'
        ).filter(
            es_examen=False,
            matricula__estado='matriculado'
        ).order_by(
            'matricula_id',
            'numero_clase',
            'fecha',
            'hora_inicio'
        )

        if es_usuario_instructor:
            clases_base = clases_base.filter(
                instructor_id=user.instructor_id
            )

        elif es_usuario_estudiante:
            clases_base = clases_base.filter(
                matricula__estudiante_id=user.estudiante_id
            )

        elif not es_admin_asistencia:
            return Response([])

        clases = clases_base.filter(
            fecha__gte=fecha_inicio,
            fecha__lte=fecha_fin
        ).order_by(
            'matricula_id',
            'numero_clase',
            'fecha',
            'hora_inicio'
        )

        asistencias = Asistencia.objects.select_related(
            'As_estudiante',
            'As_calendario'
        ).filter(
            As_calendario__in=clases
        )

        asistencias_por_clase = {
            asistencia.As_calendario_id: asistencia
            for asistencia in asistencias
        }

        resultado = {}

        for clase in clases:
            estudiante = clase.matricula.estudiante
            matricula_id = clase.matricula_id
            instructor = clase.instructor

            instructor_nombre = 'No asignado'

            if instructor:
                instructor_nombre = (
                    f"{instructor.nombre or ''} "
                    f"{instructor.apellido or ''}"
                ).strip()

                if not instructor_nombre:
                    instructor_nombre = (
                        f'Instructor {instructor.id}'
                    )

            if matricula_id not in resultado:
                resultado[matricula_id] = {
                    'matricula_id': matricula_id,
                    'nombre': f'{estudiante.nombre} {estudiante.apellido}',
                    'cedula': estudiante.cedula,
                    'tipo_curso': clase.matricula.tipo_curso,
                    'instructor_id': instructor.id if instructor else None,
                    'instructor_nombre': instructor_nombre,
                    'conductor': instructor_nombre,
                    'asistencias': {},
                    'porcentaje': 0,
                }

            asistencia = asistencias_por_clase.get(clase.id)

            if asistencia:
                estado = asistencia.estado
                asistencia_id = asistencia.id
                justificado_por_admin = asistencia.justificado_por_admin
            else:
                estado = 'pendiente'
                asistencia_id = None
                justificado_por_admin = False

            en_rango = fecha_inicio <= clase.fecha <= fecha_fin
            es_hoy = clase.fecha == hoy
            es_pasado = clase.fecha < hoy
            es_futuro = clase.fecha > hoy

            puede_marcar = (
                asistencia is None
                and clase.estado in ['pendiente', 'reprogramada']
                and (
                    (
                        es_admin_asistencia
                        and clase.fecha <= hoy
                    )
                    or (
                        es_usuario_instructor
                        and en_rango
                        and es_hoy
                        and clase.instructor_id == user.instructor_id
                    )
                )
            )

            resultado[matricula_id]['asistencias'][str(clase.numero_clase)] = {
                'id': clase.id,
                'asistencia_id': asistencia_id,
                'fecha': clase.fecha,
                'hora_inicio': clase.hora_inicio,
                'hora_fin': clase.hora_fin,
                'numero_clase': clase.numero_clase,
                'estado': estado,
                'justificado_por_admin': justificado_por_admin,
                'instructor_id': clase.instructor.id if clase.instructor else None,
                'instructor_nombre': instructor_nombre,
                'km_inicial': asistencia.km_inicial if asistencia else None,
                'km_final': asistencia.km_final if asistencia else None,
                'km_recorridos': asistencia.km_recorridos if asistencia else 0,
                'en_rango': en_rango,
                'es_hoy': es_hoy,
                'es_pasado': es_pasado,
                'es_futuro': es_futuro,
                'puede_marcar': puede_marcar,
                'bloqueado': not puede_marcar,
            }

        for item in resultado.values():
            asistencias_estudiante = item['asistencias'].values()
            total_clases_validas = 0
            total_asistidas = 0

            for asistencia in asistencias_estudiante:
                estado = asistencia.get('estado')

                if estado == 'justificado':
                    continue

                total_clases_validas += 1

                if estado == 'asistio':
                    total_asistidas += 1

            item['porcentaje'] = (
                round((total_asistidas / total_clases_validas) * 100)
                if total_clases_validas > 0
                else 0
            )
        return Response(list(resultado.values()))

    @action(detail=False, methods=['post'], url_path='marcar')
    def marcar(self, request):
        clase_id = request.data.get('clase_id')
        estado = request.data.get('estado')
        observacion = request.data.get('observacion', '')
        km_inicial = request.data.get('km_inicial')
        km_final = None
        user = request.user
        es_admin_asistencia = es_admin(user)
        es_usuario_instructor = es_instructor(user)

        if not es_admin_asistencia and not es_usuario_instructor:
            return Response(
                {'error': 'Solo el instructor o el administrador pueden marcar asistencia.'},
                status=status.HTTP_403_FORBIDDEN
            )

        if estado not in ['asistio', 'falto']:
            return Response(
                {'error': 'Solo se puede marcar asistió o faltó.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            clase = Calendario.objects.select_related(
                'matricula',
                'matricula__estudiante',
                'instructor'
            ).get(id=clase_id)
        except Calendario.DoesNotExist:
            return Response(
                {'error': 'Clase no encontrada.'},
                status=status.HTTP_404_NOT_FOUND
            )

        if clase.es_examen:
            return Response(
                {
                    'error': (
                        'Los exámenes policiales no se pueden registrar como asistencia normal. '
                        'Debe utilizarse la opción de resultado del examen policial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if es_usuario_instructor and not es_admin_asistencia:
            if clase.instructor_id != user.instructor_id:
                return Response(
                    {
                        'error': (
                            'No puedes marcar asistencia de una clase que no te pertenece.'
                        )
                    },
                    status=status.HTTP_403_FORBIDDEN
                )

        if estado == 'asistio':
            if km_inicial in [None, '']:
                    return Response(
                    {'error': 'Debe ingresar el km inicial.'},
                    status=status.HTTP_400_BAD_REQUEST
                )

            try:
                km_inicial = validar_kilometraje(
                    km_inicial,
                    'km inicial'
                )
            except ValueError as error:
                return Response(
                    {'error': str(error)},
                    status=status.HTTP_400_BAD_REQUEST
                )

            km_final = None

            if not observacion and es_admin_asistencia:
                observacion = 'Asistencia registrada por administrador.'
        else:
            km_inicial = None
            km_final = None

        hoy = timezone.localdate()

        if es_admin_asistencia:
            if clase.fecha > hoy:
                return Response(
                    {'error': 'No se puede marcar asistencia de una clase futura.'},
                    status=status.HTTP_400_BAD_REQUEST
                )
        else:
            if clase.fecha != hoy:
                return Response(
                    {'error': 'Solo se puede marcar asistencia el día exacto de la clase.'},
                    status=status.HTTP_400_BAD_REQUEST
                )

        if clase.estado not in ['pendiente', 'reprogramada']:
            return Response(
                {'error': 'Esta clase ya no está disponible para marcar asistencia.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        asistencia, created = Asistencia.objects.update_or_create(
            As_calendario=clase,
            defaults={
                'As_estudiante': clase.matricula.estudiante,
                'estado': estado,
                'observacion': observacion,
                'justificado_por_admin': False,
                'km_inicial': km_inicial,
                'km_final': km_final,
            }
        )

        if estado == 'asistio':
            clase.estado = 'completada'
        else:
            clase.estado = 'inasistencia'

        clase.save(update_fields=['estado'])
        serializer = self.get_serializer(asistencia)

        return Response({
            'success': True,
            'message': 'Asistencia registrada correctamente.',
            'data': serializer.data
        })

    @action(detail=False, methods=['post'], url_path='finalizar-km')
    def finalizar_km(self, request):
        asistencia_id = request.data.get('asistencia_id')
        km_final = request.data.get('km_final')
        user = request.user
        es_admin_asistencia = es_admin(user)
        es_usuario_instructor = es_instructor(user)

        if not es_admin_asistencia and not es_usuario_instructor:
            return Response(
                {'error': 'Solo el instructor o el administrador pueden finalizar kilometraje.'},
                status=status.HTTP_403_FORBIDDEN
            )

        if not asistencia_id:
            return Response(
                {'error': 'Debe enviar la asistencia.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            asistencia = Asistencia.objects.select_related(
                'As_calendario',
                'As_calendario__instructor'
            ).get(id=asistencia_id)
        except Asistencia.DoesNotExist:
            return Response(
                {'error': 'Asistencia no encontrada.'},
                status=status.HTTP_404_NOT_FOUND
            )

        if es_usuario_instructor and not es_admin_asistencia:
            if (
                asistencia.As_calendario.instructor_id
                != user.instructor_id
            ):
                return Response(
                    {
                        'error': (
                            'No puedes finalizar kilometraje de una clase que no te pertenece.'
                        )
                    },
                    status=status.HTTP_403_FORBIDDEN
                )

        if asistencia.estado != 'asistio':
            return Response(
                {'error': 'Solo clases asistidas pueden finalizar kilometraje.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        if km_final in [None, '']:
            return Response(
                {'error': 'Debe ingresar el km final.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            km_final = validar_kilometraje(
                km_final,
                'km final'
            )
        except ValueError as error:
            return Response(
                {'error': str(error)},
                status=status.HTTP_400_BAD_REQUEST
            )

        if asistencia.km_inicial is None:
                return Response(
                    {'error': 'La asistencia no tiene km inicial.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        if km_final < asistencia.km_inicial:
            return Response(
                {'error': 'El km final no puede ser menor al inicial.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        asistencia.km_final = km_final
        asistencia.save(
            update_fields=[
                'km_final',
                'km_recorridos',
            ]
        )

        serializer = self.get_serializer(asistencia)

        return Response({
            'success': True,
            'message': 'Kilometraje final registrado.',
            'data': serializer.data
        })

    @action(detail=False, methods=['post'], url_path='editar-km')
    def editar_km(self, request):
        asistencia_id = request.data.get('asistencia_id')
        km_inicial = request.data.get('km_inicial')
        km_final = request.data.get('km_final')
        user = request.user
        es_admin_asistencia = es_admin(user)
        es_usuario_instructor = es_instructor(user)

        if not es_admin_asistencia and not es_usuario_instructor:
            return Response(
                {'error': 'Solo el instructor o el administrador pueden editar kilometraje.'},
                status=status.HTTP_403_FORBIDDEN
            )

        if not asistencia_id:
            return Response(
                {'error': 'Debe enviar la asistencia.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            asistencia = Asistencia.objects.select_related(
                'As_calendario',
                'As_calendario__instructor'
            ).get(id=asistencia_id)
        except Asistencia.DoesNotExist:
            return Response(
                {'error': 'Asistencia no encontrada.'},
                status=status.HTTP_404_NOT_FOUND
            )

        if es_usuario_instructor and not es_admin_asistencia:
            if (
                asistencia.As_calendario.instructor_id
                != user.instructor_id
            ):
                return Response(
                    {
                        'error': (
                            'No puedes editar kilometraje de una '
                            'clase que no te pertenece.'
                        )
                    },
                    status=status.HTTP_403_FORBIDDEN
                )

        if asistencia.estado != 'asistio':
            return Response(
                {'error': 'Solo se puede editar kilometraje de clases asistidas.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        if km_inicial in [None, '']:
            return Response(
                {
                    'error': (
                        'Debe ingresar el km inicial.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            km_inicial = validar_kilometraje(
                km_inicial,
                'km inicial'
            )
        except ValueError as error:
            return Response(
                {'error': str(error)},
                status=status.HTTP_400_BAD_REQUEST
            )

        if km_final in [None, '']:
            km_final = None
        else:
            try:
                km_final = validar_kilometraje(
                    km_final,
                    'km final'
                )
            except ValueError as error:
                return Response(
                    {'error': str(error)},
                    status=status.HTTP_400_BAD_REQUEST
                )

            if km_final < km_inicial:
                return Response(
                    {
                        'error': (
                            'El km final no puede ser menor al inicial.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

        asistencia.km_inicial = km_inicial
        asistencia.km_final = km_final
        asistencia.save(
            update_fields=[
                'km_inicial',
                'km_final',
                'km_recorridos',
            ]
        )

        serializer = self.get_serializer(asistencia)

        return Response({
            'success': True,
            'message': 'Kilometraje editado correctamente.',
            'data': serializer.data
        })

    @action(detail=True, methods=['post'], url_path='justificar')
    def justificar(self, request, pk=None):
        user = request.user

        if not es_admin(user):
            return Response(
                {
                    'error': (
                        'Solo Administración o Secretaría pueden justificar una inasistencia.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        asistencia = self.get_object()

        if asistencia.estado != 'falto':
            return Response(
                {'error': 'Solo se pueden justificar clases marcadas como faltó.'},
                status=status.HTTP_400_BAD_REQUEST
            )

        observacion = str(
            request.data.get('observacion') or ''
        ).strip()

        with transaction.atomic():
            asistencia.estado = 'justificado'
            asistencia.justificado_por_admin = True

            campos_actualizados = [
                'estado',
                'justificado_por_admin',
            ]

            if observacion:
                asistencia.observacion = observacion
                campos_actualizados.append('observacion')

            asistencia.save(
                update_fields=campos_actualizados
            )

            clase_faltada = asistencia.As_calendario
            clase_faltada.estado = 'reprogramada'
            clase_faltada.save(update_fields=['estado'])
            nueva_clase = self.reprogramar_clases_por_justificacion(clase_faltada)

        serializer = self.get_serializer(asistencia)

        return Response({
            'success': True,
            'message': 'Inasistencia justificada. Se agregó un día adicional para recuperar la clase perdida.',
            'data': serializer.data,
            'nueva_clase_id': nueva_clase.id if nueva_clase else None,
        })

    def reprogramar_clases_por_justificacion(
        self,
        clase_faltada,
    ):
        return crear_clase_recuperacion(
            clase_faltada
        )

    @action(detail=False, methods=['get'], url_path='resumen')
    def resumen(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo Administración o Secretaría '
                        'pueden consultar este resumen.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        matriculas = (
            Matricula.objects
            .select_related(
                'estudiante',
                'plan_de_estudio',
            )
            .annotate(
                total_clases_resumen=Count(
                    'clases',
                    filter=Q(
                        clases__es_examen=False
                    ),
                    distinct=True,
                ),
                total_marcadas_resumen=Count(
                    'clases__asistencia',
                    filter=(
                        Q(clases__es_examen=False)
                        & ~Q(
                            clases__asistencia__estado=(
                                'justificado'
                            )
                        )
                    ),
                    distinct=True,
                ),
                presentes_resumen=Count(
                    'clases__asistencia',
                    filter=Q(
                        clases__es_examen=False,
                        clases__asistencia__estado='asistio',
                    ),
                    distinct=True,
                ),
            )
            .order_by('-id')
        )

        resultado = []

        for matricula in matriculas:
            total_marcadas = (
                matricula.total_marcadas_resumen
            )

            presentes = (
                matricula.presentes_resumen
            )

            porcentaje = (
                round(
                    (
                        presentes /
                        total_marcadas
                    ) * 100
                )
                if total_marcadas > 0
                else 0
            )

            resultado.append({
                'matricula_id': matricula.id,
                'nombre': (
                    matricula.estudiante.nombre
                ),
                'apellido': (
                    matricula.estudiante.apellido
                ),
                'cedula': (
                    matricula.estudiante.cedula
                ),
                'plan_estudio': (
                    matricula.plan_de_estudio.nombre
                    if matricula.plan_de_estudio
                    else ''
                ),
                'tipo_curso': (
                    matricula.tipo_curso
                ),
                'total_clases': (
                    matricula.total_clases_resumen
                ),
                'porcentaje': porcentaje,
            })

        return Response(resultado)

    @action(detail=False, methods=['get'], url_path='resumen-km')
    def resumen_km(self, request):
        user = request.user

        asistencias = (
            Asistencia.objects
            .select_related(
                'As_estudiante',
                'As_calendario',
                'As_calendario__instructor',
            )
            .defer(
                'As_calendario__instructor__foto_base64',
            )
            .filter(
                estado='asistio',
            )
        )

        fecha_inicio = request.query_params.get('fecha_inicio')
        fecha_fin = request.query_params.get('fecha_fin')

        if fecha_inicio and fecha_fin:
            asistencias = asistencias.filter(
                As_calendario__fecha__range=[fecha_inicio, fecha_fin]
            )
        elif fecha_inicio:
            asistencias = asistencias.filter(
                As_calendario__fecha=fecha_inicio
            )

        if es_admin(user):
            pass

        elif es_instructor(user):
            asistencias = asistencias.filter(
                As_calendario__instructor_id=user.instructor_id
            )

        elif es_estudiante(user):
            asistencias = asistencias.filter(
                As_estudiante_id=user.estudiante_id
            )

        else:
            return Response([])

        resultado = {}

        for asistencia in asistencias.iterator(
            chunk_size=500
        ):
            estudiante = asistencia.As_estudiante
            calendario = asistencia.As_calendario

            if not estudiante or not calendario or not calendario.instructor:
                continue

            instructor = calendario.instructor
            key = f"{estudiante.id}_{instructor.id}"

            if key not in resultado:
                resultado[key] = {
                    'estudiante_id': estudiante.id,
                    'estudiante_nombre': f"{estudiante.nombre} {estudiante.apellido}",
                    'cedula': estudiante.cedula,
                    'instructor_id': instructor.id,
                    'instructor_nombre': f"{instructor.nombre} {instructor.apellido}",
                    'total_clases': 0,
                    'total_km': 0,
                    'detalles': []
                }

            km = float(asistencia.km_recorridos or 0)

            resultado[key]['total_clases'] += 1
            resultado[key]['total_km'] += km

            resultado[key]['detalles'].append({
                'fecha': calendario.fecha,
                'numero_clase': calendario.numero_clase,
                'km_inicial': asistencia.km_inicial,
                'km_final': asistencia.km_final,
                'km_recorridos': asistencia.km_recorridos,
            })

        return Response(list(resultado.values()))

    @action(detail=False, methods=['get'], url_path='resumen-estudiante')
    def resumen_estudiante(self, request):
        user = request.user

        if not es_estudiante(user):
            return Response(
                {
                    'error': (
                        'Este resumen pertenece únicamente al estudiante autenticado.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        if not user.estudiante_id:
            return Response({
                'porcentaje': 0,
                'asistidas': 0,
                'total': 0,
            })

        matricula = Matricula.objects.filter(
            estudiante=user.estudiante,
            estado='matriculado'
        ).order_by('-id').first()

        if not matricula:
            return Response({
                'porcentaje': 0,
                'asistidas': 0,
                'total': 0,
            })

        clases_todas = Calendario.objects.filter(
            matricula=matricula,
            es_examen=False
        ).order_by(
            'numero_clase',
            'fecha',
            'hora_inicio'
        )

        primera_clase = clases_todas.first()

        if not primera_clase:
            return Response({
                'porcentaje': 0,
                'asistidas': 0,
                'total': 0,
            })

        inicio = datetime.combine(date.today(), primera_clase.hora_inicio)
        fin = datetime.combine(date.today(), primera_clase.hora_fin)
        horas_por_dia = int((fin - inicio).total_seconds() // 3600)

        if horas_por_dia <= 0:
            horas_por_dia = 1

        if matricula.tipo_curso == 'Principiante':
            horas_totales = 15
        else:
            horas_totales = matricula.horas_reforzamiento or 0

        total_encuentros_oficiales = math.ceil(
            float(horas_totales) / horas_por_dia
        ) if horas_totales else 0

        clases_oficiales = clases_todas.filter(
            numero_clase__lte=total_encuentros_oficiales
        )

        total = clases_oficiales.count()

        asistidas = Asistencia.objects.filter(
            As_calendario__in=clases_oficiales,
            estado='asistio'
        ).count()

        porcentaje = (
            round((asistidas / total) * 100)
            if total > 0
            else 0
        )

        return Response({
            'porcentaje': porcentaje,
            'asistidas': asistidas,
            'total': total,
        })

    @action(detail=False, methods=['get'], url_path='fechas-disponibles')
    def fechas_disponibles(self, request):
        user = request.user

        clases = Calendario.objects.filter(
            es_examen=False,
            matricula__estado='matriculado'
        ).exclude(
            estado='cancelada'
        )

        if es_instructor(user):
            clases = clases.filter(
                instructor_id=user.instructor_id
            )

        elif es_estudiante(user):
            clases = clases.filter(
                matricula__estudiante_id=user.estudiante_id
            )

        elif not es_admin(user):
            return Response([])

        fechas = clases.order_by('fecha').values_list(
            'fecha',
            flat=True
        ).distinct()

        return Response([
            fecha.isoformat()
            for fecha in fechas
        ])

class NotasViewSet(viewsets.ModelViewSet):
    queryset = Notas.objects.select_related(
        'matricula',
        'matricula__estudiante',
        'instructor',
        'plan_de_estudio',
    ).all()

    serializer_class = NotasSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    # Se permite consultar y crear. No se permite editar una nota mediante PUT o PATCH.
    http_method_names = [
        'get',
        'post',
        'head',
        'options',
    ]

    def get_queryset(self):
        """
        Limita las notas según el usuario autenticado.
        Administración y Secretaría pueden consultar todas.
        El instructor consulta las notas registradas por él.
        El estudiante consulta únicamente sus propias notas.
        """

        queryset = super().get_queryset()
        user = self.request.user

        if es_admin(user):
            return queryset

        if es_instructor(user):
            return queryset.filter(
                instructor_id=user.instructor_id
            )

        if es_estudiante(user):
            return queryset.filter(
                matricula__estudiante_id=user.estudiante_id
            )
        return queryset.none()

    def obtener_ultima_clase_practica(
        self,
        matricula,
    ):
        """
        Obtiene la última clase práctica activa de una matrícula.
        Las clases correspondientes al examen policial no se
        consideran para habilitar la nota práctica.
        Las clases canceladas tampoco se consideran.
        """

        return (
            Calendario.objects
            .select_related('instructor')
            .filter(
                matricula=matricula,
                es_examen=False,
            )
            .exclude(
                estado='cancelada'
            )
            .order_by(
                '-fecha',
                '-hora_fin',
                '-id',
            )
            .first()
        )

    @action(detail=False, methods=['get'], url_path='agrupadas')
    def agrupadas(self, request):
        """
        Pagina por matrícula y mantiene juntas sus notas práctica y
        teórica. La búsqueda, el filtro y el resumen se calculan sobre
        todos los resultados permitidos para el usuario autenticado.
        """

        notas_permitidas = self.get_queryset()
        notas_coincidentes = notas_permitidas

        buscar = request.query_params.get(
            'buscar',
            ''
        ).strip()[:100]

        tipo_curso = request.query_params.get(
            'tipo_curso',
            ''
        ).strip()

        if buscar:
            notas_coincidentes = notas_coincidentes.filter(
                Q(
                    matricula__estudiante__nombre__icontains=buscar
                )
                | Q(
                    matricula__estudiante__apellido__icontains=buscar
                )
                | Q(
                    matricula__estudiante__cedula__icontains=buscar
                )
                | Q(
                    instructor__nombre__icontains=buscar
                )
                | Q(
                    instructor__apellido__icontains=buscar
                )
                | Q(
                    plan_de_estudio__nombre__icontains=buscar
                )
            )

        if tipo_curso and tipo_curso.lower() != 'todas':
            notas_coincidentes = notas_coincidentes.filter(
                matricula__tipo_curso__iexact=tipo_curso
            )

        ids_coincidentes = (
            notas_coincidentes
            .order_by()
            .values('matricula_id')
        )

        matriculas = (
            Matricula.objects
            .select_related('estudiante')
            .filter(
                id__in=Subquery(ids_coincidentes)
            )
            .order_by('-id')
        )

        nota_practica = (
            notas_permitidas
            .filter(
                matricula_id=OuterRef('pk'),
                tipo_nota='practico',
            )
            .order_by(
                '-fecha_registro',
                '-id',
            )
            .values('nota')[:1]
        )

        nota_teorica = (
            notas_permitidas
            .filter(
                matricula_id=OuterRef('pk'),
                tipo_nota='teorico',
            )
            .order_by(
                '-fecha_registro',
                '-id',
            )
            .values('nota')[:1]
        )

        matriculas_resumen = matriculas.annotate(
            nota_practica_num=Cast(
                Subquery(nota_practica),
                FloatField(),
            ),
            nota_teorica_num=Cast(
                Subquery(nota_teorica),
                FloatField(),
            ),
        )

        datos_resumen = matriculas_resumen.aggregate(
            total=Count('id'),
            aprobados=Count(
                'id',
                filter=(
                    Q(nota_practica_num__isnull=False)
                    & Q(nota_teorica_num__isnull=False)
                    & Q(nota_practica_num__gte=80)
                    & Q(nota_teorica_num__gte=80)
                ),
            ),
            reprobados=Count(
                'id',
                filter=(
                    Q(nota_practica_num__isnull=False)
                    & Q(nota_teorica_num__isnull=False)
                    & (
                        Q(nota_practica_num__lt=80)
                        | Q(nota_teorica_num__lt=80)
                    )
                ),
            ),
            pendientes=Count(
                'id',
                filter=(
                    Q(nota_practica_num__isnull=True)
                    | Q(nota_teorica_num__isnull=True)
                ),
            ),
            suma_practica=Sum(
                'nota_practica_num'
            ),
            suma_teorica=Sum(
                'nota_teorica_num'
            ),
            cantidad_practica=Count(
                'nota_practica_num'
            ),
            cantidad_teorica=Count(
                'nota_teorica_num'
            ),
        )

        cantidad_notas = (
            datos_resumen['cantidad_practica']
            + datos_resumen['cantidad_teorica']
        )

        suma_notas = (
            (datos_resumen['suma_practica'] or 0)
            + (datos_resumen['suma_teorica'] or 0)
        )

        resumen = {
            'total': datos_resumen['total'],
            'aprobados': datos_resumen['aprobados'],
            'reprobados': datos_resumen['reprobados'],
            'pendientes': datos_resumen['pendientes'],
            'promedio': (
                f'{suma_notas / cantidad_notas:.1f}'
                if cantidad_notas
                else '0.0'
            ),
        }

        pagina = self.paginate_queryset(
            matriculas
        )

        matriculas_pagina = (
            pagina
            if pagina is not None
            else list(matriculas)
        )

        ids_pagina = [
            matricula.id
            for matricula in matriculas_pagina
        ]

        notas_pagina = (
            notas_permitidas
            .filter(
                matricula_id__in=ids_pagina
            )
            .order_by(
                '-fecha_registro',
                '-id',
            )
        )

        notas_serializadas = self.get_serializer(
            notas_pagina,
            many=True,
        ).data

        agrupadas = {}

        for nota in notas_serializadas:
            matricula_id = nota['matricula']

            if matricula_id not in agrupadas:
                agrupadas[matricula_id] = {
                    **nota,
                    'nota_practica': None,
                    'nota_teorica': None,
                    'comentario_practico': '',
                    'comentario_teorico': '',
                }

            if nota['tipo_nota'] == 'practico':
                agrupadas[matricula_id][
                    'nota_practica'
                ] = nota['nota']

                agrupadas[matricula_id][
                    'comentario_practico'
                ] = nota['comentario'] or ''

            if nota['tipo_nota'] == 'teorico':
                agrupadas[matricula_id][
                    'nota_teorica'
                ] = nota['nota']

                agrupadas[matricula_id][
                    'comentario_teorico'
                ] = nota['comentario'] or ''

        resultados = [
            agrupadas[matricula.id]
            for matricula in matriculas_pagina
            if matricula.id in agrupadas
        ]

        if pagina is None:
            return Response({
                'count': len(resultados),
                'next': None,
                'previous': None,
                'results': resultados,
                'resumen': resumen,
            })

        return Response({
            'count': (
                self.paginator
                .page
                .paginator
                .count
            ),
            'next': (
                self.paginator.get_next_link()
            ),
            'previous': (
                self.paginator.get_previous_link()
            ),
            'results': resultados,
            'resumen': resumen,
        })

    @action(detail=False, methods=['get'], url_path='estudiantes-disponibles')
    def estudiantes_disponibles(self, request):
        """
        Devuelve únicamente estudiantes que:
        1. Están asignados al instructor autenticado.
        2. No tienen todavía nota práctica.
        3. Ya llegaron al último día de clases prácticas.
        4. Tienen al instructor autenticado en la última clase.
        """

        user = request.user

        if not es_instructor(user):
            return Response(
                [],
                status=status.HTTP_200_OK
            )

        hoy = timezone.localdate()

        matriculas_con_nota_practica = (
            Notas.objects
            .filter(
                tipo_nota='practico'
            )
            .values_list(
                'matricula_id',
                flat=True,
            )
        )

        clases_practicas_activas = (
            Calendario.objects
            .filter(
                es_examen=False,
            )
            .exclude(
                estado='cancelada'
            )
            .select_related(
                'instructor'
            )
            .order_by(
                '-fecha',
                '-hora_fin',
                '-id',
            )
        )

        matriculas = (
            Matricula.objects
            .select_related(
                'estudiante',
                'plan_de_estudio',
            )
            .filter(
                estado='matriculado',
                clases__instructor_id=user.instructor_id,
                clases__es_examen=False,
                clases__estado__in=[
                    'pendiente',
                    'completada',
                    'inasistencia',
                    'reprogramada',
                ],
            )
            .exclude(
                id__in=matriculas_con_nota_practica
            )
            .prefetch_related(
                Prefetch(
                    'clases',
                    queryset=clases_practicas_activas,
                    to_attr='clases_practicas_activas',
                )
            )
            .distinct()
            .order_by(
                'estudiante__nombre',
                'estudiante__apellido',
            )
        )

        resultado = []

        for matricula in matriculas:
            clases_practicas = getattr(
                matricula,
                'clases_practicas_activas',
                [],
            )

            if not clases_practicas:
                continue

            ultima_clase = clases_practicas[0]

            # Solo el instructor de la última clase práctica puede registrar la nota.
            if (
                ultima_clase.instructor_id
                != user.instructor_id
            ):
                continue

            # Antes del último día, el estudiante no aparece. En el último día y posteriormente, permanece disponible hasta registrar la nota.
            if ultima_clase.fecha > hoy:
                continue

            resultado.append({
                'id': matricula.id,
                'estudiante_id': matricula.estudiante_id,
                'estudiante_nombre': (
                    f'{matricula.estudiante.nombre} '
                    f'{matricula.estudiante.apellido}'
                ).strip(),
                'estudiante_cedula': (
                    matricula.estudiante.cedula
                ),
                'tipo_curso': matricula.tipo_curso,
                'plan_nombre': (
                    matricula.plan_de_estudio.nombre
                    if matricula.plan_de_estudio
                    else 'Sin plan asignado'
                ),
                'fecha_inscripcion': (
                    matricula.fecha_registro
                ),
                'ultima_fecha_curso': (
                    ultima_clase.fecha
                ),
            })

        return Response(
            resultado,
            status=status.HTTP_200_OK
        )

    def create(self, request, *args, **kwargs):
        """
        Registra exclusivamente una nota práctica.
        La operación se valida nuevamente aunque el estudiante
        ya haya aparecido en el buscador del frontend. Así no se
        puede saltar la regla enviando un POST manual.
        """

        user = request.user

        if not es_instructor(user):
            return Response(
                {
                    'error': (
                        'Solo un instructor puede registrar la nota práctica.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        matricula_id = request.data.get('matricula')

        if not matricula_id:
            return Response(
                {
                    'error': (
                        'Debe seleccionar una matrícula.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        hoy = timezone.localdate()

        with transaction.atomic():
            try:
                matricula = (
                    Matricula.objects
                    .select_for_update()
                    .select_related(
                        'estudiante',
                        'plan_de_estudio',
                    )
                    .get(
                        id=matricula_id
                    )
                )
            except (
                Matricula.DoesNotExist,
                ValueError,
                TypeError,
            ):
                return Response(
                    {
                        'error': (
                            'Matrícula no encontrada.'
                        )
                    },
                    status=status.HTTP_404_NOT_FOUND
                )

            ultima_clase = (
                self.obtener_ultima_clase_practica(
                    matricula
                )
            )

            if ultima_clase is None:
                return Response(
                    {
                        'error': (
                            'La matrícula no tiene clases prácticas asignadas.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if (
                ultima_clase.instructor_id
                != user.instructor_id
            ):
                return Response(
                    {
                        'error': (
                            'Solo el instructor asignado a la última clase práctica puede registrar esta nota.'
                        )
                    },
                    status=status.HTTP_403_FORBIDDEN
                )

            if ultima_clase.fecha > hoy:
                return Response(
                    {
                        'error': (
                            'La nota práctica estará disponible '
                            'a partir del último día del curso: '
                            f'{ultima_clase.fecha.strftime("%d/%m/%Y")}.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if Notas.objects.filter(
                matricula=matricula,
                tipo_nota='practico',
            ).exists():
                return Response(
                    {
                        'error': (
                            'Ya existe una nota práctica registrada para este estudiante.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            plan_estudio = matricula.plan_de_estudio

            if not plan_estudio:
                return Response(
                    {
                        'error': (
                            'La matrícula no tiene un plan de estudio asignado.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            data = request.data.copy()
            data['matricula'] = matricula.id
            data['instructor'] = user.instructor_id
            data['plan_de_estudio'] = plan_estudio.id
            data['tipo_nota'] = 'practico'

            serializer = self.get_serializer(
                data=data
            )

            serializer.is_valid(
                raise_exception=True
            )

            nota = serializer.save()

            matricula_finalizada = (
                actualizar_estado_matricula_por_notas(
                    nota.matricula
                )
            )

            if matricula_finalizada:
                desactivar_usuarios_estudiante(
                    nota.matricula.estudiante
                )

        headers = self.get_success_headers(
            serializer.data
        )

        return Response(
            serializer.data,
            status=status.HTTP_201_CREATED,
            headers=headers,
        )

@api_view(['POST'])
@authentication_classes([])
@permission_classes([AllowAny])
@throttle_classes([LoginRateThrottle])
def login(request):
    username = request.data.get('username')
    password = request.data.get('password')

    if not isinstance(username, str):
        return Response(
            {
                'error': (
                    'El nombre de usuario y la contraseña '
                    'son obligatorios.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    if not isinstance(password, str):
        return Response(
            {
                'error': (
                    'El nombre de usuario y la contraseña '
                    'son obligatorios.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    username = username.strip()

    if not username or not password:
        return Response(
            {
                'error': (
                    'El nombre de usuario y la contraseña '
                    'son obligatorios.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    usuario_existente = (
        Usuario.objects
        .select_related('rol')
        .filter(username=username)
        .first()
    )

    # MySQL normalmente utiliza una comparación que no diferencia entre mayúsculas y minúsculas. Por eso comprobamos también
    # el texto exacto almacenado en el usuario.
    if (
        usuario_existente is None
        or usuario_existente.username != username
        or not usuario_existente.is_active
    ):
        return Response(
            {
                'error': 'Credenciales inválidas.'
            },
            status=status.HTTP_401_UNAUTHORIZED
        )

    user = authenticate(
        request=request,
        username=username,
        password=password
    )

    if (
        user is None
        or user.pk != usuario_existente.pk
        or user.username != username
        or not user.is_active
    ):
        return Response(
            {
                'error': 'Credenciales inválidas.'
            },
            status=status.HTTP_401_UNAUTHORIZED
        )

    with transaction.atomic():
        # Se elimina el token anterior para que un token antiguo no continúe funcionando después de un nuevo ingreso.
        Token.objects.filter(user=user).delete()

        token = Token.objects.create(user=user)

        update_last_login(
            sender=None,
            user=user
        )

    return Response(
        {
            'token': token.key,
            'user_id': user.id,
            'username': user.username,
            'email': user.email,
            'first_name': user.first_name,
            'last_name': user.last_name,
            'rol': (
                user.rol.nombre
                if user.rol
                else 'sin rol'
            ),
        },
        status=status.HTTP_200_OK
    )

@api_view(['POST'])
@permission_classes([IsAuthenticated])
def cerrar_sesion(request):
    """
    Elimina el token utilizado en la solicitud actual.
    Después de ejecutar esta acción, el mismo token ya no
    podrá utilizarse para acceder a rutas protegidas.
    """

    token_actual = request.auth

    if isinstance(token_actual, Token):
        token_actual.delete()
    else:
        Token.objects.filter(
            user=request.user
        ).delete()

    return Response(
        {
            'mensaje': 'Sesión cerrada correctamente.'
        },
        status=status.HTTP_200_OK
    )

class DashboardGananciasView(APIView):
    """Endpoint para obtener ganancias mensuales y matriculados"""
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'No tienes permiso para ver información del dashboard.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        try:
            hoy = timezone.localdate()
            anio_param = request.query_params.get('anio')

            try:
                anio = (
                    int(anio_param)
                    if anio_param
                    else hoy.year
                )
            except (
                ValueError,
                TypeError,
            ):
                anio = hoy.year

            if anio < 1 or anio > 9998:
                return Response(
                    {
                        'error': (
                            'El año indicado no es válido.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            inicio_anio = date(anio, 1, 1)
            inicio_anio_siguiente = date(
                anio + 1,
                1,
                1,
            )

            datos_por_mes = (
                Recibo.objects
                .filter(
                    fecha_pago__gte=inicio_anio,
                    fecha_pago__lt=inicio_anio_siguiente,
                    tipo_pago__in=TIPOS_RECIBO_INGRESO,
                )
                .annotate(
                    mes=TruncMonth('fecha_pago')
                )
                .values('mes')
                .annotate(
                    total=Sum('monto_pagado'),
                    matriculados=Count(
                        'matricula_id',
                        filter=Q(
                            tipo_pago__in=[
                                'completo',
                                'beneficio',
                            ]
                        ),
                        distinct=True,
                    ),
                )
                .order_by('mes')
            )

            resumen_por_mes = {
                fila['mes'].month: fila
                for fila in datos_por_mes
                if fila['mes'] is not None
            }

            meses_resultado = []

            for mes in range(1, 13):
                datos_mes = resumen_por_mes.get(
                    mes,
                    {}
                )

                meses_resultado.append({
                    'mes': f'{anio}-{mes:02d}',
                    'total': float(
                        datos_mes.get('total')
                        or Decimal('0')
                    ),
                    'matriculados': (
                        datos_mes.get('matriculados')
                        or 0
                    ),
                })

            return Response(meses_resultado)

        except Exception:
            logger.exception(
                'Error en DashboardGananciasView.'
            )

            return Response(
                {
                    'error': (
                        'No se pudieron cargar las '
                        'ganancias del dashboard.'
                    )
                },
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def _get_nombre_mes(self, mes_numero):
        meses = {
            1: "Enero", 2: "Febrero", 3: "Marzo", 4: "Abril",
            5: "Mayo", 6: "Junio", 7: "Julio", 8: "Agosto",
            9: "Septiembre", 10: "Octubre", 11: "Noviembre", 12: "Diciembre"
        }
        return meses.get(mes_numero, "")

class DashboardResumenView(APIView):
    """Endpoint para resumen del dashboard"""
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get(self, request):


        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para ver información del dashboard.'},
                status=status.HTTP_403_FORBIDDEN
            )
        try:
            hoy = timezone.localdate()

            inicio_mes = (
                timezone.localtime(timezone.now())
                .replace(
                    day=1,
                    hour=0,
                    minute=0,
                    second=0,
                    microsecond=0,
                )
            )

            if inicio_mes.month == 12:
                inicio_mes_siguiente = inicio_mes.replace(
                    year=inicio_mes.year + 1,
                    month=1,
                )
            else:
                inicio_mes_siguiente = inicio_mes.replace(
                    month=inicio_mes.month + 1,
                )

            inicio_mes_fecha = inicio_mes.date()
            inicio_mes_siguiente_fecha = inicio_mes_siguiente.date()

            # Total histórico de matrículas registradas.
            # Cada matrícula cuenta, aunque un estudiante
            # se haya matriculado más de una vez.
            total_matriculados = (
                Matricula.objects.count()
            )

            # Estudiantes que tienen al menos una matrícula no finalizada.
            # Si finalizan una matrícula, dejan de aparecer aquí.
            # Si luego reciben otra matrícula nueva, vuelven a aparecer.
            estudiantes_activos = (
                Matricula.objects
                .exclude(
                    estado='finalizado'
                )
                .values(
                    'estudiante_id'
                )
                .distinct()
                .count()
            )

            # Matrículas finalizadas durante el mes actual.
            # Se usa rango de fechas para evitar problemas de zona horaria.
            egresados_mes = Matricula.objects.filter(
                estado='finalizado',
                fecha_finalizacion__gte=inicio_mes,
                fecha_finalizacion__lt=inicio_mes_siguiente,
            ).count()

            # Ingresos del mes actual.
            # Recibo.fecha_pago es DateField, por eso se filtra con fechas.
            ingresos_mes = (
                Recibo.objects
                .filter(
                    fecha_pago__gte=inicio_mes_fecha,
                    fecha_pago__lt=inicio_mes_siguiente_fecha,
                    tipo_pago__in=TIPOS_RECIBO_INGRESO,
                )
                .aggregate(
                    total=Sum('monto_pagado')
                )['total']
                or Decimal('0')
            )

            # Ingresos totales históricos.
            ingresos_totales = (
                Recibo.objects
                .filter(
                    tipo_pago__in=TIPOS_RECIBO_INGRESO,
                )
                .aggregate(
                    total=Sum('monto_pagado')
                )['total']
                or Decimal('0')
            )

            return Response({
                'total_matriculados': total_matriculados,
                'estudiantes_activos': estudiantes_activos,
                'egresados_mes': egresados_mes,
                'ingresos_mes': float(ingresos_mes),
                'ingresos_totales': float(ingresos_totales),
            })

        except Exception:
            logger.exception(
                'Error en DashboardResumenView.'
            )

            return Response(
                {
                    'error': (
                        'No se pudo cargar el resumen del dashboard.'
                    )
                },
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

class DashboardIngresosMensualesView(APIView):
    """Endpoint específico para ingresos mensuales."""

    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get(self, request):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'No tienes permiso para ver '
                        'información del dashboard.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        try:
            hoy = timezone.localdate()

            anio_inicio = hoy.year
            mes_inicio = hoy.month - 5

            while mes_inicio <= 0:
                mes_inicio += 12
                anio_inicio -= 1

            inicio_periodo = date(
                anio_inicio,
                mes_inicio,
                1,
            )

            if hoy.month == 12:
                fin_periodo = date(
                    hoy.year + 1,
                    1,
                    1,
                )
            else:
                fin_periodo = date(
                    hoy.year,
                    hoy.month + 1,
                    1,
                )

            ingresos_agrupados = (
                Recibo.objects
                .filter(
                    fecha_pago__gte=inicio_periodo,
                    fecha_pago__lt=fin_periodo,
                    tipo_pago__in=TIPOS_RECIBO_INGRESO,
                )
                .annotate(
                    mes=TruncMonth('fecha_pago')
                )
                .values('mes')
                .annotate(
                    total=Sum('monto_pagado')
                )
                .order_by('mes')
            )

            ingresos_por_mes = {
                (
                    fila['mes'].year,
                    fila['mes'].month,
                ): (
                    fila['total']
                    or Decimal('0')
                )
                for fila in ingresos_agrupados
                if fila['mes'] is not None
            }

            meses_resultado = []

            for i in range(5, -1, -1):
                anio = hoy.year
                mes = hoy.month - i

                while mes <= 0:
                    mes += 12
                    anio -= 1

                total = ingresos_por_mes.get(
                    (anio, mes),
                    Decimal('0'),
                )

                meses_resultado.append({
                    'mes': f'{anio}-{mes:02d}',
                    'total': float(total),
                    'nombre_mes': self._get_nombre_mes(
                        mes
                    ),
                })

            return Response(meses_resultado)

        except Exception:
            logger.exception(
                'Error en DashboardIngresosMensualesView.'
            )

            return Response(
                {
                    'error': (
                        'No se pudieron cargar los '
                        'ingresos mensuales.'
                    )
                },
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    def _get_nombre_mes(self, mes_numero):
        meses = {
            1: 'Enero',
            2: 'Febrero',
            3: 'Marzo',
            4: 'Abril',
            5: 'Mayo',
            6: 'Junio',
            7: 'Julio',
            8: 'Agosto',
            9: 'Septiembre',
            10: 'Octubre',
            11: 'Noviembre',
            12: 'Diciembre',
        }

        return meses.get(
            mes_numero,
            ''
        )

class ProgresoTemaViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = ProgresoTema.objects.select_related(
        'matricula',
        'matricula__estudiante',
        'tema',
        'tema__plan_estudio'
    ).prefetch_related(
        'tema__subtemas'
    ).all()

    serializer_class = ProgresoTemaSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]


    def get_queryset(self):
        user = self.request.user
        rol = user.rol_nombre if hasattr(user, 'rol_nombre') else ''

        queryset = ProgresoTema.objects.select_related(
            'matricula',
            'matricula__estudiante',
            'tema',
            'tema__plan_estudio'
        ).prefetch_related(
            'tema__subtemas'
        ).filter(
            tema__plan_estudio_id=F(
                'matricula__plan_de_estudio_id'
            ),
            tema__activo=True,
        ).exclude(
            matricula_id__in=obtener_ids_matriculas_egresadas()
        ).order_by(
            'matricula_id',
            'orden_general',
            'id'
        )

        if es_admin(user):
            return queryset

        if rol == 'estudiante' and hasattr(user, 'estudiante') and user.estudiante:
            return queryset.filter(matricula__estudiante=user.estudiante)

        if rol == 'instructor' and hasattr(user, 'instructor') and user.instructor:
            return queryset.filter(
                matricula__clases__instructor=user.instructor
            ).distinct()

        return ProgresoTema.objects.none()

    def usuario_puede_acceder_matricula(self, user, matricula):
        if es_admin(user):
            return True

        if es_estudiante(user):
            return (
                matricula.estudiante_id ==
                user.estudiante_id
            )

        if es_instructor(user):
            return Calendario.objects.filter(
                matricula=matricula,
                instructor_id=user.instructor_id,
                es_examen=False
            ).exclude(
                estado='cancelada'
            ).exists()

        return False

    def list(self, request, *args, **kwargs):
        queryset = self.filter_queryset(
            self.get_queryset()
        )

        page = self.paginate_queryset(queryset)

        if page is not None:
            serializer = self.get_serializer(
                page,
                many=True
            )
            return self.get_paginated_response(
                serializer.data
            )

        serializer = self.get_serializer(
            queryset,
            many=True
        )

        return Response(serializer.data)

    def retrieve(self, request, *args, **kwargs):
        progreso = self.get_object()
        serializer = self.get_serializer(progreso)

        return Response(serializer.data)

    def es_modo_diario(self, matricula):
        return False

    def validar_curso_con_checks(self, matricula):
        if matricula_usa_checks(matricula):
            return None

        return Response(
            {
                'success': False,
                'usa_checks': False,
                'error': (
                    'Los estudiantes de cursos Intermedio y Avanzado '
                    'pueden consultar el plan de estudio, pero no '
                    'utilizan checks.'
                ),
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    def obtener_clase_habilitada_hoy(self, matricula):
        ahora = timezone.localtime()
        hoy = timezone.localdate()

        clases = Calendario.objects.filter(
            matricula=matricula,
            es_examen=False,
            fecha=hoy
        ).exclude(
            estado='cancelada'
        ).order_by(
            'hora_inicio',
            'id'
        )

        for clase in clases:
            if not clase.hora_inicio:
                continue

            inicio = datetime.combine(clase.fecha, clase.hora_inicio)

            if timezone.is_naive(inicio):
                inicio = timezone.make_aware(
                    inicio,
                    timezone.get_current_timezone()
                )

            momento_habilitado = inicio - timedelta(minutes=10)

            if ahora >= momento_habilitado:
                return clase
        return None

    def obtener_check_diario_actual(self, progreso, crear=False):
        clase = self.obtener_clase_habilitada_hoy(progreso.matricula)

        if not clase:
            return None

        if crear:
            check, _ = ProgresoClaseTema.objects.get_or_create(
                calendario=clase,
                progreso_tema=progreso
            )
            return check

        return ProgresoClaseTema.objects.filter(
            calendario=clase,
            progreso_tema=progreso
        ).first()

    def preparar_contexto_diario(self, progreso):
        if not self.es_modo_diario(progreso.matricula):
            return progreso

        clase = self.obtener_clase_habilitada_hoy(progreso.matricula)
        progreso.calendario_actual = clase
        progreso.habilitado_hoy = bool(clase)

        if clase:
            check, _ = ProgresoClaseTema.objects.get_or_create(
                calendario=clase,
                progreso_tema=progreso
            )
            progreso.check_dia_actual = check
        else:
            progreso.check_dia_actual = None

        return progreso

    def todas_las_clases_diarias_completadas(self, progreso):
        clases = Calendario.objects.filter(
            matricula=progreso.matricula,
            es_examen=False
        ).exclude(
            estado='cancelada'
        )

        total_clases = clases.count()

        if total_clases == 0:
            return False

        total_checks_completados = ProgresoClaseTema.objects.filter(
            progreso_tema=progreso,
            calendario__in=clases,
            estudiante_completado=True,
            instructor_completado=True,
            completado=True
        ).count()

        return total_checks_completados >= total_clases

    def actualizar_estado_global_diario(self, progreso):
        clases = Calendario.objects.filter(
            matricula=progreso.matricula,
            es_examen=False
        ).exclude(
            estado='cancelada'
        )

        total_clases = clases.count()

        if total_clases == 0:
            progreso.estudiante_completado = False
            progreso.instructor_completado = False
            progreso.completado = False
            progreso.save(
                update_fields=[
                    'estudiante_completado',
                    'instructor_completado',
                    'completado',
                ]
            )
            return

        total_checks_completados = ProgresoClaseTema.objects.filter(
            progreso_tema=progreso,
            calendario__in=clases,
            estudiante_completado=True,
            instructor_completado=True,
            completado=True
        ).count()

        plan_diario_completo = total_checks_completados >= total_clases
        progreso.estudiante_completado = plan_diario_completo
        progreso.instructor_completado = plan_diario_completo
        progreso.completado = plan_diario_completo

        if plan_diario_completo:
            progreso.desbloqueado = True

            ahora = timezone.now()

            if not progreso.fecha_estudiante:
                progreso.fecha_estudiante = ahora

            if not progreso.fecha_instructor:
                progreso.fecha_instructor = ahora

            if not progreso.fecha_completado:
                progreso.fecha_completado = ahora
        else:
            progreso.fecha_completado = None

        progreso.save(
            update_fields=[
                'estudiante_completado',
                'instructor_completado',
                'completado',
                'desbloqueado',
                'fecha_estudiante',
                'fecha_instructor',
                'fecha_completado',
            ]
        )

    def normalizar_desbloqueo_diario(self, matricula):
        progresos = list(
            ProgresoTema.objects.filter(
                matricula=matricula
            ).order_by(
                'orden_general',
                'id'
            )
        )

        if not progresos:
            return

        clase = self.obtener_clase_habilitada_hoy(matricula)
        habilitado = bool(clase)

        for index, progreso in enumerate(progresos):
            progreso.desbloqueado = habilitado and index == 0
            progreso.save(update_fields=['desbloqueado'])

            if habilitado and index == 0:
                ProgresoClaseTema.objects.get_or_create(
                    calendario=clase,
                    progreso_tema=progreso
                )

    def normalizar_desbloqueo_principiante(self, matricula):
        hoy = timezone.localdate()
        progresos = list(
            ProgresoTema.objects.filter(
                matricula=matricula
            ).order_by(
                'orden_general',
                'id'
            )
        )

        if not progresos:
            return

        clases = Calendario.objects.filter(
            matricula=matricula,
            es_examen=False
        ).exclude(
            estado='cancelada'
        ).order_by(
            'fecha',
            'hora_inicio'
        )

        primera_clase = clases.first()

        if not primera_clase:
            ProgresoTema.objects.filter(
                matricula=matricula
            ).update(
                desbloqueado=False
            )
            return

        if primera_clase.fecha > hoy:
            ProgresoTema.objects.filter(
                matricula=matricula
            ).update(
                desbloqueado=False
            )
            return

        clases_hasta_hoy = clases.filter(
            fecha__lte=hoy
        )

        limite_temas = 0

        for clase in clases_hasta_hoy:
            if clase.hora_inicio and clase.hora_fin:
                inicio = datetime.combine(
                    clase.fecha,
                    clase.hora_inicio
                )

                fin = datetime.combine(
                    clase.fecha,
                    clase.hora_fin
                )

                horas_por_dia = int(
                    (fin - inicio).total_seconds() // 3600
                )

                if horas_por_dia <= 0:
                    horas_por_dia = 1
            else:
                horas_por_dia = 1

            limite_temas += horas_por_dia

        limite_temas = min(
            limite_temas,
            len(progresos)
        )

        progresos_actualizar = []

        for index, progreso in enumerate(progresos):
            nuevo_estado = index < limite_temas

            if progreso.desbloqueado != nuevo_estado:
                progreso.desbloqueado = nuevo_estado
                progresos_actualizar.append(progreso)

        if progresos_actualizar:
            ProgresoTema.objects.bulk_update(
                progresos_actualizar,
                ['desbloqueado']
            )

    def normalizar_desbloqueo(self, matricula):
        if self.es_modo_diario(matricula):
            self.normalizar_desbloqueo_diario(matricula)
            return

        self.normalizar_desbloqueo_principiante(matricula)

    def _actualizar_progreso(self, progreso):
        progreso.completado = (
            progreso.estudiante_completado and
            progreso.instructor_completado
        )

        if progreso.completado and not progreso.fecha_completado:
            progreso.fecha_completado = timezone.now()

        progreso.save()

        self._crear_notificacion_pendiente(progreso)

    def _obtener_horas_clase_actual(self, progreso):
        clase = Calendario.objects.filter(
            matricula=progreso.matricula,
            es_examen=False
        ).order_by('fecha', 'hora_inicio').first()

        if not clase:
            return 1

        inicio = datetime.combine(clase.fecha, clase.hora_inicio)
        fin = datetime.combine(clase.fecha, clase.hora_fin)
        horas = int((fin - inicio).total_seconds() // 3600)

        return max(horas, 1)

    def _desbloquear_siguiente(self, progreso_actual):
        limite_temas = self._obtener_horas_clase_actual(progreso_actual)

        progresos = list(
            ProgresoTema.objects.filter(
                matricula=progreso_actual.matricula
            ).select_related(
                'tema'
            ).order_by(
                'orden_general',
                'id'
            )
        )

        for progreso in progresos:
            progreso.desbloqueado = False
            progreso.save(update_fields=['desbloqueado'])

        pendientes = [
            progreso for progreso in progresos
            if not progreso.completado
        ]
        bloque_actual = pendientes[:limite_temas]

        bloque_incompleto = any(
            not (
                item.estudiante_completado and
                item.instructor_completado
            )
            for item in bloque_actual
        )

        if bloque_incompleto:
            return

        usuario_estudiante = progreso_actual.matricula.estudiante.usuarios.first()
        for progreso in pendientes[:limite_temas]:
            progreso.desbloqueado = True
            progreso.save(update_fields=['desbloqueado'])

    def _obtener_usuario_estudiante(self, progreso):
        estudiante = progreso.matricula.estudiante

        usuario = estudiante.usuarios.first()

        return usuario

    def _crear_notificacion_pendiente(self, progreso):
        estudiante = progreso.matricula.estudiante
        estudiante_nombre = f"{estudiante.nombre} {estudiante.apellido}".strip()
        usuario_estudiante = self._obtener_usuario_estudiante(progreso)

        if not usuario_estudiante:
            return

        Notificacion.objects.filter(
            estudiante=usuario_estudiante,
            progreso_tema=progreso,
            tema=progreso.tema,
            tipo__in=['falta_estudiante', 'falta_instructor'],
            leida=False
        ).update(leida=True)

        if progreso.estudiante_completado and progreso.instructor_completado:
            return

        clase = Calendario.objects.filter(
            matricula=progreso.matricula,
            instructor__isnull=False
        ).select_related('instructor').first()

        instructor_nombre = "Instructor no asignado"

        if clase and clase.instructor:
            instructor_nombre = (
                f"{clase.instructor.nombre or ''} "
                f"{clase.instructor.apellido or ''}"
            ).strip() or "Instructor no asignado"

        if progreso.estudiante_completado and not progreso.instructor_completado:
            Notificacion.objects.update_or_create(
                estudiante=usuario_estudiante,
                progreso_tema=progreso,
                tema=progreso.tema,
                tipo='falta_instructor',
                defaults={
                    'mensaje': (
                        f'El instructor "{instructor_nombre}" no ha dado check '
                        f'al tema "{progreso.tema.titulo}" del estudiante '
                        f'"{estudiante_nombre}".'
                    ),
                    'leida': False,
                    'fecha_creacion': timezone.now(),
                }
            )

        elif progreso.instructor_completado and not progreso.estudiante_completado:
            Notificacion.objects.update_or_create(
                estudiante=usuario_estudiante,
                progreso_tema=progreso,
                tema=progreso.tema,
                tipo='falta_estudiante',
                defaults={
                    'mensaje': (
                        f'El estudiante "{estudiante_nombre}" no ha dado check '
                        f'al tema "{progreso.tema.titulo}" marcado por el instructor '
                        f'"{instructor_nombre}".'
                    ),
                    'leida': False,
                    'fecha_creacion': timezone.now(),
                }
            )

    @action(detail=False, methods=['post'], url_path='actualizar-desbloqueos')
    def actualizar_desbloqueos(self, request):
        matricula_ids = request.data.get('matricula_ids')

        # Mantiene compatibilidad con llamadas antiguas.
        if matricula_ids is None:
            matricula_id = request.data.get('matricula_id')
            matricula_ids = [matricula_id] if matricula_id else []

        if not isinstance(matricula_ids, list):
            return Response(
                {
                    'success': False,
                    'error': 'matricula_ids debe ser una lista.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            matricula_ids = list(dict.fromkeys(
                int(matricula_id)
                for matricula_id in matricula_ids
                if matricula_id
            ))
        except (TypeError, ValueError):
            return Response(
                {
                    'success': False,
                    'error': 'Las matrículas enviadas no son válidas.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if not matricula_ids:
            return Response(
                {
                    'success': False,
                    'error': 'Debe enviar al menos una matrícula.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if len(matricula_ids) > 100:
            return Response(
                {
                    'success': False,
                    'error': 'Solo se permiten 100 matrículas por solicitud.'
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        matriculas_encontradas = {
            matricula.id: matricula
            for matricula in Matricula.objects.select_related(
                'estudiante'
            ).filter(id__in=matricula_ids)
        }

        ids_no_encontrados = [
            matricula_id
            for matricula_id in matricula_ids
            if matricula_id not in matriculas_encontradas
        ]

        if ids_no_encontrados:
            return Response(
                {
                    'success': False,
                    'error': 'Una o más matrículas no existen.',
                    'matricula_ids': ids_no_encontrados,
                },
                status=status.HTTP_404_NOT_FOUND
            )

        matriculas = [
            matriculas_encontradas[matricula_id]
            for matricula_id in matricula_ids
        ]

        # Primero valida todas. Así no actualiza unas antes de
        # descubrir que el usuario no tiene acceso a otra.
        for matricula in matriculas:
            if not self.usuario_puede_acceder_matricula(
                request.user,
                matricula
            ):
                return Response(
                    {
                        'success': False,
                        'error': (
                            'No tiene permiso para actualizar el progreso '
                            'de una o más matrículas.'
                        )
                    },
                    status=status.HTTP_403_FORBIDDEN
                )

        actualizadas = 0

        with transaction.atomic():
            for matricula in matriculas:
                # Los checks continúan siendo solamente para Principiante.
                if not matricula_usa_checks(matricula):
                    continue

                self.normalizar_desbloqueo(matricula)
                actualizadas += 1

        return Response({
            'success': True,
            'actualizadas': actualizadas,
            'message': 'Desbloqueos actualizados correctamente.'
        })

    @action(detail=True, methods=['post'], url_path='marcar-estudiante')
    def marcar_estudiante(self, request, pk=None):
        if not es_estudiante(request.user):
            return Response(
                {
                    'success': False,
                    'error': (
                        'Solo el estudiante puede marcar su check del plan de estudio.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        try:
            progreso = self.get_object()
            bloqueo = self.validar_curso_con_checks(
                progreso.matricula
            )

            if bloqueo:
                return bloqueo

            if self.es_modo_diario(progreso.matricula):
                self.normalizar_desbloqueo(progreso.matricula)
                progreso.refresh_from_db()
                progreso = self.preparar_contexto_diario(progreso)

                if not getattr(progreso, 'habilitado_hoy', False):
                    return Response({
                        'success': False,
                        'error': 'Este tema aún no está disponible para la clase de hoy.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                check_dia = self.obtener_check_diario_actual(
                    progreso,
                    crear=True
                )

                if not check_dia:
                    return Response({
                        'success': False,
                        'error': 'No hay clase habilitada para marcar este tema.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                if check_dia.estudiante_completado:
                    return Response({
                        'success': False,
                        'error': 'Ya marcaste este tema como recibido en la clase de hoy.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                with transaction.atomic():
                    check_dia.estudiante_completado = True
                    check_dia.fecha_estudiante = timezone.now()
                    check_dia.save()

                    if check_dia.completado:
                        self.actualizar_estado_global_diario(progreso)
                        mensaje = "Tema completado correctamente para la clase de hoy."
                    else:
                        mensaje = "Tema marcado como recibido para la clase de hoy. Falta confirmación del instructor."

                progreso.refresh_from_db()
                progreso = self.preparar_contexto_diario(progreso)

                serializer = self.get_serializer(progreso)

                return Response({
                    'success': True,
                    'message': mensaje,
                    'data': serializer.data
                })

            if not progreso.desbloqueado:
                return Response({
                    'success': False,
                    'error': 'Este tema aún no está disponible.'
                }, status=status.HTTP_400_BAD_REQUEST)

            if progreso.estudiante_completado:
                return Response({
                    'success': False,
                    'error': 'Ya habías marcado este tema como recibido'
                }, status=status.HTTP_400_BAD_REQUEST)

            with transaction.atomic():
                progreso.estudiante_completado = True
                progreso.fecha_estudiante = timezone.now()

                if progreso.estudiante_completado and progreso.instructor_completado:
                    progreso.completado = True

                    if not progreso.fecha_completado:
                        progreso.fecha_completado = timezone.now()

                progreso.save()

                if progreso.estudiante_completado and progreso.instructor_completado:
                    Notificacion.objects.filter(
                        estudiante=self._obtener_usuario_estudiante(progreso),
                        progreso_tema=progreso,
                        tema=progreso.tema,
                        tipo__in=['falta_estudiante', 'falta_instructor'],
                        leida=False
                    ).update(leida=True)
                    self.normalizar_desbloqueo(progreso.matricula)
                else:
                    self._crear_notificacion_pendiente(progreso)

                if progreso.instructor_completado:
                    mensaje = "Tema completado correctamente."
                else:
                    mensaje = "Tema marcado como recibido. Falta confirmación del instructor."

            serializer = self.get_serializer(progreso)

            return Response({
                'success': True,
                'message': mensaje,
                'data': serializer.data
            })

        except Exception:
            logger.exception(
                'Error inesperado al marcar el check del estudiante. '
                'progreso_id=%s usuario_id=%s',
                pk,
                request.user.id,
            )

            return Response(
                {
                    'success': False,
                    'error': (
                        'No se pudo registrar el check del estudiante. '
                        'Inténtelo nuevamente.'
                    )
                },
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )


    @action(detail=True, methods=['post'], url_path='marcar-instructor')
    def marcar_instructor(self, request, pk=None):
        if not es_instructor(request.user):
            return Response(
                {
                    'success': False,
                    'error': (
                        'Solo el instructor puede marcar '
                        'el check correspondiente al instructor.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        try:
            progreso = self.get_object()
            bloqueo = self.validar_curso_con_checks(
                progreso.matricula
            )

            if bloqueo:
                return bloqueo

            if self.es_modo_diario(progreso.matricula):
                self.normalizar_desbloqueo(progreso.matricula)
                progreso.refresh_from_db()
                progreso = self.preparar_contexto_diario(progreso)

                if not getattr(progreso, 'habilitado_hoy', False):
                    return Response({
                        'success': False,
                        'error': 'Este tema aún no está disponible para la clase de hoy.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                check_dia = self.obtener_check_diario_actual(
                    progreso,
                    crear=True
                )

                if not check_dia:
                    return Response({
                        'success': False,
                        'error': 'No hay clase habilitada para marcar este tema.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                if check_dia.instructor_completado:
                    return Response({
                        'success': False,
                        'error': 'Ya marcaste este tema como dado en la clase de hoy.'
                    }, status=status.HTTP_400_BAD_REQUEST)

                with transaction.atomic():
                    check_dia.instructor_completado = True
                    check_dia.fecha_instructor = timezone.now()
                    check_dia.save()

                    if check_dia.completado:
                        self.actualizar_estado_global_diario(progreso)
                        mensaje = "Tema completado correctamente para la clase de hoy."
                    else:
                        mensaje = "Tema marcado como dado para la clase de hoy. Falta confirmación del estudiante."

                progreso.refresh_from_db()
                progreso = self.preparar_contexto_diario(progreso)

                serializer = self.get_serializer(progreso)

                return Response({
                    'success': True,
                    'message': mensaje,
                    'data': serializer.data
                })

            if not progreso.desbloqueado:
                return Response({
                    'success': False,
                    'error': 'Este tema aún no está disponible.'
                }, status=status.HTTP_400_BAD_REQUEST)

            if progreso.instructor_completado:
                return Response({
                    'success': False,
                    'error': 'Ya habías marcado este tema como dado'
                }, status=status.HTTP_400_BAD_REQUEST)

            with transaction.atomic():
                progreso.instructor_completado = True
                progreso.fecha_instructor = timezone.now()

                if progreso.estudiante_completado and progreso.instructor_completado:
                    progreso.completado = True

                    if not progreso.fecha_completado:
                        progreso.fecha_completado = timezone.now()

                progreso.save()

                if progreso.estudiante_completado and progreso.instructor_completado:
                    Notificacion.objects.filter(
                        estudiante=self._obtener_usuario_estudiante(progreso),
                        progreso_tema=progreso,
                        tema=progreso.tema,
                        tipo__in=['falta_estudiante', 'falta_instructor'],
                        leida=False
                    ).update(leida=True)
                    self.normalizar_desbloqueo(progreso.matricula)
                else:
                    self._crear_notificacion_pendiente(progreso)

                if progreso.estudiante_completado:
                    mensaje = "Tema completado correctamente."
                else:
                    mensaje = "Tema marcado como dado. Falta confirmación del estudiante."

            serializer = self.get_serializer(progreso)

            return Response({
                'success': True,
                'message': mensaje,
                'data': serializer.data
            })

        except Exception:
            logger.exception(
                'Error inesperado al marcar el check del instructor. '
                'progreso_id=%s usuario_id=%s',
                pk,
                request.user.id,
            )

            return Response(
                {
                    'success': False,
                    'error': (
                        'No se pudo registrar el check del instructor. '
                        'Inténtelo nuevamente.'
                    )
                },
                status=status.HTTP_500_INTERNAL_SERVER_ERROR
            )

    @action(detail=True, methods=['post'], url_path='admin-forzar')
    def admin_forzar(self, request, pk=None):
        rol = str(getattr(request.user, 'rol_nombre', '') or '').lower()

        if not es_admin(request.user):
            return Response({
                'success': False,
                'error': (
                    'Solo administración puede realizar esta acción.'
                )
            }, status=status.HTTP_403_FORBIDDEN)

        progreso = self.get_object()

        bloqueo = self.validar_curso_con_checks(
            progreso.matricula
        )

        if bloqueo:
            return bloqueo

        tipo_check = request.data.get('tipo')
        valor = request.data.get('valor')

        if isinstance(valor, str):
            valor = valor.lower() == 'true'

        if tipo_check not in ['estudiante', 'instructor']:
            return Response({
                'success': False,
                'error': 'Tipo de check inválido. Use "estudiante" o "instructor"'
            }, status=status.HTTP_400_BAD_REQUEST)

        if valor not in [True, False]:
            return Response({
                'success': False,
                'error': 'Valor inválido. Use true o false'
            }, status=status.HTTP_400_BAD_REQUEST)

        if self.es_modo_diario(progreso.matricula):
            self.normalizar_desbloqueo(progreso.matricula)
            progreso.refresh_from_db()
            progreso = self.preparar_contexto_diario(progreso)

            if not getattr(progreso, 'habilitado_hoy', False):
                return Response({
                    'success': False,
                    'error': 'Este tema aún no está disponible para la clase de hoy.'
                }, status=status.HTTP_400_BAD_REQUEST)

            check_dia = self.obtener_check_diario_actual(
                progreso,
                crear=True
            )

            if not check_dia:
                return Response({
                    'success': False,
                    'error': 'No hay clase habilitada para modificar este tema.'
                }, status=status.HTTP_400_BAD_REQUEST)

            with transaction.atomic():
                old_estudiante = check_dia.estudiante_completado
                old_instructor = check_dia.instructor_completado

                if tipo_check == 'estudiante':
                    check_dia.estudiante_completado = valor
                    check_dia.fecha_estudiante = timezone.now() if valor else None
                else:
                    check_dia.instructor_completado = valor
                    check_dia.fecha_instructor = timezone.now() if valor else None

                if not valor:
                    check_dia.fecha_completado = None

                check_dia.save()

                if check_dia.completado:
                    self.actualizar_estado_global_diario(progreso)

                try:
                    HistorialPlanEstudio.objects.create(
                        progreso_tema=progreso,
                        usuario=request.user,
                        accion='admin_forzar_diario',
                        valor_anterior_estudiante=old_estudiante,
                        valor_anterior_instructor=old_instructor,
                        valor_nuevo_estudiante=check_dia.estudiante_completado,
                        valor_nuevo_instructor=check_dia.instructor_completado
                    )
                except Exception as e:
                    print(f"Error al guardar historial diario: {e}")

            progreso.refresh_from_db()
            progreso = self.preparar_contexto_diario(progreso)

            serializer = self.get_serializer(progreso)

            check_nombre = "estudiante" if tipo_check == 'estudiante' else "instructor"

            return Response({
                'success': True,
                'message': f'Check diario de {check_nombre} actualizado correctamente',
                'data': serializer.data
            })

        with transaction.atomic():
            old_estudiante = progreso.estudiante_completado
            old_instructor = progreso.instructor_completado

            if tipo_check == 'estudiante':
                progreso.estudiante_completado = valor
                progreso.fecha_estudiante = timezone.now() if valor else None
            else:
                progreso.instructor_completado = valor
                progreso.fecha_instructor = timezone.now() if valor else None

            progreso.fecha_admin_edit = timezone.now()

            if not valor:
                progreso.fecha_completado = None

            self._actualizar_progreso(progreso)
            self.normalizar_desbloqueo(progreso.matricula)

            try:
                HistorialPlanEstudio.objects.create(
                    progreso_tema=progreso,
                    usuario=request.user,
                    accion='admin_forzar',
                    valor_anterior_estudiante=old_estudiante,
                    valor_anterior_instructor=old_instructor,
                    valor_nuevo_estudiante=progreso.estudiante_completado,
                    valor_nuevo_instructor=progreso.instructor_completado
                )
            except Exception as e:
                print(f"Error al guardar historial: {e}")

            usuario_estudiante = progreso.matricula.estudiante.usuarios.first()

            if usuario_estudiante:
                try:
                    Notificacion.objects.update_or_create(
                        estudiante=usuario_estudiante,
                        progreso_tema=progreso,
                        tema=progreso.tema,
                        tipo='tema_desbloqueado',
                        defaults={
                            'mensaje': (
                                f"Nuevo tema desbloqueado: "
                                f"'{progreso.tema.titulo}'."
                            ),
                            'leida': False
                        }
                    )
                except Exception as e:
                    print(f"Error creando notificación: {e}")

        serializer = self.get_serializer(progreso)

        check_nombre = "estudiante" if tipo_check == 'estudiante' else "instructor"

        return Response({
            'success': True,
            'message': f'Check de {check_nombre} actualizado exitosamente',
            'data': serializer.data
        })

    @action(detail=False, methods=['get'], url_path='verificar-plan-completado')
    def verificar_plan_completado(self, request):
        user = self.request.user
        rol = user.rol_nombre if hasattr(user, 'rol_nombre') else ''

        resultado = []

        if es_admin(user):
            matriculas = Matricula.objects.filter(
                estado='matriculado'
            )

        elif hasattr(user, 'instructor') and user.instructor:
            matriculas = Matricula.objects.filter(
                clases__instructor=user.instructor,
                estado='matriculado'
            ).distinct()

        elif hasattr(user, 'estudiante') and user.estudiante:
            matriculas = Matricula.objects.filter(
                estudiante=user.estudiante,
                estado='matriculado'
            )

        else:
            return Response({
                'success': False,
                'error': 'Usuario no tiene permisos para acceder a esta información'
            }, status=status.HTTP_403_FORBIDDEN)

        for matricula in matriculas:
            if not matricula.plan_de_estudio_id:
                continue

            validacion_plan = validar_plan_completado_para_examen(matricula)

            usa_checks = validacion_plan.get(
                'usa_checks',
                True
            )

            progresos = obtener_progresos_plan_actual(matricula)
            total_temas = progresos.count()

            if usa_checks:
                temas_completados = progresos.filter(
                    Q(completado=True)
                    |
                    Q(
                        estudiante_completado=True,
                        instructor_completado=True,
                    )
                ).distinct().count()

                porcentaje = (
                    round(
                        (temas_completados / total_temas) * 100
                    )
                    if total_temas > 0
                    else 0
                )
            else:
                temas_completados = None
                porcentaje = None

            resultado.append({
                'matricula_id': matricula.id,
                'plan_nombre': (
                    matricula.plan_de_estudio.nombre
                    if matricula.plan_de_estudio
                    else 'Sin plan'
                ),
                'tipo_curso': matricula.tipo_curso,
                'usa_checks': usa_checks,
                'total_temas': total_temas,
                'temas_completados': temas_completados,
                'porcentaje': porcentaje,
                'progreso': validacion_plan['progreso'],
                'plan_completado': (
                    validacion_plan['completo']
                    if usa_checks
                    else None
                ),
                'puede_presentar_examen': True,
                'estudiante_nombre': (
                    f'{matricula.estudiante.nombre} '
                    f'{matricula.estudiante.apellido}'
                ).strip(),
                'estudiante_cedula': (
                    matricula.estudiante.cedula
                ),
                'estudiante_id': (
                    matricula.estudiante.id
                ),
                'fecha_inscripcion': (
                    matricula.fecha_registro
                ),
            })

        return Response(resultado)

    @action(detail=True, methods=['get'], url_path='plan-completado')
    def plan_completado(self, request, pk=None):
        matricula = get_object_or_404(
            Matricula.objects.select_related(
                'estudiante',
                'plan_de_estudio'
            ),
            id=pk
        )

        if not self.usuario_puede_acceder_matricula(
            request.user,
            matricula
        ):
            return Response(
                {
                    'success': False,
                    'error': (
                        'No tiene permiso para consultar el progreso de esta matrícula.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        validacion_plan = (
            validar_plan_completado_para_examen(
                matricula
            )
        )

        usa_checks = validacion_plan.get(
            'usa_checks',
            True
        )

        if usa_checks:
            mensaje = (
                'Los checks son de seguimiento opcional y no bloquean el examen teórico. '
                f"Progreso actual: {validacion_plan['progreso']}."
            )
        else:
            mensaje = (
                'Este curso muestra el plan como contenido informativo y no requiere checks.'
            )

        return Response({
            'success': True,
            'matricula_id': matricula.id,
            'tipo_curso': matricula.tipo_curso,
            'usa_checks': usa_checks,
            'plan_completado': (
                validacion_plan['completo']
                if usa_checks
                else None
            ),
            'puede_presentar_examen': True,
            'progreso': validacion_plan['progreso'],
            'mensaje': mensaje,
        })

class NotificacionViewSet(viewsets.ReadOnlyModelViewSet):
    """ViewSet para notificaciones del administrador"""

    queryset = Notificacion.objects.all()
    serializer_class = NotificacionSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        if es_admin(user):
            return Notificacion.objects.all().order_by(
                '-fecha_creacion'
            )

        return Notificacion.objects.none()

    @action(detail=True, methods=['post'], url_path='marcar-leida')
    def marcar_leida(self, request, pk=None):
        notificacion = self.get_object()
        notificacion.leida = True
        notificacion.save(update_fields=['leida'])

        return Response({
            'success': True,
            'message': 'Notificación marcada como leída'
        })

    @action(detail=False, methods=['get'], url_path='admin-pendientes')
    def admin_pendientes(self, request):
        user = request.user

        if not es_admin(user):
            return Response(
                {
                    'error': (
                        'Solo el administrador puede ver estas notificaciones.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        limite_tiempo = timezone.now() - timedelta(
            hours=23
        )

        Notificacion.objects.filter(
            leida=False,
            fecha_creacion__lt=limite_tiempo
        ).update(
            leida=True
        )

        notificaciones = (
            Notificacion.objects
            .filter(
                leida=False,
                tipo__in=[
                    'falta_estudiante',
                    'falta_instructor',
                ],
                fecha_creacion__gte=limite_tiempo,
            )
            .select_related(
                'estudiante',
                'tema',
                'progreso_tema',
                'progreso_tema__matricula',
                'progreso_tema__matricula__estudiante',
                'progreso_tema__tema',
            )
            .order_by('-fecha_creacion')
        )

        resultado = []
        ids_cerrar = []

        for notificacion in notificaciones:
            progreso = notificacion.progreso_tema

            # Las notificaciones anteriores a esta migración no tienen
            # una matrícula confiable, por eso no deben mostrarse.
            if not progreso:
                ids_cerrar.append(notificacion.id)
                continue

            if not matricula_usa_checks(progreso.matricula):
                ids_cerrar.append(notificacion.id)
                continue

            if (
                progreso.estudiante_completado
                and progreso.instructor_completado
            ):
                ids_cerrar.append(notificacion.id)
                continue

            if (
                notificacion.tipo == 'falta_estudiante'
                and progreso.estudiante_completado
            ):
                ids_cerrar.append(notificacion.id)
                continue

            if (
                notificacion.tipo == 'falta_instructor'
                and progreso.instructor_completado
            ):
                ids_cerrar.append(notificacion.id)
                continue

            estudiante = progreso.matricula.estudiante

            estudiante_nombre = (
                f'{estudiante.nombre or ""} '
                f'{estudiante.apellido or ""}'
            ).strip()

            falta_instructor = (
                notificacion.tipo == 'falta_instructor'
            )

            resultado.append({
                'id': notificacion.id,
                'tipo': notificacion.tipo,
                'tipo_texto': (
                    'Falta check del instructor'
                    if falta_instructor
                    else 'Falta check del estudiante'
                ),
                'quien_falta': (
                    'Instructor'
                    if falta_instructor
                    else 'Estudiante'
                ),
                'quien_espera': (
                    'Estudiante'
                    if falta_instructor
                    else 'Instructor'
                ),
                'mensaje': notificacion.mensaje,
                'estudiante': estudiante_nombre,
                'tema': progreso.tema.titulo,
                'matricula_id': progreso.matricula_id,
                'tipo_curso': progreso.matricula.tipo_curso,
                'fecha_creacion': notificacion.fecha_creacion,
            })

            if len(resultado) >= 10:
                break

        if ids_cerrar:
            Notificacion.objects.filter(
                id__in=ids_cerrar
            ).update(
                leida=True
            )

        return Response(resultado)

class DashboardPlanViewSet(viewsets.ViewSet):
    """Resumen del plan para el estudiante autenticado."""

    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'head',
        'options',
    ]

    @action(detail=False, methods=['get'], url_path='mi-progreso')
    def mi_progreso(self, request):
        user = request.user

        estudiante = getattr(
            user,
            'estudiante',
            None
        )

        if not estudiante:
            return Response({
                'porcentaje': 0,
                'temas_completados': 0,
                'total_temas': 0,
                'unidad': 'temas',
                'tipo_curso': None,
                'matricula_id': None,
                'aplica_progreso': False,
                'examen_disponible': False,
            })

        matricula = Matricula.objects.select_related(
            'estudiante',
            'plan_de_estudio',
        ).filter(
            estudiante=estudiante,
            estado='matriculado'
        ).order_by(
            '-id'
        ).first()

        if not matricula:
            return Response({
                'porcentaje': 0,
                'temas_completados': 0,
                'total_temas': 0,
                'unidad': 'temas',
                'tipo_curso': None,
                'matricula_id': None,
                'aplica_progreso': False,
                'examen_disponible': False,
            })

        tipo = str(
            matricula.tipo_curso or ''
        ).strip().lower()

        progresos = ProgresoTema.objects.filter(
            matricula=matricula
        ).order_by(
            'orden_general',
            'id'
        )

        total_temas = progresos.count()

        # Intermedio y Avanzado no tienen porcentaje.
        if tipo in ['intermedio', 'avanzado']:
            return Response({
                'porcentaje': None,
                'temas_completados': 0,
                'total_temas': total_temas,
                'unidad': 'temas',
                'tipo_curso': matricula.tipo_curso,
                'matricula_id': matricula.id,
                'aplica_progreso': False,
                'examen_disponible': True,
                'mensaje_examen': (
                    'Este curso no usa checks como requisito para el examen teórico.'
                ),
            })

        if total_temas == 0:
            return Response({
                'porcentaje': 0,
                'temas_completados': 0,
                'total_temas': 0,
                'unidad': 'temas',
                'tipo_curso': matricula.tipo_curso,
                'matricula_id': matricula.id,
                'aplica_progreso': True,
                'examen_disponible': True,
                'mensaje_examen': (
                    'El examen teórico no depende de los checks '
                    'del plan de estudio.'
                ),
            })

        temas_completados = progresos.filter(
            estudiante_completado=True,
            instructor_completado=True
        ).count()

        porcentaje = round(
            (temas_completados / total_temas) * 100
        )

        return Response({
            'porcentaje': porcentaje,
            'temas_completados': temas_completados,
            'total_temas': total_temas,
            'unidad': 'temas',
            'tipo_curso': matricula.tipo_curso,
            'matricula_id': matricula.id,
            'aplica_progreso': True,
            'examen_disponible': True,
        })

class PreguntaExamenTeoricoViewSet(viewsets.ModelViewSet):
    queryset = PreguntaExamenTeorico.objects.prefetch_related('opciones').all()
    serializer_class = PreguntaExamenTeoricoSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'delete',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user
        rol = user.rol_nombre if hasattr(user, 'rol_nombre') else ""

        queryset = PreguntaExamenTeorico.objects.prefetch_related(
            'opciones'
        ).all().order_by('-fecha_creacion')

        activa = self.request.query_params.get('activa')

        if activa is not None:
            if activa.lower() == 'true':
                queryset = queryset.filter(activa=True)
            elif activa.lower() == 'false':
                queryset = queryset.filter(activa=False)

        if es_admin(user):
            return queryset

        return queryset.none()

    def perform_create(self, serializer):
        user = self.request.user

        if not es_admin(user):
            raise serializers.ValidationError(
                'Solo administración puede crear preguntas del examen teórico.'
            )

        serializer.save()

    def perform_update(self, serializer):
        user = self.request.user

        if not es_admin(user):
            raise serializers.ValidationError(
                (
                    'Solo administración puede editar '
                    'preguntas del examen teórico.'
                )
            )

        pregunta = serializer.instance

        tiene_intento_activo = (
            pregunta.asignaciones_examen
            .filter(
                intento__estado__in=[
                    'habilitado',
                    'iniciado',
                ]
            )
            .exists()
        )

        if tiene_intento_activo:
            raise serializers.ValidationError(
                {
                    'error': (
                        'Esta pregunta pertenece a un examen que un estudiante todavía tiene activo. '
                        'No puede editarse hasta que el intento sea enviado.'
                    )
                }
            )

        serializer.save()

    def destroy(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {
                    'error': (
                        'Solo administración puede '
                        'eliminar preguntas.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN,
            )

        pregunta = self.get_object()

        fue_utilizada = (
            pregunta.asignaciones_examen
            .exists()
        )

        if fue_utilizada:
            if pregunta.activa:
                pregunta.activa = False
                pregunta.save(
                    update_fields=[
                        'activa',
                    ]
                )

            return Response(
                {
                    'message': (
                        'La pregunta ya fue utilizada en un '
                        'examen. Se retiró de los exámenes nuevos, '
                        'pero se conservó en el historial.'
                    ),
                    'desactivada': True,
                },
                status=status.HTTP_200_OK,
            )

        pregunta.delete()

        return Response(
            {
                'message': (
                    'La pregunta y sus opciones fueron '
                    'eliminadas correctamente.'
                ),
                'desactivada': False,
            },
            status=status.HTTP_200_OK,
        )

class ExamenTeoricoViewSet(viewsets.ReadOnlyModelViewSet):
    queryset = ExamenTeorico.objects.select_related(
        'matricula',
        'matricula__estudiante',
        'habilitado_por',
    ).all()

    serializer_class = ExamenTeoricoSerializer
    permission_classes = [IsAuthenticated]
    pagination_class = PaginacionOpcional
    http_method_names = [
        'get',
        'post',
        'head',
        'options',
    ]

    def get_queryset(self):
        user = self.request.user

        queryset = (
            ExamenTeorico.objects
            .select_related(
                'matricula',
                'matricula__estudiante',
                'habilitado_por',
            )
            .all()
            .order_by('-id')
        )

        if es_admin(user):
            return queryset

        if es_instructor(user):
            return queryset.filter(
                habilitado_por_id=user.instructor_id
            )

        if es_estudiante(user):
            return queryset.filter(
                matricula__estudiante_id=user.estudiante_id
            )

        return queryset.none()

    def create(self, request, *args, **kwargs):
        return Response(
            {
                'error': (
                    'Para habilitar un examen usa '
                    '/api/examen-teorico/habilitar/.'
                )
            },
            status=status.HTTP_405_METHOD_NOT_ALLOWED
        )

    def obtener_ultima_clase_practica(self, matricula):
        return (
            Calendario.objects
            .filter(
                matricula=matricula,
                es_examen=False,
            )
            .exclude(
                estado='cancelada'
            )
            .order_by(
                '-fecha',
                '-hora_fin',
                '-id',
            )
            .first()
        )

    def obtener_instructor_actual_id(self, matricula):
        ultima_clase = self.obtener_ultima_clase_practica(
            matricula
        )

        if not ultima_clase:
            return None

        return ultima_clase.instructor_id

    def obtener_ordenes_previos(self, examen):
        ordenes = []

        intentos = (
            IntentoExamenTeorico.objects
            .filter(
                examen=examen
            )
            .prefetch_related(
                'preguntas_asignadas'
            )
            .order_by(
                'numero_intento'
            )
        )

        for intento in intentos:
            ids = list(
                intento.preguntas_asignadas
                .order_by('orden')
                .values_list(
                    'pregunta_id',
                    flat=True,
                )
            )

            if ids:
                ordenes.append(
                    tuple(ids)
                )

        return set(ordenes)

    def seleccionar_orden_preguntas(
        self,
        examen,
        total_requerido,
    ):
        ids_preguntas = list(
            PreguntaExamenTeorico.objects
            .filter(
                activa=True
            )
            .values_list(
                'id',
                flat=True,
            )
        )

        if not ids_preguntas:
            return []

        total = min(
            len(ids_preguntas),
            total_requerido,
        )

        ids_preguntas = ids_preguntas[:]
        ordenes_previos = self.obtener_ordenes_previos(
            examen
        )

        if total <= 1 and ordenes_previos:
            raise ValueError(
                'No hay suficientes preguntas activas para '
                'generar un orden diferente al intento anterior.'
            )

        for _ in range(50):
            random.shuffle(ids_preguntas)

            seleccion = tuple(
                ids_preguntas[:total]
            )

            if seleccion not in ordenes_previos:
                return list(seleccion)

        seleccion_base = ids_preguntas[:total]

        for desplazamiento in range(1, total):
            seleccion = tuple(
                seleccion_base[desplazamiento:]
                + seleccion_base[:desplazamiento]
            )

            if seleccion not in ordenes_previos:
                return list(seleccion)

        raise ValueError(
            'No fue posible generar un orden de preguntas diferente para este estudiante.'
        )

    def crear_intento_habilitado(
        self,
        examen,
        total_requerido=30,
    ):
        ultimo_intento = (
            IntentoExamenTeorico.objects
            .filter(
                examen=examen
            )
            .order_by(
                '-numero_intento'
            )
            .first()
        )

        if (
            ultimo_intento
            and ultimo_intento.estado in [
                'habilitado',
                'iniciado',
            ]
        ):
            return ultimo_intento

        if (
            ultimo_intento
            and ultimo_intento.estado == 'realizado'
            and ultimo_intento.nota is not None
            and ultimo_intento.nota >= 80
        ):
            raise ValueError(
                'El estudiante ya aprobó el examen teórico.'
            )

        siguiente_numero = (
            ultimo_intento.numero_intento + 1
            if ultimo_intento
            else 1
        )

        ids_preguntas = self.seleccionar_orden_preguntas(
            examen,
            total_requerido,
        )

        if not ids_preguntas:
            raise ValueError(
                'No existen preguntas activas para el examen teórico.'
            )

        intento = IntentoExamenTeorico.objects.create(
            examen=examen,
            numero_intento=siguiente_numero,
            estado='habilitado',
            fecha_habilitado=timezone.now(),
        )

        asignaciones = [
            PreguntaIntentoExamenTeorico(
                intento=intento,
                pregunta_id=pregunta_id,
                orden=indice,
            )
            for indice, pregunta_id in enumerate(
                ids_preguntas,
                start=1,
            )
        ]

        PreguntaIntentoExamenTeorico.objects.bulk_create(
            asignaciones
        )

        return intento

    def obtener_intento_actual(self, examen):
        return (
            IntentoExamenTeorico.objects
            .filter(
                examen=examen,
                estado__in=[
                    'habilitado',
                    'iniciado',
                ],
            )
            .order_by(
                '-numero_intento'
            )
            .first()
        )

    def obtener_preguntas_del_intento(self, intento):
        asignaciones = (
            PreguntaIntentoExamenTeorico.objects
            .select_related(
                'pregunta'
            )
            .prefetch_related(
                'pregunta__opciones'
            )
            .filter(
                intento=intento
            )
            .order_by(
                'orden'
            )
        )

        return [
            asignacion.pregunta
            for asignacion in asignaciones
        ]

    @action(detail=False, methods=['post'], url_path='habilitar')
    def habilitar(self, request):
        user = request.user

        if not es_instructor(user):
            return Response(
                {
                    'error': (
                        'Solo un instructor puede habilitar '
                        'el examen teórico.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        matricula_id = request.data.get('matricula_id')

        if not matricula_id:
            return Response(
                {
                    'error': (
                        'Debe enviar la matrícula del estudiante.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        try:
            matricula = (
                Matricula.objects
                .select_related(
                    'estudiante',
                    'plan_de_estudio',
                )
                .get(
                    id=matricula_id
                )
            )
        except Matricula.DoesNotExist:
            return Response(
                {
                    'error': 'Matrícula no encontrada.'
                },
                status=status.HTTP_404_NOT_FOUND
            )

        if matricula.estado not in [
            'matriculado',
            'finalizado',
        ]:
            return Response(
                {
                    'error': (
                        'La matrícula debe encontrarse activa o finalizada para habilitar el examen.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        instructor_actual_id = (
            self.obtener_instructor_actual_id(
                matricula
            )
        )

        if not instructor_actual_id:
            return Response(
                {
                    'error': (
                        'La matrícula todavía no tiene clases prácticas asignadas.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        if instructor_actual_id != user.instructor_id:
            return Response(
                {
                    'error': (
                        'Solo el instructor actualmente asignado al estudiante puede habilitar el examen.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        preguntas_disponibles = (
            PreguntaExamenTeorico.objects
            .filter(
                activa=True
            )
            .count()
        )

        if preguntas_disponibles == 0:
            return Response(
                {
                    'error': (
                        'No existen preguntas activas para el examen teórico.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        with transaction.atomic():
            examen, created = (
                ExamenTeorico.objects
                .select_for_update()
                .get_or_create(
                    matricula=matricula,
                    defaults={
                        'habilitado_por_id': user.instructor_id,
                        'estado': 'habilitado',
                        'fecha_habilitado': timezone.now(),
                    }
                )
            )

            if (
                not created
                and examen.estado == 'realizado'
                and examen.nota is not None
                and examen.nota >= 80
            ):
                return Response(
                    {
                        'error': (
                            'El estudiante ya aprobó el examen teórico.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            examen.estado = 'habilitado'
            examen.nota = None
            examen.fecha_realizado = None
            examen.fecha_habilitado = timezone.now()
            examen.habilitado_por_id = user.instructor_id
            examen.save(
                update_fields=[
                    'estado',
                    'nota',
                    'fecha_realizado',
                    'fecha_habilitado',
                    'habilitado_por',
                ]
            )

            try:
                intento = self.crear_intento_habilitado(
                    examen
                )
            except ValueError as error:
                return Response(
                    {
                        'error': str(error)
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

        serializer = self.get_serializer(
            examen
        )

        return Response(
            {
                'message': (
                    'Examen teórico habilitado correctamente.'
                ),
                'examen': serializer.data,
                'intento': intento.numero_intento,
            },
            status=status.HTTP_200_OK
        )

    @action(
        detail=False,
        methods=['get'],
        url_path='mi-examen',
    )
    def mi_examen(self, request):
        user = request.user

        if not es_estudiante(user):
            return Response(
                {
                    'error': (
                        'Solo el estudiante puede consultar su examen teórico.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        matricula = (
            Matricula.objects
            .select_related(
                'estudiante',
                'plan_de_estudio',
            )
            .filter(
                estudiante_id=user.estudiante_id,
                estado__in=[
                    'matriculado',
                    'finalizado',
                ]
            )
            .order_by(
                '-id'
            )
            .first()
        )

        if not matricula:
            return Response(
                {
                    'disponible': False,
                    'message': (
                        'No tienes una matrícula activa para realizar el examen teórico.'
                    ),
                },
                status=status.HTTP_200_OK
            )

        tipo_curso = str(
            matricula.tipo_curso or ''
        ).strip().lower()

        examen = (
            ExamenTeorico.objects
            .select_related(
                'matricula',
                'matricula__estudiante',
                'habilitado_por',
            )
            .filter(
                matricula=matricula
            )
            .first()
        )

        if tipo_curso in [
            'intermedio',
            'avanzado',
        ]:
            instructor_id = (
                self.obtener_instructor_actual_id(
                    matricula
                )
            )

            if not examen:
                examen = ExamenTeorico.objects.create(
                    matricula=matricula,
                    habilitado_por_id=instructor_id,
                    estado='habilitado',
                    fecha_habilitado=timezone.now(),
                )

            elif (
                examen.estado == 'pendiente'
                or (
                    examen.estado == 'realizado'
                    and examen.nota is not None
                    and examen.nota < 80
                )
            ):
                examen.estado = 'habilitado'
                examen.nota = None
                examen.fecha_realizado = None
                examen.fecha_habilitado = timezone.now()

                campos_actualizados = [
                    'estado',
                    'nota',
                    'fecha_realizado',
                    'fecha_habilitado',
                ]

                if instructor_id:
                    examen.habilitado_por_id = instructor_id
                    campos_actualizados.append(
                        'habilitado_por'
                    )

                examen.save(
                    update_fields=campos_actualizados
                )

        if not examen:
            return Response(
                {
                    'disponible': False,
                    'message': (
                        'Todavía no tienes examen teórico habilitado.'
                    ),
                },
                status=status.HTTP_200_OK
            )

        if (
            examen.estado == 'realizado'
            and examen.nota is not None
            and examen.nota >= 80
        ):
            return Response(
                {
                    'disponible': False,
                    'realizado': True,
                    'message': (
                        'Ya aprobaste el examen teórico.'
                    ),
                    'examen': self.get_serializer(
                        examen
                    ).data,
                },
                status=status.HTTP_200_OK
            )

        if examen.estado != 'habilitado':
            return Response(
                {
                    'disponible': False,
                    'realizado': False,
                    'message': (
                        'Tu examen teórico todavía no está disponible.'
                    ),
                },
                status=status.HTTP_200_OK
            )

        with transaction.atomic():
            examen = (
                ExamenTeorico.objects
                .select_for_update()
                .get(
                    id=examen.id
                )
            )

            try:
                intento = self.crear_intento_habilitado(
                    examen
                )
            except ValueError as error:
                return Response(
                    {
                        'disponible': False,
                        'message': str(error),
                    },
                    status=status.HTTP_200_OK
                )

            if intento.estado == 'habilitado':
                intento.estado = 'iniciado'
                intento.fecha_iniciado = timezone.now()

                intento.save(
                    update_fields=[
                        'estado',
                        'fecha_iniciado',
                    ]
                )

        preguntas = self.obtener_preguntas_del_intento(
            intento
        )

        preguntas_serializer = (
            PreguntaExamenEstudianteSerializer(
                preguntas,
                many=True
            )
        )

        return Response(
            {
                'disponible': True,
                'realizado': False,
                'tipo_curso': matricula.tipo_curso,
                'matricula_id': matricula.id,
                'examen': self.get_serializer(
                    examen
                ).data,
                'intento': intento.numero_intento,
                'preguntas': preguntas_serializer.data,
            },
            status=status.HTTP_200_OK
        )

    @action(
        detail=True,
        methods=['post'],
        url_path='enviar',
    )
    def enviar(self, request, pk=None):
        user = request.user

        if not es_estudiante(user):
            return Response(
                {
                    'error': (
                        'Solo el estudiante puede enviar el examen teórico.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        examen = self.get_object()

        if (
            examen.matricula.estudiante_id
            != user.estudiante_id
        ):
            return Response(
                {
                    'error': (
                        'No puedes enviar un examen que no te pertenece.'
                    )
                },
                status=status.HTTP_403_FORBIDDEN
            )

        if examen.estado != 'habilitado':
            return Response(
                {
                    'error': (
                        'Este examen no está habilitado o ya fue realizado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        serializer = RespuestaEnviarExamenSerializer(
            data=request.data
        )

        serializer.is_valid(
            raise_exception=True
        )

        respuestas = serializer.validated_data[
            'respuestas'
        ]

        respuestas_ids = {
            respuesta['pregunta_id']
            for respuesta in respuestas
        }

        intento = self.obtener_intento_actual(
            examen
        )

        if not intento:
            return Response(
                {
                    'error': (
                        'No existe un intento activo para este examen.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        asignaciones = list(
            PreguntaIntentoExamenTeorico.objects
            .select_related(
                'pregunta'
            )
            .filter(
                intento=intento
            )
            .order_by(
                'orden'
            )
        )

        total_requerido = len(
            asignaciones
        )

        if total_requerido == 0:
            return Response(
                {
                    'error': (
                        'El intento no tiene preguntas asignadas.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        preguntas_asignadas_ids = {
            asignacion.pregunta_id
            for asignacion in asignaciones
        }

        if respuestas_ids != preguntas_asignadas_ids:
            return Response(
                {
                    'error': (
                        'Las respuestas enviadas no coinciden con las preguntas asignadas al intento.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        opciones = (
            OpcionPreguntaExamenTeorico.objects
            .filter(
                id__in=[
                    respuesta['opcion_id']
                    for respuesta in respuestas
                ],
                pregunta_id__in=preguntas_asignadas_ids,
            )
            .select_related(
                'pregunta'
            )
        )

        opciones_por_id = {
            opcion.id: opcion
            for opcion in opciones
        }

        respuestas_preparadas = []

        for respuesta in respuestas:
            pregunta_id = respuesta['pregunta_id']
            opcion_id = respuesta['opcion_id']
            opcion = opciones_por_id.get(
                opcion_id
            )

            if (
                not opcion
                or opcion.pregunta_id != pregunta_id
            ):
                return Response(
                    {
                        'error': (
                            'Una opción no pertenece a la pregunta indicada.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            respuestas_preparadas.append({
                'pregunta_id': pregunta_id,
                'opcion': opcion,
            })

        plan = examen.matricula.plan_de_estudio

        if not plan:
            return Response(
                {
                    'error': (
                        'La matrícula no tiene plan '
                        'de estudio asignado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )
        instructor = examen.habilitado_por

        if not instructor:
            instructor_id = self.obtener_instructor_actual_id(
                examen.matricula
            )

            if instructor_id:
                instructor = (
                    Instructor.objects
                    .filter(
                        id=instructor_id
                    )
                    .first()
                )

        if not instructor:
            return Response(
                {
                    'error': (
                        'No se puede registrar la nota porque la matrícula todavía no '
                        'tiene instructor asignado.'
                    )
                },
                status=status.HTTP_400_BAD_REQUEST
            )

        with transaction.atomic():
            examen = (
                ExamenTeorico.objects
                .select_for_update()
                .get(
                    id=examen.id
                )
            )

            intento = (
                IntentoExamenTeorico.objects
                .select_for_update()
                .filter(
                    examen=examen,
                    estado__in=[
                        'habilitado',
                        'iniciado',
                    ],
                )
                .order_by(
                    '-numero_intento'
                )
                .first()
            )

            if not intento:
                return Response(
                    {
                        'error': (
                            'El intento ya fue enviado '
                            'o dejó de estar disponible.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            if examen.estado != 'habilitado':
                return Response(
                    {
                        'error': (
                            'El examen ya fue enviado '
                            'o dejó de estar habilitado.'
                        )
                    },
                    status=status.HTTP_400_BAD_REQUEST
                )

            correctas = 0
            respuestas_crear = []

            for item in respuestas_preparadas:
                opcion = item['opcion']
                es_correcta = opcion.es_correcta

                if es_correcta:
                    correctas += 1

                respuestas_crear.append(
                    RespuestaExamenTeorico(
                        examen=examen,
                        intento=intento,
                        pregunta_id=item['pregunta_id'],
                        opcion_seleccionada=opcion,
                        correcta=es_correcta,
                    )
                )

            RespuestaExamenTeorico.objects.bulk_create(
                respuestas_crear
            )

            nota_final = round(
                (
                    correctas
                    / total_requerido
                ) * 100,
                2
            )

            intento.estado = 'realizado'
            intento.nota = nota_final
            intento.fecha_realizado = timezone.now()

            intento.save(
                update_fields=[
                    'estado',
                    'nota',
                    'fecha_realizado',
                ]
            )

            examen.estado = 'realizado'
            examen.nota = nota_final
            examen.fecha_realizado = timezone.now()
            examen.habilitado_por = instructor

            examen.save(
                update_fields=[
                    'estado',
                    'nota',
                    'fecha_realizado',
                    'habilitado_por',
                ]
            )

            comentario_nota = (
                'Examen teórico aprobado automáticamente por el sistema.'
                if nota_final >= 80
                else (
                    'Examen teórico reprobado. '
                    'Puede realizarlo nuevamente.'
                )
            )

            nota_teorica = (
                Notas.objects
                .select_for_update()
                .filter(
                    matricula=examen.matricula,
                    tipo_nota='teorico',
                )
                .order_by(
                    '-fecha_registro',
                    '-id',
                )
                .first()
            )

            if nota_teorica:
                nota_teorica.instructor = instructor
                nota_teorica.plan_de_estudio = plan
                nota_teorica.nota = f'{nota_final:.2f}'
                nota_teorica.comentario = comentario_nota

                nota_teorica.save(
                    update_fields=[
                        'instructor',
                        'plan_de_estudio',
                        'nota',
                        'comentario',
                    ]
                )
            else:
                nota_teorica = Notas.objects.create(
                    matricula=examen.matricula,
                    tipo_nota='teorico',
                    instructor=instructor,
                    plan_de_estudio=plan,
                    nota=f'{nota_final:.2f}',
                    comentario=comentario_nota,
                )

            matricula_finalizada = (
                actualizar_estado_matricula_por_notas(
                    nota_teorica.matricula
                )
            )

            if matricula_finalizada:
                desactivar_usuarios_estudiante(
                    nota_teorica.matricula.estudiante
                )

        return Response(
            {
                'message': (
                    'Examen enviado y calificado correctamente.'
                ),
                'intento': intento.numero_intento,
                'total_preguntas': total_requerido,
                'correctas': correctas,
                'nota': nota_final,
                'resultado': (
                    'Aprobado'
                    if nota_final >= 80
                    else 'Reprobado'
                ),
            },
            status=status.HTTP_200_OK
        )

    @action(detail=True, methods=['get'], url_path='respuestas')
    def respuestas(self, request, pk=None):
        examen = self.get_object()

        intento = (
            IntentoExamenTeorico.objects
            .filter(
                examen=examen,
                estado='realizado',
            )
            .order_by(
                '-numero_intento'
            )
            .first()
        )

        respuestas = (
            RespuestaExamenTeorico.objects
            .select_related(
                'pregunta',
                'opcion_seleccionada',
                'intento',
            )
            .filter(
                examen=examen
            )
        )

        if intento:
            respuestas = respuestas.filter(
                intento=intento
            )

        serializer = RespuestaExamenTeoricoSerializer(
            respuestas,
            many=True
        )

        return Response(
            serializer.data
        )

class PerfilView(APIView):
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def get_foto_url(self, request, instructor):
        return obtener_foto_instructor(instructor)

    def serializar_instructor(self, request, instructor, incluir_foto=True):
        categoria = instructor.categoria_instructor or ""

        data = {
            "id": instructor.id,
            "nombre": instructor.nombre,
            "apellido": instructor.apellido,
            "telefono": instructor.numero_telefono or "",
            "direccion": instructor.direccion or "",
            "edad": instructor.edad,
            "experiencia": instructor.experiencia or "",
            "categoria": categoria,
            "categoria_instructor": categoria,
            "cedula": instructor.cedula or "",
            "nacionalidad": instructor.nacionalidad or "",
            "centro_trabajo": instructor.centro_trabajo or "",
            "cargo": instructor.cargo or "",
        }

        if incluir_foto:
            data["foto"] = self.get_foto_url(request, instructor)
        else:
            data["foto"] = None

        return data

    def serializar_estudiante(self, estudiante):
        return {
            "id": estudiante.id,
            "nombre": estudiante.nombre,
            "apellido": estudiante.apellido,
            "cedula": estudiante.cedula,
            "telefono": estudiante.telefono_movil,
            "direccion": estudiante.direccion,
            "edad": estudiante.edad,
            "nivel_educativo": estudiante.nivel_educativo,
            "nacionalidad": estudiante.nacionalidad or "",
            "nombre_emergencia": estudiante.nombre_emergencia or "",
            "telefono_emergencia": estudiante.telefono_emergencia or "",
            "activo": estudiante.activo,
        }

    def get(self, request):
        usuario = request.user
        rol = usuario.rol_nombre

        data = {
            "rol": rol,
            "mi_perfil": None,
            "instructor": None,
            "estudiantes": [],
            "instructores": [],
            "matricula_activa": False,
        }

        if rol in ['admin', 'administrador', 'secretaria']:
            instructores = Instructor.objects.all().order_by('-id')

            matriculas_activas = Matricula.objects.exclude(
                id__in=obtener_ids_matriculas_egresadas()
            )

            estudiantes = Estudiante.objects.filter(
                id__in=matriculas_activas.values(
                    'estudiante_id'
                )
            ).distinct().order_by('-id')

            data["instructores"] = [
                self.serializar_instructor(
                    request,
                    instructor,
                    incluir_foto=False
                )
                for instructor in instructores
            ]

            data["estudiantes"] = [
                self.serializar_estudiante(estudiante)
                for estudiante in estudiantes
            ]
            return Response(data)

        if rol == 'instructor':
            instructor = usuario.instructor

            if instructor:
                data["mi_perfil"] = self.serializar_instructor(request, instructor)

                matriculas_activas = Matricula.objects.exclude(
                    id__in=obtener_ids_matriculas_egresadas()
                )

                estudiantes_ids = Calendario.objects.filter(
                    instructor=instructor,
                    es_examen=False,
                    matricula_id__in=matriculas_activas.values(
                        'id'
                    ),
                ).exclude(
                    estado='cancelada'
                ).values_list(
                    'matricula__estudiante_id',
                    flat=True
                ).distinct()

                estudiantes = Estudiante.objects.filter(id__in=estudiantes_ids)

                data["estudiantes"] = [
                    self.serializar_estudiante(estudiante)
                    for estudiante in estudiantes
                ]
            return Response(data)

        if rol == 'estudiante':
            estudiante = usuario.estudiante

            if estudiante:
                data["mi_perfil"] = self.serializar_estudiante(estudiante)

                matricula_activa = (
                    Matricula.objects
                    .filter(
                        estudiante=estudiante
                    )
                    .exclude(
                        estado='finalizado'
                    )
                    .order_by(
                        '-id'
                    )
                    .first()
                )

                data["matricula_activa"] = bool(matricula_activa)

                if matricula_activa:
                    clase = (
                        Calendario.objects
                        .filter(
                            matricula=matricula_activa,
                            es_examen=False,
                        )
                        .exclude(
                            estado='cancelada'
                        )
                        .select_related(
                            'instructor'
                        )
                        .order_by(
                            'fecha',
                            'hora_inicio',
                            'id',
                        )
                        .first()
                    )

                    if clase and clase.instructor:
                        data["instructor"] = self.serializar_instructor(
                            request,
                            clase.instructor
                        )
            return Response(data)
        return Response(data)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def exportar_reporte_instructores_policial(request):
    if not es_admin(request.user):
        return Response(
            {
                'error': (
                    'Solo Administración o Secretaría pueden '
                    'generar el reporte policial de instructores.'
                )
            },
            status=status.HTTP_403_FORBIDDEN
        )

    fecha_desde = request.query_params.get('desde')
    fecha_hasta = request.query_params.get('hasta')
    fecha_desde_parseada = (
        parse_date(fecha_desde)
        if fecha_desde
        else None
    )

    fecha_hasta_parseada = (
        parse_date(fecha_hasta)
        if fecha_hasta
        else None
    )

    if fecha_desde and not fecha_desde_parseada:
        return Response(
            {
                'error': (
                    'La fecha inicial no tiene un formato válido. '
                    'Utilice AAAA-MM-DD.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    if fecha_hasta and not fecha_hasta_parseada:
        return Response(
            {
                'error': (
                    'La fecha final no tiene un formato válido. '
                    'Utilice AAAA-MM-DD.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    if (
        fecha_desde_parseada
        and fecha_hasta_parseada
        and fecha_desde_parseada > fecha_hasta_parseada
    ):
        return Response(
            {
                'error': (
                    'La fecha inicial no puede ser posterior '
                    'a la fecha final.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    zona_horaria = timezone.get_current_timezone()
    inicio_rango = None
    fin_rango = None

    if fecha_desde_parseada:
        inicio_rango = timezone.make_aware(
            datetime.combine(
                fecha_desde_parseada,
                datetime.min.time()
            ),
            zona_horaria
        )

    if fecha_hasta_parseada:
        fin_rango = timezone.make_aware(
            datetime.combine(
                fecha_hasta_parseada + timedelta(days=1),
                datetime.min.time()
            ),
            zona_horaria
        )

    ruta_plantilla = os.path.join(
        settings.BASE_DIR,
        'app_escuela',
        'plantilla',
        'INFORME TRANSITO POLICIA NAC.xlsm'
    )

    if not os.path.exists(ruta_plantilla):
        return Response(
            {'error': f'No existe la plantilla: {ruta_plantilla}'},
            status=404
        )

    wb = load_workbook(
        ruta_plantilla,
        keep_vba=True
    )

    def copiar_estilo_fila(ws, fila_origen, fila_destino, columnas):
        for columna in range(1, columnas + 1):
            origen = ws.cell(row=fila_origen, column=columna)
            destino = ws.cell(row=fila_destino, column=columna)

            if type(destino).__name__ == 'MergedCell':
                continue

            if origen.has_style:
                destino._style = copy(origen._style)

            destino.font = copy(origen.font)
            destino.fill = copy(origen.fill)
            destino.border = copy(origen.border)
            destino.alignment = copy(origen.alignment)
            destino.number_format = origen.number_format
            destino.protection = copy(origen.protection)

        ws.row_dimensions[fila_destino].height = ws.row_dimensions[fila_origen].height

    def aplicar_estilo_tabla_manual(ws, fila, columna_inicio, columna_fin, es_par):
        azul = "B8CCE4"
        azul_claro = "DCE6F1"
        blanco = "FFFFFF"

        fill = PatternFill(
            fill_type="solid",
            fgColor=azul if es_par else azul_claro
        )

        borde = Side(style="thin", color="FFFFFF")

        for columna in range(columna_inicio, columna_fin + 1):
            celda = ws.cell(row=fila, column=columna)

            if type(celda).__name__ == "MergedCell":
                continue

            celda.fill = fill
            celda.border = Border(
                left=borde,
                right=borde,
                top=borde,
                bottom=borde
            )
            celda.alignment = Alignment(
                horizontal="center",
                vertical="center",
                wrap_text=True
            )

    nombre_hoja = 'LISTADO INSTRUCTORES'

    if nombre_hoja not in wb.sheetnames:
        return Response(
            {'error': f'No existe la hoja {nombre_hoja}. Hojas disponibles: {wb.sheetnames}'},
            status=400
        )

    ws = wb[nombre_hoja]
    fila_inicio = 5
    fila = fila_inicio
    imagenes_conservar = []
    archivos_temporales = []

    for imagen_excel in ws._images:
        try:
            columna_imagen = imagen_excel.anchor._from.col + 1
            fila_imagen = imagen_excel.anchor._from.row + 1

            if columna_imagen == 2 and fila_imagen >= fila_inicio:
                continue

            imagenes_conservar.append(imagen_excel)

        except Exception:
            imagenes_conservar.append(imagen_excel)

    ws._images = imagenes_conservar

    while fila <= ws.max_row:

        for columna in range(1, 18):
            if columna != 2:
                ws.cell(row=fila, column=columna).value = None

        fila += 1

    instructores = Instructor.objects.order_by(
        'nombre',
        'apellido'
    )

    fila = fila_inicio

    for index, instructor in enumerate(instructores, start=1):
        copiar_estilo_fila(
            ws,
            fila_inicio,
            fila,
            17
        )

        ws.row_dimensions[fila].height = 65

        nombre_completo = (
            f"{instructor.nombre or ''} "
            f"{instructor.apellido or ''}"
        ).strip()

        ws.cell(row=fila, column=1, value=index)
        ruta_foto = obtener_ruta_foto_instructor_para_excel(instructor)

        if ruta_foto:
            archivos_temporales.append(ruta_foto)

            try:
                imagen = ExcelImage(ruta_foto)
                imagen.width = 42
                imagen.height = 42

                ws.row_dimensions[fila].height = 48
                ws.column_dimensions['B'].width = 10
                ws.add_image(imagen, f'B{fila}')

            except Exception:
                pass

        ws.cell(row=fila, column=3, value=nombre_completo)

        ws.cell(
            row=fila,
            column=4,
            value=instructor.cedula or ""
        )

        ws.cell(
            row=fila,
            column=5,
            value=instructor.nacionalidad or ""
        )

        ws.cell(
            row=fila,
            column=6,
            value=instructor.direccion or ""
        )

        ws.cell(
            row=fila,
            column=7,
            value=instructor.numero_telefono or ""
        )

        ws.cell(
            row=fila,
            column=8,
            value=instructor.nivel_escolar or ""
        )

        ws.cell(
            row=fila,
            column=9,
            value=instructor.categoria_instructor or ""
        )

        ws.cell(
            row=fila,
            column=10,
            value=instructor.antecedentes_penales or ""
        )

        ws.cell(
            row=fila,
            column=11,
            value=instructor.centro_trabajo or ""
        )

        ws.cell(
            row=fila,
            column=12,
            value=instructor.cargo or ""
        )

        ws.cell(
            row=fila,
            column=13,
            value=instructor.curso_aprobado_instructor or ""
        )

        ws.cell(
            row=fila,
            column=14,
            value=(
                instructor.fecha_ingreso.strftime('%d/%m/%Y')
                if instructor.fecha_ingreso else ""
            )
        )

        ws.cell(
            row=fila,
            column=15,
            value=(
                instructor.fecha_salida.strftime('%d/%m/%Y')
                if instructor.fecha_salida else ""
            )
        )

        ws.cell(
            row=fila,
            column=16,
            value=instructor.motivo_salida or ""
        )

        ws.cell(
            row=fila,
            column=17,
            value=instructor.infracciones_resoluciones or ""
        )

        fila += 1

    nombre_hoja_ingresos = 'REPORTE DE INGRESOS'

    if nombre_hoja_ingresos not in wb.sheetnames:
        return Response(
            {'error': f'No existe la hoja {nombre_hoja_ingresos}. Hojas disponibles: {wb.sheetnames}'},
            status=400
        )

    ws_ingresos = wb[nombre_hoja_ingresos]

    texto_fecha = ""

    if fecha_desde and fecha_hasta:
        texto_fecha = (
            f"Desde {fecha_desde} hasta {fecha_hasta}"
        )

    elif fecha_desde:
        texto_fecha = (
            f"Desde {fecha_desde}"
        )

    elif fecha_hasta:
        texto_fecha = (
            f"Hasta {fecha_hasta}"
        )

    ws_ingresos.cell(row=3, column=2, value=timezone.now().strftime('%d/%m/%Y'))

    fila_inicio_ingresos = 5
    fila_modelo_ingresos_azul = 5
    fila_modelo_ingresos_clara = 6
    fila = fila_inicio_ingresos

    while fila <= ws_ingresos.max_row:
        for columna in range(1, 16):
            ws_ingresos.cell(row=fila, column=columna).value = None

        fila += 1

    matriculas = (
        Matricula.objects
        .select_related(
            "estudiante",
            "categoria",
        )
        .annotate(
            fecha_finalizacion_reporte=models.Max(
                "clases__fecha",
                filter=(
                    Q(clases__es_examen=False)
                    & ~Q(clases__estado="cancelada")
                ),
            )
        )
    )

    if inicio_rango:
        matriculas = matriculas.filter(
            fecha_registro__gte=inicio_rango
        )

    if fin_rango:
        matriculas = matriculas.filter(
            fecha_registro__lt=fin_rango
        )

    matriculas = matriculas.order_by(
        'fecha_registro',
        'id'
    )

    fila = fila_inicio_ingresos

    for matricula in matriculas:
        fila_modelo = (
            fila_modelo_ingresos_azul
            if (fila - fila_inicio_ingresos) % 2 == 0
            else fila_modelo_ingresos_clara
        )

        copiar_estilo_fila(
            ws_ingresos,
            fila_modelo,
            fila,
            15
        )

        aplicar_estilo_tabla_manual(
            ws_ingresos,
            fila,
            1,
            15,
            (fila - fila_inicio_ingresos) % 2 == 0
        )

        estudiante = matricula.estudiante

        fecha_finalizacion = (
            matricula.fecha_finalizacion_reporte
        )

        if matricula.tipo_curso == "Principiante":
            horas_practicas = 15
        else:
            horas_practicas = matricula.horas_reforzamiento or 0

        horas_totales = round(float(horas_practicas) / 0.6)
        horas_teoricas = round(
            horas_totales - float(horas_practicas)
        )

        ws_ingresos.cell(
            row=fila,
            column=1,
            value=estudiante.codigo_estudiante or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=2,
            value=f"{estudiante.nombre or ''} {estudiante.apellido or ''}".strip()
        )

        ws_ingresos.cell(
            row=fila,
            column=3,
            value=estudiante.nacionalidad or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=4,
            value=estudiante.cedula or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=5,
            value=estudiante.direccion or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=6,
            value=estudiante.telefono_movil or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=7,
            value=estudiante.nivel_educativo or ""
        )

        ws_ingresos.cell(
            row=fila,
            column=8,
            value="X" if matricula.tipo_curso == "Principiante" else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=9,
            value="X" if matricula.tipo_curso == "Intermedio" else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=10,
            value="X" if matricula.tipo_curso == "Avanzado" else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=11,
            value=matricula.categoria.nombre if matricula.categoria else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=12,
            value=matricula.fecha_registro.strftime('%d/%m/%Y') if matricula.fecha_registro else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=13,
            value=fecha_finalizacion.strftime('%d/%m/%Y') if fecha_finalizacion else ""
        )

        ws_ingresos.cell(
            row=fila,
            column=14,
            value=horas_teoricas
        )

        ws_ingresos.cell(
            row=fila,
            column=15,
            value=horas_practicas
        )

        fila += 1

    ultima_fila_ingresos = fila - 1

    nombre_hoja_egresos = 'REPORTE DE EGRESOS'

    if nombre_hoja_egresos not in wb.sheetnames:
        return Response(
            {'error': f'No existe la hoja {nombre_hoja_egresos}. Hojas disponibles: {wb.sheetnames}'},
            status=400
        )

    ws_egresos = wb[nombre_hoja_egresos]
    ws_egresos.cell(
        row=5,
        column=2,
        value=timezone.now().strftime('%d/%m/%Y')
    )

    fila_inicio_egresos = 8
    fila_estilo_egresos = 8
    fila = fila_inicio_egresos

    while fila <= ws_egresos.max_row:
        for columna in range(1, 20):
            ws_egresos.cell(row=fila, column=columna).value = None

        fila += 1

    egresados = (
        Matricula.objects
        .select_related(
            "estudiante",
            "categoria",
            "plan_de_estudio",
        )
        .prefetch_related(
            Prefetch(
                "notas",
                queryset=(
                    Notas.objects
                    .filter(
                        tipo_nota__in=[
                            "teorico",
                            "practico",
                        ]
                    )
                    .order_by(
                        "-fecha_registro",
                        "-id",
                    )
                ),
                to_attr="notas_reporte_policial",
            )
        )
        .annotate(
            fecha_inicio_reporte=models.Min(
                "clases__fecha",
                filter=(
                    Q(clases__es_examen=False)
                    & ~Q(clases__estado="cancelada")
                ),
            ),
            fecha_egreso_reporte=models.Max(
                "clases__fecha",
                filter=(
                    Q(clases__es_examen=False)
                    & ~Q(clases__estado="cancelada")
                ),
            ),
        )
        .filter(
            tipo_curso__in=[
                "Principiante",
                "Intermedio",
                "Avanzado",
            ],
            fecha_egreso_reporte__isnull=False,
        )
    )

    if fecha_desde_parseada:
        egresados = egresados.filter(
            fecha_egreso_reporte__gte=(
                fecha_desde_parseada
            )
        )

    if fecha_hasta_parseada:
        egresados = egresados.filter(
            fecha_egreso_reporte__lte=(
                fecha_hasta_parseada
            )
        )

    egresados = egresados.order_by(
        "estudiante__apellido",
        "estudiante__nombre",
        "id",
    )

    fila = fila_inicio_egresos

    for matricula in egresados:
        notas_por_tipo = {}

        for nota in matricula.notas_reporte_policial:
            if nota.tipo_nota not in notas_por_tipo:
                notas_por_tipo[nota.tipo_nota] = nota

        nota_teorica_obj = notas_por_tipo.get(
            "teorico"
        )
        nota_practica_obj = notas_por_tipo.get(
            "practico"
        )

        if not nota_teorica_obj or not nota_practica_obj:
            continue

        nota_teorica = certificado_numero(
            nota_teorica_obj.nota
        )
        nota_practica = certificado_numero(
            nota_practica_obj.nota
        )

        if nota_teorica is None or nota_practica is None:
            continue

        if (
            nota_teorica < Decimal("80")
            or nota_practica < Decimal("80")
        ):
            continue

        copiar_estilo_fila(
            ws_egresos,
            fila_estilo_egresos,
            fila,
            20
        )

        aplicar_estilo_tabla_manual(
            ws_egresos,
            fila,
            1,
            20,
            (fila - fila_inicio_egresos) % 2 == 0
        )

        estudiante = matricula.estudiante

        fecha_inicio = (
            matricula.fecha_inicio_reporte
        )
        fecha_finalizacion = (
            matricula.fecha_egreso_reporte
        )

        if matricula.tipo_curso == "Principiante":
            horas_practicas = 15
        else:
            horas_practicas = matricula.horas_reforzamiento or 0

        horas_totales = round(float(horas_practicas) / 0.6)
        horas_teoricas = round(
            horas_totales - float(horas_practicas)
        )

        ws_egresos.cell(row=fila, column=1, value="")
        ws_egresos.cell(
            row=fila,
            column=2,
            value=f"{estudiante.nombre or ''} {estudiante.apellido or ''}".strip()
        )

        ws_egresos.cell(
            row=fila,
            column=3,
            value=estudiante.nacionalidad or ""
        )

        ws_egresos.cell(
            row=fila,
            column=4,
            value=estudiante.cedula or ""
        )

        ws_egresos.cell(
            row=fila,
            column=5,
            value=estudiante.telefono_movil or ""
        )

        ws_egresos.cell(
            row=fila,
            column=6,
            value=estudiante.nivel_educativo or ""
        )

        ws_egresos.cell(
            row=fila,
            column=7,
            value="X" if matricula.tipo_curso == "Principiante" else ""
        )

        ws_egresos.cell(
            row=fila,
            column=8,
            value="X" if matricula.tipo_curso == "Intermedio" else ""
        )

        ws_egresos.cell(
            row=fila,
            column=9,
            value="X" if matricula.tipo_curso == "Avanzado" else ""
        )

        ws_egresos.cell(
            row=fila,
            column=10,
            value=str(matricula.categoria) if matricula.categoria else ""
        )

        ws_egresos.cell(
            row=fila,
            column=11,
            value=(
                fecha_inicio.strftime("%d/%m/%Y")
                if fecha_inicio
                else ""
            )
        )

        ws_egresos.cell(
            row=fila,
            column=12,
            value=fecha_finalizacion.strftime('%d/%m/%Y') if fecha_finalizacion else ""
        )

        ws_egresos.cell(row=fila, column=13, value=horas_teoricas)
        ws_egresos.cell(row=fila, column=14, value=horas_practicas)
        ws_egresos.cell(row=fila, column=15, value=nota_teorica or "")
        ws_egresos.cell(row=fila, column=16, value=nota_practica or "")
        ws_egresos.cell(row=fila, column=17, value="")
        ws_egresos.cell(row=fila, column=18, value="")
        ws_egresos.cell(row=fila, column=19, value="")
        ws_egresos.cell(row=fila, column=20, value="")

        fila += 1

    ultima_fila_egresos = fila - 1

    response = HttpResponse(
        content_type='application/vnd.ms-excel.sheet.macroEnabled.12'
    )

    response[
        'Content-Disposition'
    ] = 'attachment; filename="INFORME_TRANSITO_POLICIA_NAC.xlsm"'

    try:
        wb.save(response)
    finally:
        for ruta in archivos_temporales:
            try:
                if ruta and os.path.exists(ruta):
                    os.remove(ruta)
            except Exception:
                pass
    return response

class PagoInstructorViewSet(viewsets.ModelViewSet):
    queryset = PagoInstructor.objects.all().order_by('-fecha_inicio', '-id')
    serializer_class = PagoInstructorSerializer
    permission_classes = [IsAuthenticated]

    def get_queryset(self):
        queryset = (
            PagoInstructor.objects
            .all()
            .order_by(
                '-fecha_inicio',
                '-id',
            )
        )

        if es_admin(self.request.user):
            return queryset

        return queryset.none()

    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def _validar_permiso_administrativo(self, request):
        if es_admin(request.user):
            return None

        return Response(
            {
                'error': (
                    'Solo Administración o Secretaría pueden '
                    'registrar, modificar o eliminar pagos '
                    'de instructores.'
                )
            },
            status=status.HTTP_403_FORBIDDEN
        )


    def create(self, request, *args, **kwargs):
        respuesta_permiso = (
            self._validar_permiso_administrativo(
                request
            )
        )

        if respuesta_permiso:
            return respuesta_permiso

        return super().create(
            request,
            *args,
            **kwargs
        )


    def update(self, request, *args, **kwargs):
        respuesta_permiso = (
            self._validar_permiso_administrativo(
                request
            )
        )

        if respuesta_permiso:
            return respuesta_permiso

        return super().update(
            request,
            *args,
            **kwargs
        )

    def partial_update(self, request, *args, **kwargs):
        respuesta_permiso = (
            self._validar_permiso_administrativo(
                request
            )
        )

        if respuesta_permiso:
            return respuesta_permiso

        return super().partial_update(
            request,
            *args,
            **kwargs
        )

    def destroy(self, request, *args, **kwargs):
        respuesta_permiso = (
            self._validar_permiso_administrativo(
                request
            )
        )

        if respuesta_permiso:
            return respuesta_permiso

        return super().destroy(
            request,
            *args,
            **kwargs
        )

    def perform_create(self, serializer):
        activo = serializer.validated_data.get('activo', True)

        if activo:
            PagoInstructor.objects.filter(activo=True).update(
                activo=False,
                fecha_fin=timezone.now().date()
            )
        serializer.save()

    def perform_update(self, serializer):
        activo = serializer.validated_data.get('activo', None)

        if activo is True:
            PagoInstructor.objects.exclude(
                id=serializer.instance.id
            ).filter(
                activo=True
            ).update(
                activo=False,
                fecha_fin=timezone.now().date()
            )
        serializer.save()

class CargoInstitucionalViewSet(viewsets.ModelViewSet):
    queryset = CargoInstitucional.objects.all().order_by('tipo', 'nombre')
    serializer_class = CargoInstitucionalSerializer
    permission_classes = [IsAuthenticated]
    http_method_names = [
        'get',
        'post',
        'put',
        'patch',
        'head',
        'options',
    ]

    def create(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para crear este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )
        return super().create(request, *args, **kwargs)

    def update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )
        return super().update(request, *args, **kwargs)

    def partial_update(self, request, *args, **kwargs):
        if not es_admin(request.user):
            return Response(
                {'error': 'No tienes permiso para editar este registro.'},
                status=status.HTTP_403_FORBIDDEN
            )
        return super().partial_update(request, *args, **kwargs)

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def reporte_induccion_instructores(request):
    if not es_admin(request.user):
        return Response(
            {
                'error': (
                    'Solo Administración o Secretaría pueden generar el reporte de inducción de instructores.'
                )
            },
            status=status.HTTP_403_FORBIDDEN
        )

    fecha_desde = request.query_params.get('desde')
    fecha_hasta = request.query_params.get('hasta')
    instructor_id = request.query_params.get('instructor')

    if not instructor_id:
        return Response(
            {'error': 'Debe seleccionar un instructor.'},
            status=400
        )

    pago_instructor = PagoInstructor.objects.filter(
        activo=True
    ).order_by(
        '-fecha_inicio',
        '-id'
    ).first()

    if not pago_instructor:
        return Response(
            {'error': 'No hay una tarifa activa configurada para el pago por hora del instructor.'},
            status=400
        )

    tarifa_por_hora = pago_instructor.monto_por_alumno

    try:
        instructor = Instructor.objects.get(id=instructor_id)
    except Instructor.DoesNotExist:
        return Response(
            {'error': 'Instructor no encontrado.'},
            status=404
        )

    fecha_desde_parseada = (
        parse_date(fecha_desde)
        if fecha_desde
        else None
    )

    fecha_hasta_parseada = (
        parse_date(fecha_hasta)
        if fecha_hasta
        else None
    )

    if fecha_desde and not fecha_desde_parseada:
        return Response(
            {
                'error': (
                    'La fecha desde no tiene un formato válido.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    if fecha_hasta and not fecha_hasta_parseada:
        return Response(
            {
                'error': (
                    'La fecha hasta no tiene un formato válido.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    if (
        fecha_desde_parseada
        and fecha_hasta_parseada
        and fecha_desde_parseada > fecha_hasta_parseada
    ):
        return Response(
            {
                'error': (
                    'La fecha desde no puede ser posterior a la fecha hasta.'
                )
            },
            status=status.HTTP_400_BAD_REQUEST
        )

    ultima_clase_valida = (
        Calendario.objects
        .filter(
            matricula_id=OuterRef("pk"),
            es_examen=False,
        )
        .exclude(
            estado="cancelada"
        )
        .order_by(
            "-fecha",
            "-hora_fin",
            "-id",
        )
    )

    matriculas = (
        Matricula.objects
        .select_related(
            "estudiante",
            "categoria",
        )
        .annotate(
            instructor_responsable_id=Subquery(
                ultima_clase_valida
                .values("instructor_id")[:1]
            ),
            fecha_finalizacion_reporte=Subquery(
                ultima_clase_valida
                .values("fecha")[:1]
            ),
            hora_finalizacion_reporte=Subquery(
                ultima_clase_valida
                .values("hora_fin")[:1]
            ),
        )
        .filter(
            instructor_responsable_id=instructor_id,
            fecha_finalizacion_reporte__isnull=False,
            hora_finalizacion_reporte__isnull=False,
        )
    )

    if fecha_desde_parseada:
        matriculas = matriculas.filter(
            fecha_finalizacion_reporte__gte=(
                fecha_desde_parseada
            )
        )

    if fecha_hasta_parseada:
        matriculas = matriculas.filter(
            fecha_finalizacion_reporte__lte=(
                fecha_hasta_parseada
            )
        )

    matriculas = matriculas.distinct().order_by(
        "fecha_finalizacion_reporte",
        "estudiante__nombre",
        "estudiante__apellido",
        "id",
    )

    datos = []
    total = Decimal('0')

    for matricula in matriculas:
        estudiante = matricula.estudiante

        recibo = Recibo.objects.filter(
            matricula=matricula
        ).order_by(
            '-fecha_pago',
            '-id'
        ).first()

        if matricula.tipo_curso == 'Principiante':
            horas_practicas = Decimal('15')
        else:
            horas_practicas = Decimal(str(matricula.horas_reforzamiento or 0))

        monto = horas_practicas * tarifa_por_hora
        total += monto

        if matricula.tipo_curso in ['Intermedio', 'Avanzado']:
            observacion = f'Reforzamiento {horas_practicas} horas'
        else:
            observacion = 'Curso principiante 15 horas prácticas'

        datos.append({
            'matricula_id': matricula.id,
                        'estudiante': (
                f'{estudiante.nombre or ""} '
                f'{estudiante.apellido or ""}'
            ).strip(),

            'fecha_matricula': (
                timezone.localtime(
                    matricula.fecha_registro
                ).strftime('%d/%m/%Y')
                if matricula.fecha_registro
                else ''
            ),

            "fecha_finalizacion": (
                matricula.fecha_finalizacion_reporte.strftime(
                    "%d/%m/%Y"
                )
                if matricula.fecha_finalizacion_reporte
                else ""
            ),

            "fecha": (
                matricula.fecha_finalizacion_reporte.strftime(
                    "%d/%m/%Y"
                )
                if matricula.fecha_finalizacion_reporte
                else ""
            ),

            'numero_recibo': (
                recibo.numero_recibo
                if recibo and recibo.numero_recibo
                else ''
            ),
            'horas': float(horas_practicas),
            'tarifa_hora': float(tarifa_por_hora),
            'cobro': float(monto),
            'observaciones': observacion,
        })

    gerente = CargoInstitucional.objects.filter(
        tipo='gerente',
        activo=True
    ).first()

    director = CargoInstitucional.objects.filter(
        tipo='director',
        activo=True
    ).first()

    return Response({
        'instructor': {
            'id': instructor.id,
            'nombre': f'{instructor.nombre or ""} {instructor.apellido or ""}'.strip(),
        },
        'fecha_desde': fecha_desde,
        'fecha_hasta': fecha_hasta,
        'fecha_emision': timezone.now().strftime('%d/%m/%Y'),
        'tarifa_hora': float(tarifa_por_hora),
        'total': float(total),
        'estudiantes': datos,
        'firmas': {
            'gerente_nombre': gerente.nombre if gerente else '',
            'gerente_cargo': gerente.cargo if gerente else '',
            'director_nombre': director.nombre if director else '',
            'director_cargo': director.cargo if director else '',
        },
    })

@api_view(['GET'])
@permission_classes([IsAuthenticated])
def reporte_kilometros_instructor(request):
    if not es_admin(request.user):
        return Response(
            {
                'error': (
                    'Solo Administración o Secretaría pueden '
                    'generar el reporte de kilómetros.'
                )
            },
            status=status.HTTP_403_FORBIDDEN
        )

    fecha_desde = request.query_params.get('desde')
    fecha_hasta = request.query_params.get('hasta')
    instructor_id = request.query_params.get('instructor')

    if not instructor_id:
        return Response(
            {'error': 'Debe seleccionar un instructor.'},
            status=400
        )

    try:
        instructor = Instructor.objects.get(id=instructor_id)
    except Instructor.DoesNotExist:
        return Response(
            {'error': 'Instructor no encontrado.'},
            status=404
        )

    asistencias = Asistencia.objects.select_related(
        'As_estudiante',
        'As_calendario',
        'As_calendario__instructor',
        'As_calendario__matricula',
    ).filter(
        As_calendario__instructor_id=instructor_id,
        estado='asistio',
        km_inicial__isnull=False,
        km_final__isnull=False,
    )

    if fecha_desde:
        asistencias = asistencias.filter(
            As_calendario__fecha__gte=fecha_desde
        )

    if fecha_hasta:
        asistencias = asistencias.filter(
            As_calendario__fecha__lte=fecha_hasta
        )

    asistencias = asistencias.order_by(
        'As_calendario__fecha',
        'As_calendario__hora_inicio',
        'id'
    )

    wb = Workbook()
    ws = wb.active
    ws.title = 'Kilómetros por Instructor'

    thin = Side(style='thin', color='000000')
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    header_fill = PatternFill(fill_type='solid', fgColor='D9EAF7')
    titulo_fill = PatternFill(fill_type='solid', fgColor='FFFFFF')

    ws.merge_cells('A1:J1')
    ws['A1'] = 'Instituto de Formación y Capacitación “Adiact”'
    ws['A1'].font = Font(bold=True, size=16)
    ws['A1'].alignment = Alignment(horizontal='center')
    ws.merge_cells('A2:J2')
    ws['A2'] = 'Somos expertos en Formación y Capacitación del Talento Humano'
    ws['A2'].alignment = Alignment(horizontal='center')
    ws.merge_cells('A3:J3')
    ws['A3'] = 'Ética, Integridad, Dedicación y Solidaridad'
    ws['A3'].font = Font(bold=True)
    ws['A3'].alignment = Alignment(horizontal='center')
    ws.merge_cells('A5:J5')
    ws['A5'] = 'REPORTE DE KILÓMETROS RECORRIDOS POR INSTRUCTOR'
    ws['A5'].font = Font(bold=True, size=14)
    ws['A5'].alignment = Alignment(horizontal='center')
    ws['A5'].fill = titulo_fill
    ws['A7'] = 'Instructor:'
    ws['B7'] = f'{instructor.nombre or ""} {instructor.apellido or ""}'.strip()
    ws['A8'] = 'Fecha de emisión:'
    ws['B8'] = timezone.localdate()
    ws['B8'].number_format = 'dd/mm/yyyy'

    fecha_desde_excel = (
        parse_date(fecha_desde)
        if fecha_desde
        else None
    )

    fecha_hasta_excel = (
        parse_date(fecha_hasta)
        if fecha_hasta
        else None
    )

    ws['D8'] = 'Desde:'
    ws['E8'] = fecha_desde_excel or 'Inicio'
    ws['F8'] = 'Hasta:'
    ws['G8'] = fecha_hasta_excel or 'Fin'

    if fecha_desde_excel:
        ws['E8'].number_format = 'dd/mm/yyyy'

    if fecha_hasta_excel:
        ws['G8'].number_format = 'dd/mm/yyyy'

    encabezados = [
        'No.',
        'Fecha',
        'Instructor',
        'Estudiante',
        'Cédula',
        'Clase No.',
        'Hora Inicio',
        'Hora Fin',
        'Km Inicial',
        'Km Final',
        'Km Recorridos',
    ]

    fila_encabezado = 10

    for col, encabezado in enumerate(encabezados, start=1):
        celda = ws.cell(row=fila_encabezado, column=col, value=encabezado)
        celda.font = Font(bold=True)
        celda.fill = header_fill
        celda.border = border
        celda.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)

    fila_inicio_datos = fila_encabezado + 1
    fila = fila_inicio_datos

    for index, asistencia in enumerate(asistencias, start=1):
        clase = asistencia.As_calendario
        estudiante = asistencia.As_estudiante

        # Convertir explícitamente a número.
        km_inicial = float(asistencia.km_inicial or 0)
        km_final = float(asistencia.km_final or 0)

        valores = [
            index,
            clase.fecha if clase.fecha else None,
            f'{instructor.nombre or ""} {instructor.apellido or ""}'.strip(),
            f'{estudiante.nombre or ""} {estudiante.apellido or ""}'.strip(),
            estudiante.cedula or '',
            int(clase.numero_clase or 0),
            clase.hora_inicio if clase.hora_inicio else None,
            clase.hora_fin if clase.hora_fin else None,
            km_inicial,
            km_final,
        ]

        # Columnas A hasta J.
        for col, valor in enumerate(valores, start=1):
            celda = ws.cell(
                row=fila,
                column=col,
                value=valor
            )

            celda.border = border
            celda.alignment = Alignment(
                horizontal='center',
                vertical='center',
                wrap_text=True
            )

        # Fecha real de Excel.
        ws.cell(
            row=fila,
            column=2
        ).number_format = 'dd/mm/yyyy'

        # Horas reales de Excel.
        ws.cell(
            row=fila,
            column=7
        ).number_format = 'hh:mm'

        ws.cell(
            row=fila,
            column=8
        ).number_format = 'hh:mm'

        # Kilometrajes como números.
        ws.cell(
            row=fila,
            column=9
        ).number_format = '#,##0.00'

        ws.cell(
            row=fila,
            column=10
        ).number_format = '#,##0.00'

        # Columna K: Km recorridos = Km final - Km inicial.
        celda_km_recorridos = ws.cell(
            row=fila,
            column=11,
            value=f'=J{fila}-I{fila}'
        )

        celda_km_recorridos.number_format = '#,##0.00'
        celda_km_recorridos.border = border
        celda_km_recorridos.alignment = Alignment(
            horizontal='center',
            vertical='center'
        )

        fila += 1

    fila_total = fila

    ws.cell(
        row=fila_total,
        column=10,
        value='TOTAL KM:'
    )

    celda_total = ws.cell(
        row=fila_total,
        column=11
    )

    if fila_total > fila_inicio_datos:
        celda_total.value = (
            f'=SUM(K{fila_inicio_datos}:K{fila_total - 1})'
        )
    else:
        celda_total.value = 0

    celda_total.number_format = '#,##0.00'

    for col in range(10, 12):
        celda = ws.cell(
            row=fila_total,
            column=col
        )
        celda.font = Font(bold=True)
        celda.border = border
        celda.fill = header_fill
        celda.alignment = Alignment(
            horizontal='center',
            vertical='center'
        )

    fila = fila_total
    fila_footer = fila_total + 3

    ws.merge_cells(start_row=fila_footer, start_column=1, end_row=fila_footer, end_column=12)
    ws.cell(
        row=fila_footer,
        column=1,
        value='Gasolinera Uno Sutiaba 1 cuadra al norte ½ cuadra al oeste. León, Nicaragua'
    )

    ws.merge_cells(start_row=fila_footer + 1, start_column=1, end_row=fila_footer + 1, end_column=12)
    ws.cell(
        row=fila_footer + 1,
        column=1,
        value='Teléfonos: 2311-1333 y 8966-3770. email: institutoadiact@esesa.com.ni'
    )

    ws.cell(row=fila_footer, column=1).alignment = Alignment(horizontal='center')
    ws.cell(row=fila_footer + 1, column=1).alignment = Alignment(horizontal='center')

    columnas = {
        'A': 8,
        'B': 15,
        'C': 28,
        'D': 28,
        'E': 18,
        'F': 12,
        'G': 14,
        'H': 14,
        'I': 14,
        'J': 14,
        'K': 16,
        'L': 28,
    }

    for columna, ancho in columnas.items():
        ws.column_dimensions[columna].width = ancho

    response = HttpResponse(
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )

    response['Content-Disposition'] = 'attachment; filename="reporte_kilometros_instructor.xlsx"'
    wb.calculation.fullCalcOnLoad = True
    wb.calculation.forceFullCalc = True
    wb.calculation.calcMode = "auto"
    wb.save(response)

    return response

def certificado_numero(valor):
    try:
        if valor is None or valor == "":
            return None

        return Decimal(str(valor).replace(",", "."))
    except (InvalidOperation, ValueError, TypeError):
        return None

def certificado_fecha_larga(fecha):
    if not fecha:
        return ""

    meses = {
        1: "enero",
        2: "febrero",
        3: "marzo",
        4: "abril",
        5: "mayo",
        6: "junio",
        7: "julio",
        8: "agosto",
        9: "septiembre",
        10: "octubre",
        11: "noviembre",
        12: "diciembre",
    }
    return f"{fecha.day:02d} de {meses[fecha.month]} del {fecha.year}"

def certificado_mes_anio(fecha):
    if not fecha:
        return ""

    meses = {
        1: "enero",
        2: "febrero",
        3: "marzo",
        4: "abril",
        5: "mayo",
        6: "junio",
        7: "julio",
        8: "agosto",
        9: "septiembre",
        10: "octubre",
        11: "noviembre",
        12: "diciembre",
    }

    return (
        f"{meses[fecha.month]} "
        f"del año {fecha.year}"
    )

def certificado_periodo(fecha_inicio, fecha_fin):
    if not fecha_inicio and not fecha_fin:
        return ""

    if not fecha_inicio:
        return (
            f"hasta el {certificado_fecha_larga(fecha_fin)}"
        )

    if not fecha_fin:
        return (
            f"desde el {certificado_fecha_larga(fecha_inicio)}"
        )

    meses = {
        1: "enero",
        2: "febrero",
        3: "marzo",
        4: "abril",
        5: "mayo",
        6: "junio",
        7: "julio",
        8: "agosto",
        9: "septiembre",
        10: "octubre",
        11: "noviembre",
        12: "diciembre",
    }

    if (
        fecha_inicio.year == fecha_fin.year
        and fecha_inicio.month == fecha_fin.month
    ):
        return (
            f"del {fecha_inicio.day:02d} al "
            f"{fecha_fin.day:02d} de "
            f"{meses[fecha_fin.month]} del "
            f"{fecha_fin.year}"
        )

    if fecha_inicio.year == fecha_fin.year:
        return (
            f"del {fecha_inicio.day:02d} de "
            f"{meses[fecha_inicio.month]} al "
            f"{fecha_fin.day:02d} de "
            f"{meses[fecha_fin.month]} del "
            f"{fecha_fin.year}"
        )

    return (
        f"del {certificado_fecha_larga(fecha_inicio)} "
        f"al {certificado_fecha_larga(fecha_fin)}"
    )


def certificado_obtener_datos(desde, hasta):
    desde_fecha = parse_date(desde)
    hasta_fecha = parse_date(hasta)

    if (
        not desde_fecha
        or not hasta_fecha
        or desde_fecha > hasta_fecha
    ):
        return []

    matriculas = (
        Matricula.objects
        .select_related(
            "estudiante",
            "categoria",
            "plan_de_estudio",
        )
        .prefetch_related(
            Prefetch(
                "notas",
                queryset=(
                    Notas.objects
                    .filter(
                        tipo_nota__in=[
                            "teorico",
                            "practico",
                        ]
                    )
                    .order_by(
                        "-fecha_registro",
                        "-id",
                    )
                ),
                to_attr="notas_certificado",
            )
        )
        .annotate(
            fecha_inicio_curso=models.Min(
                "clases__fecha",
                filter=(
                    Q(clases__es_examen=False)
                    & ~Q(clases__estado="cancelada")
                ),
            ),
            fecha_finalizacion_curso=models.Max(
                "clases__fecha",
                filter=(
                    Q(clases__es_examen=False)
                    & ~Q(clases__estado="cancelada")
                ),
            ),
        )
        .filter(
            tipo_curso="Principiante",
            fecha_finalizacion_curso__gte=desde_fecha,
            fecha_finalizacion_curso__lte=hasta_fecha,
        )
        .order_by(
            "estudiante__apellido",
            "estudiante__nombre",
            "id",
        )
    )

    resultados = []

    for matricula in matriculas:
        notas_por_tipo = {}

        for nota in matricula.notas_certificado:
            if nota.tipo_nota not in notas_por_tipo:
                notas_por_tipo[nota.tipo_nota] = nota

        nota_teorica_obj = notas_por_tipo.get("teorico")
        nota_practica_obj = notas_por_tipo.get("practico")

        if not nota_teorica_obj or not nota_practica_obj:
            continue

        nota_teorica = certificado_numero(
            nota_teorica_obj.nota
        )

        nota_practica = certificado_numero(
            nota_practica_obj.nota
        )

        if nota_teorica is None or nota_practica is None:
            continue

        if (
            nota_teorica < Decimal("80")
            or nota_practica < Decimal("80")
        ):
            continue

        fecha_finalizacion = (
            matricula.fecha_finalizacion_curso
        )

        estudiante = matricula.estudiante

        categoria = (
            matricula.categoria.nombre
            if matricula.categoria
            else ""
        )

        resultados.append({
            "id": matricula.id,
            "estudiante": (
                f"{estudiante.nombre} "
                f"{estudiante.apellido}"
            ).strip().upper(),
            "cedula": estudiante.cedula,
            "categoria": categoria,
            "tipo_curso": matricula.tipo_curso,

            # Primera clase regular válida del curso.
            "fecha_inicio": matricula.fecha_inicio_curso,

            # Fecha real guardada cuando finalizó la matrícula.
            "fecha_egreso": fecha_finalizacion,

            "nota_teorica": nota_teorica,
            "nota_practica": nota_practica,
        })

    return resultados

COLOR_AZUL_CERTIFICADO = RGBColor(53, 86, 112)
COLOR_DORADO_CERTIFICADO = RGBColor(226, 174, 94)
COLOR_NEGRO_CERTIFICADO = RGBColor(20, 20, 20)
COLOR_BLANCO_CERTIFICADO = RGBColor(255, 255, 255)


def certificado_ppt_rectangulo(
    slide,
    x,
    y,
    ancho,
    alto,
    color_relleno=None,
    color_borde=None,
    grosor_borde=1,
):
    forma = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(ancho),
        Inches(alto),
    )

    if color_relleno is None:
        forma.fill.background()
    else:
        forma.fill.solid()
        forma.fill.fore_color.rgb = color_relleno

    if color_borde is None:
        forma.line.fill.background()
    else:
        forma.line.color.rgb = color_borde
        forma.line.width = Pt(grosor_borde)

    return forma


def certificado_ppt_texto(
    slide,
    texto,
    x,
    y,
    ancho,
    alto,
    tamano=10,
    negrita=False,
    alineacion=PP_ALIGN.CENTER,
    color=COLOR_NEGRO_CERTIFICADO,
):
    caja = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(ancho),
        Inches(alto),
    )

    marco = caja.text_frame
    marco.clear()
    marco.word_wrap = True
    marco.vertical_anchor = MSO_ANCHOR.MIDDLE

    marco.margin_left = Inches(0.03)
    marco.margin_right = Inches(0.03)
    marco.margin_top = Inches(0.01)
    marco.margin_bottom = Inches(0.01)

    parrafo = marco.paragraphs[0]
    parrafo.alignment = alineacion
    parrafo.space_before = Pt(0)
    parrafo.space_after = Pt(0)
    parrafo.line_spacing = 1

    texto_run = parrafo.add_run()
    texto_run.text = str(texto)

    texto_run.font.name = "Arial"
    texto_run.font.size = Pt(tamano)
    texto_run.font.bold = negrita
    texto_run.font.color.rgb = color

    return caja


def certificado_ppt_agregar_esquinas(
    slide,
    x,
    y,
    ancho,
    alto,
):
    # Las decoraciones quedan dentro del certificado,
    # sobre las esquinas del área blanca.
    posiciones = [
        (x + 0.42, y + 0.42),
        (x + ancho - 0.71, y + 0.42),
        (x + 0.42, y + alto - 0.71),
        (
            x + ancho - 0.71,
            y + alto - 0.71,
        ),
    ]

    for esquina_x, esquina_y in posiciones:
        certificado_ppt_rectangulo(
            slide,
            esquina_x,
            esquina_y,
            0.22,
            0.22,
            color_relleno=None,
            color_borde=COLOR_NEGRO_CERTIFICADO,
            grosor_borde=1.2,
        )

        certificado_ppt_rectangulo(
            slide,
            esquina_x + 0.07,
            esquina_y + 0.07,
            0.22,
            0.22,
            color_relleno=None,
            color_borde=COLOR_NEGRO_CERTIFICADO,
            grosor_borde=0.9,
        )

def certificado_ppt_texto_principal(
    slide,
    item,
    x,
    y,
    ancho,
    alto,
):
    if item["fecha_inicio"]:
        periodo = (
            f"del {certificado_fecha_larga(item['fecha_inicio'])} "
            f"al {certificado_fecha_larga(item['fecha_egreso'])}"
        )
    else:
        periodo = (
            f"hasta el "
            f"{certificado_fecha_larga(item['fecha_egreso'])}"
        )

    caja = slide.shapes.add_textbox(
        Inches(x),
        Inches(y),
        Inches(ancho),
        Inches(alto),
    )

    marco = caja.text_frame
    marco.clear()
    marco.word_wrap = True
    marco.vertical_anchor = MSO_ANCHOR.TOP

    marco.margin_left = Inches(0.02)
    marco.margin_right = Inches(0.02)
    marco.margin_top = Inches(0.01)
    marco.margin_bottom = Inches(0.01)

    parrafo = marco.paragraphs[0]
    parrafo.alignment = PP_ALIGN.LEFT
    parrafo.space_before = Pt(0)
    parrafo.space_after = Pt(0)
    parrafo.line_spacing = 1

    fragmentos = [
        (
            "Ha cumplido con el plan de instrucción "
            "teórico y práctico aprobado por la DSTN, "
            "de Principiante, para optar a la Licencia "
            "de Conducir de Tipo Ordinaria en la Categoría ",
            False,
        ),
        (
            item["categoria"] or "__________",
            True,
        ),
        (
            f" impartido en el periodo comprendido "
            f"{periodo}, habiendo obtenido las "
            f"siguientes calificaciones:",
            False,
        ),
    ]

    for texto, subrayado in fragmentos:
        texto_run = parrafo.add_run()
        texto_run.text = texto
        texto_run.font.name = "Arial"
        texto_run.font.size = Pt(9.6)
        texto_run.font.bold = True
        texto_run.font.underline = subrayado
        texto_run.font.color.rgb = COLOR_NEGRO_CERTIFICADO

    return caja


def certificado_ppt_agregar_certificado(
    slide,
    item,
    logo_path,
    auto_path,
    posicion_y,
    nombre_gerente,
):
    x = 0.20
    y = posicion_y
    ancho = 8.10
    alto = 5.15

    # Fondo dorado exterior.
    certificado_ppt_rectangulo(
        slide,
        x,
        y,
        ancho,
        alto,
        color_relleno=COLOR_DORADO_CERTIFICADO,
        color_borde=COLOR_DORADO_CERTIFICADO,
        grosor_borde=1,
    )

    # Banda azul gruesa sobre el marco dorado.
    certificado_ppt_rectangulo(
        slide,
        x + 0.28,
        y + 0.28,
        ancho - 0.56,
        alto - 0.56,
        color_relleno=COLOR_AZUL_CERTIFICADO,
        color_borde=COLOR_AZUL_CERTIFICADO,
        grosor_borde=1,
    )

    # Área blanca interior.
    certificado_ppt_rectangulo(
        slide,
        x + 0.42,
        y + 0.42,
        ancho - 0.84,
        alto - 0.84,
        color_relleno=COLOR_BLANCO_CERTIFICADO,
        color_borde=COLOR_AZUL_CERTIFICADO,
        grosor_borde=1.3,
    )

    # Línea azul interior fina.
    certificado_ppt_rectangulo(
        slide,
        x + 0.50,
        y + 0.50,
        ancho - 1.00,
        alto - 1.00,
        color_relleno=None,
        color_borde=COLOR_AZUL_CERTIFICADO,
        grosor_borde=0.8,
    )

    certificado_ppt_agregar_esquinas(
        slide,
        x,
        y,
        ancho,
        alto,
    )

    # Líneas doradas verticales interiores.
    certificado_ppt_rectangulo(
        slide,
        x + 0.62,
        y + 0.72,
        0.01,
        alto - 1.44,
        color_relleno=COLOR_DORADO_CERTIFICADO,
        color_borde=None,
    )

    certificado_ppt_rectangulo(
        slide,
        x + ancho - 0.63,
        y + 0.72,
        0.01,
        alto - 1.44,
        color_relleno=COLOR_DORADO_CERTIFICADO,
        color_borde=None,
    )

    # Título.
    certificado_ppt_texto(
        slide,
        "-ESCUELA DE MANEJO EL CACIQUE ADIACT-",
        x + 0.55,
        y + 0.54,
        ancho - 1.10,
        0.25,
        tamano=18,
        negrita=True,
    )

    # Franja dorada.
    certificado_ppt_rectangulo(
        slide,
        x + 2.30,
        y + 0.82,
        3.50,
        0.43,
        color_relleno=COLOR_DORADO_CERTIFICADO,
        color_borde=None,
    )

    certificado_ppt_texto(
        slide,
        "CERTIFICA QUE:",
        x + 2.45,
        y + 0.83,
        3.20,
        0.39,
        tamano=17,
        negrita=True,
    )

    # Imágenes ya definidas en static/certificados.
    slide.shapes.add_picture(
        logo_path,
        Inches(x + 0.78),
        Inches(y + 0.78),
        width=Inches(0.80),
    )

    slide.shapes.add_picture(
        auto_path,
        Inches(x + ancho - 1.80),
        Inches(y + 0.81),
        width=Inches(1.00),
    )

    tamano_nombre = (
        15
        if len(item["estudiante"]) > 38
        else 17
    )

    certificado_ppt_texto(
        slide,
        item["estudiante"],
        x + 1.00,
        y + 1.29,
        ancho - 2.00,
        0.38,
        tamano=tamano_nombre,
        negrita=True,
    )

    certificado_ppt_texto(
        slide,
        f"Cédula: {item['cedula']}.",
        x + 1.50,
        y + 1.66,
        ancho - 3.00,
        0.30,
        tamano=12,
        negrita=True,
    )

    certificado_ppt_texto_principal(
        slide,
        item,
        x + 0.72,
        y + 1.94,
        ancho - 1.44,
        0.78,
    )

    certificado_ppt_texto(
        slide,
        (
            f"Evaluación Teórica: "
            f"{int(item['nota_teorica'])} puntos.\n"
            f"Evaluación Práctica: "
            f"{int(item['nota_practica'])} puntos."
        ),
        x + 1.65,
        y + 2.58,
        ancho - 3.30,
        0.50,
        tamano=10.5,
        negrita=True,
    )

    certificado_ppt_texto(
        slide,
        (
            "Registrado en el Asiento No. _____ "
            "del Folio No. ______ "
            "del Libro ___."
        ),
        x + 0.90,
        y + 3.03,
        ancho - 1.80,
        0.25,
        tamano=9.8,
        negrita=True,
    )

    fecha_emision = item["fecha_egreso"]

    certificado_ppt_texto(
        slide,
        (
            "Dado en la ciudad de León a los "
            f"{fecha_emision.day:02d} días del mes de "
            f"{certificado_mes_anio(fecha_emision)}."
        ),
        x + 0.85,
        y + 3.25,
        ancho - 1.70,
        0.25,
        tamano=9.8,
        negrita=True,
    )

    # Línea de firma.
    certificado_ppt_rectangulo(
        slide,
        x + 2.30,
        y + 3.95,
        ancho - 4.60,
        0.01,
        color_relleno=COLOR_NEGRO_CERTIFICADO,
        color_borde=COLOR_NEGRO_CERTIFICADO,
        grosor_borde=0.5,
    )

    certificado_ppt_texto(
        slide,
        nombre_gerente,
        x + 1.70,
        y + 3.99,
        ancho - 3.40,
        0.27,
        tamano=10.2,
        negrita=True,
    )

    certificado_ppt_texto(
        slide,
        "Director.",
        x + 2.80,
        y + 4.25,
        ancho - 5.60,
        0.22,
        tamano=9.8,
        negrita=True,
    )


def certificado_crear_powerpoint(certificados):
    presentacion = Presentation()

    gerente_configurado = (
        CargoInstitucional.objects
        .filter(
            tipo="gerente",
            activo=True,
        )
        .order_by("-id")
        .first()
    )

    if (
        not gerente_configurado
        or not str(
            gerente_configurado.nombre or ""
        ).strip()
    ):
        raise ValueError(
            "No existe un gerente activo configurado."
        )

    nombre_sin_prefijo = str(
        gerente_configurado.nombre
    ).strip()

    if nombre_sin_prefijo.lower().startswith("lic."):
        nombre_sin_prefijo = (
            nombre_sin_prefijo[4:].strip()
        )

    nombre_gerente = (
        f"Lic. "
        f"{nombre_sin_prefijo.rstrip('.').upper()}."
    )

    # Una diapositiva equivale a una hoja tamaño carta.
    presentacion.slide_width = Inches(8.5)
    presentacion.slide_height = Inches(11)

    logo_path = os.path.join(
        settings.BASE_DIR,
        "static",
        "certificados",
        "logo.png",
    )

    auto_path = os.path.join(
        settings.BASE_DIR,
        "static",
        "certificados",
        "auto.png",
    )

    if not os.path.exists(logo_path):
        raise FileNotFoundError(
            f"No se encontró el logo en: {logo_path}"
        )

    if not os.path.exists(auto_path):
        raise FileNotFoundError(
            f"No se encontró la imagen del auto en: {auto_path}"
        )

    diapositiva_vacia = presentacion.slide_layouts[6]

    for indice in range(0, len(certificados), 2):
        grupo = certificados[indice:indice + 2]
        slide = presentacion.slides.add_slide(
            diapositiva_vacia
        )

        fondo = slide.background.fill
        fondo.solid()
        fondo.fore_color.rgb = COLOR_BLANCO_CERTIFICADO

        posiciones_y = [
            0.20,
            5.65,
        ]

        for posicion, item in enumerate(grupo):
            certificado_ppt_agregar_certificado(
                slide,
                item,
                logo_path,
                auto_path,
                posiciones_y[posicion],
                nombre_gerente,
            )

    return presentacion




@api_view(["GET"])
@permission_classes([IsAuthenticated])
def certificados_egresados(request):
    if not es_admin(request.user):
        return Response(
            {"detail": "No tienes permiso para consultar certificados."},
            status=403,
        )

    desde = request.GET.get("desde")
    hasta = request.GET.get("hasta")

    if not desde or not hasta:
        return Response(
            {"detail": "Debe enviar fecha desde y fecha hasta."},
            status=400,
        )

    certificados = certificado_obtener_datos(desde, hasta)
    data = []

    for item in certificados:
        data.append({
            "id": item["id"],
            "estudiante": item["estudiante"],
            "cedula": item["cedula"],
            "categoria": item["categoria"],
            "tipo_curso": item["tipo_curso"],
            "fecha_inicio": item["fecha_inicio"].isoformat() if item["fecha_inicio"] else "",
            "fecha_egreso": item["fecha_egreso"].isoformat() if item["fecha_egreso"] else "",
            "nota_teorica": int(item["nota_teorica"]),
            "nota_practica": int(item["nota_practica"]),
        })
    return Response(data)

@api_view(["GET"])
@permission_classes([IsAuthenticated])
def certificados_egresados_powerpoint(request):
    if not es_admin(request.user):
        return Response(
            {
                "detail": (
                    "No tienes permiso para generar "
                    "certificados."
                )
            },
            status=403,
        )

    desde = request.GET.get("desde")
    hasta = request.GET.get("hasta")

    if not desde or not hasta:
        return Response(
            {
                "detail": (
                    "Debe enviar fecha desde "
                    "y fecha hasta."
                )
            },
            status=400,
        )

    certificados = certificado_obtener_datos(
        desde,
        hasta,
    )

    if not certificados:
        return Response(
            {
                "detail": (
                    "No hay estudiantes del curso "
                    "principiante con nota teórica y "
                    "práctica mayor o igual a 80 en "
                    "ese rango de fechas."
                )
            },
            status=404,
        )

    try:
        presentacion = certificado_crear_powerpoint(
            certificados
        )
    except FileNotFoundError as error:
        return Response(
            {"detail": str(error)},
            status=500,
        )
    except ValueError as error:
        return Response(
            {"detail": str(error)},
            status=400,
        )

    archivo = BytesIO()
    presentacion.save(archivo)
    archivo.seek(0)

    response = HttpResponse(
        archivo.getvalue(),
        content_type=(
            "application/vnd.openxmlformats-officedocument."
            "presentationml.presentation"
        ),
    )

    response["Content-Disposition"] = (
        f'attachment; filename='
        f'"certificados_{desde}_{hasta}.pptx"'
    )

    return response
